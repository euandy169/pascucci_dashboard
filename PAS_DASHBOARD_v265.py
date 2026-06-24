"""
PASCUCCI Consumer Intelligence Dashboard
Consumer Data-driven Marketing for PASCUCCI

═══════════════════════════════════════════════════════════
버전 이력 (Version History)
───────────────────────────────────────────────────────────
v2.6.5(beta) (2026-06-24)
  [기능] 트렌드 조기경보 탭 확장 — '식음료 인접 트렌드 조기경보'(커피 외 디저트·베이커리·카페푸드·식재료 VoC)
         별도 레이어로 분리(브랜드 NSS·SOV 미반영). 급상승 키워드 레이더·카테고리 볼륨/감성·대표 원문.
v2.6.4(beta) (2026-06-23)
  [기능] 2단계 데이터 필터 — ①소스(블로그/커뮤니티) ②콘텐츠 품질(VoC만/홍보·보도성만) + 소스별 특성 비교
         선택한 렌즈로 모든 분석 범위 재설정. 블로그 혜택·홍보성/보도 글 휴리스틱 태깅(_src/_ctype)
v2.6.3(beta) (2026-06-23)
  [기능] 데이터 로딩에 gzip(.csv.gz/.gz) 지원 추가, 최신 파일 선택을 파일명 날짜(YYYYMMDD) 우선으로 robust화
         (+news_master 파일은 메인 데이터 후보에서 제외)
v2.6.2(beta) (2026-06-22)
  [보안] 접근 제한 — 공유 비밀번호 관문(옵션 C). Secrets의 app_password로 사업부 전용 접근, 세션당 1회 입력 후 통과. 미설정 시 경고 후 개방.
v2.6.1(beta) (2026-06-21)
  [수정] 리포트 차트 누락 해결 — matplotlib 부재 시 네이티브 PPT 차트로 폴백(빈자리 방지)
  [기능] DOCX 리포트도 PPTX 동일 스킴(Exec Summary·5챕터·결과/의미/함의/시나리오)으로 재작성, 메뉴 '(개발 중)' 제거
  [기능] 원두선물·환율 우측 그래프 원/달러(USD)→원/유로(EUR)로 변경
v2.6.0(beta) (2026-06-21)
  [기능] 분석 리포트 전면 재설계 — 사전저작 코멘트(report_commentary.json) 기반 컨설팅형 리포트
         표지·Executive Summary·5개 챕터(결과·의미·함의·시나리오)·마지막 장 / 차트 이미지화 / 지문 동기화 가드
v2.5.7(beta) (2026-06-21)
  [수정] 리포트 템플릿 탐색 강화(대소문자·하위폴더·철자 무시 재귀) + 미발견 시 진단 슬라이드
v2.5.6(beta) (2026-06-20)
  [기능] 분석 리포트 PPTX를 브랜드 템플릿(insight_report_templete.pptx) 기반으로 전면 개편
         (표지 분석주제 자동작성·본문 데이터·마지막 장 고정 / 맑은 고딕, 챕터10·헤드라인18·본문12·차트9pt)
v2.5.5(beta) (2026-06-20)
  [기능] 분석 리포트(PPTX/DOCX)에 '2025 서울시민 먹거리조사' 외부 근거 자동 인용 연동
         (reference/seoul_food_survey_ref.json 로더 + 섹션→토픽 매핑 + 코멘트 삽입)
v2.5.4(beta) (2026-06-16)
  [디자인] 분석영역 버튼 텍스트 볼드 강제, 하위 분석탭 활성색 옅은 블루로(텍스트 가독성),
           소비자 VoC·분석 리포트 다운로드에 '(개발 중)' 표기
v2.5.3(beta) (2026-06-16)
  [디자인] 카테고리 설명 풋노트화·버튼 텍스트 확대/볼드, 활성 카테고리 버튼 옅은 블루,
           세부 분석탭 확대·볼드, PASCUCCI Brand Status 명칭, 메인 최상단 타이틀
  [기능] 원두선물·환율 명칭 정리, 포지셔닝맵 사분면 직관화, 사이드바 업로더 제거
v2.5.2(beta) (2026-06-16)
  [디자인] 분석영역 대형 카테고리 버튼 UI, 카테고리 설명문 개편, 파스쿠찌 중심화,
           페르소나 가독성(상세보기 통합·이모지·키워드 정제), 브랜드 이미지맵 명칭
  [기능] 프로모션/CA맵 원문·키워드 매칭 정확도 개선, CA맵 축 재점검·파스쿠찌 강조
v2.5.1(beta) (2026-06-15)
  [버그] 로열티 차트 파스쿠찌 막대 누적(중복행) 수정, 드라이버 원문 매칭 정확도 개선
         (속성 문장 스니펫 추출 + 감성어 동반 조건 + 관련도 정렬)
  [디자인] 사이드바 버튼 한글화(컨슈머 인텔리전스 대시보드)
v2.4.2 (2026-06-15)
  [견고성] 데이터 자동탐색에 data/(점없음)·리포 루트 추가, 미발견 시 탐색경로 안내,
           로딩 오류를 화면에 명시(Streamlit Cloud 자가진단)
v2.4.1 (2026-06-15)
  [기능] 데이터/.data·뉴스/.news GitHub 폴더 자동 로딩(최신 파일), 키워드 버블 가독성,
         원두가격·환율 풋노트, 뉴스룸 유형별 대표기사 정리, 드라이버 원문·MD/굿즈,
         프로모션 브랜드별 시계열, 페르소나 Pain-point·Key Buying Factor
v2.3.6 (2026-06-13)
  [기능] PAS 뉴스룸 브리핑 보강 — 날짜·브랜드·유형별 이모지 건수 통계 + 요약(G칼럼) 기반 종합 브리핑 기사 분리
v2.3.5 (2026-06-13)
  [기능] 원두·환율: 원화환산 가격 삭제·외부시세 버튼 복원, PAS 뉴스룸(명칭)·오늘의 뉴스 브리핑(1000자 요약)
v2.3.4 (2026-06-13)
  [기능] 소비맥락+세분화 화면 통합, 포지셔닝 레이더에 파스쿠찌(굵은 라인) 추가,
         페르소나 네이밍 실무형·크리에이티브 개편, 분석 리포트 PPT(시사점·Lessons 포함)
  [디자인] CID 버튼 분리·업로드 상단 이동, 시작화면 안내문구 삭제, 원두·환율 모니터링 명칭
v2.3.3 (2026-06-13)
  [기능] News Monitoring 자동 연동 — 최신 news_master_*.csv 자동 탐색·표시(필터·키워드·링크)
v2.3.2 (2026-06-13)
  [디버그] 토픽 카드 IndexError(cluster_colors) 수정, 토픽 클러스터링 제목 중복 제거,
           PPTX 미설치 안내 개선, 원자재 폴백 출처 설명 추가, News Monitoring 명칭,
           CID 버튼 명칭변경·리포트 상단 이동
v2.3.1 (2026-06-12)
  [디자인] 사이드바 타이틀 확대·부제 삭제, 분석 리포트로 명칭변경·위치이동, MI 버튼
  [기능] TARGET_BRANDS 파스쿠찌 포함(전 기능), Market Intelligence 모듈(원자재·환율·뉴스),
         소비자 VoC 모니터링(TBD), 분석 리포트 PPT·시장맥락 고도화
  [프로모션] 페이행사 추가, 최근 급증 프로모션·브랜드 cross 체크
  [LDA] 토픽 클러스터링 merge, 토픽 6~8개 세분화
  [고객] 페르소나 8~10개, 모듈화·라이프스타일·인접 브랜드 표시
  [키워드] 시그니처·트렌디(두쫀쿠·우베 등) 보강
v2.1.2 (2026-06-05)
  [디자인] 사이드바 버전·Copyright 표기, 필터 무채색, 푸터 Copyright
  [키워드] 디저트(아이스크림·젤라또), 프로모션(기프티콘·모바일상품권·MD) 보강
  [텍스트마이닝] 리스크탐지 브랜드 필터 정합성 수정, 무의미 부사·수식어
                  분석 집계 제외(불용어 확장)
  [기능] NSS 시사점 다각화, 찐팬지수 설명·해설 보강, 분석결과 docx 다운로드,
         촉진 반응강도 지표화 + 최근 프로모션·브랜드, 촉진활동 세부·원문 확인
v2.1.1 (이전) 19개 탭, 4대 고도화 분석
═══════════════════════════════════════════════════════════

데이터: 82Cook / 클리앙 / 네이트판 Community VoC

실행 방법:
  pip install streamlit plotly pandas numpy scikit-learn scipy python-docx
  streamlit run PAS_DASHBOARD_v212.py
"""

APP_VERSION = "v2.6.5(beta)"
APP_DATE = "2026-06-15"
COPYRIGHT = "Copyright 2026 PAS DIVISION Paris Croissant Co., Ltd. All Rights Reserved"

import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import numpy as np
import json
import re
from collections import Counter
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import LatentDirichletAllocation, PCA
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
import scipy.linalg as la
import io
import os
import glob
try:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
try:
    import pptx as _pptx_check  # noqa: F401
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ══════════════════════════════════════════════════════
# 기본 설정
# ══════════════════════════════════════════════════════
st.set_page_config(
    page_title="PASCUCCI Consumer Intelligence Dashboard",
    page_icon="☕",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
    .stApp { background-color: #F7F9FC; }
    .main .block-container { background-color: #F7F9FC; padding-top: 1.5rem; max-width: 1400px; }

    [data-testid="stSidebar"] { background-color: #EDF1F5 !important; border-right: 2px solid #CBD5E0; }
    [data-testid="stSidebar"] p,
    [data-testid="stSidebar"] span,
    [data-testid="stSidebar"] label,
    [data-testid="stSidebar"] div { color: #2D3748 !important; }
    [data-testid="stSidebar"] h1,
    [data-testid="stSidebar"] h2,
    [data-testid="stSidebar"] h3 { color: #1B2A4A !important; font-size: 1rem; font-weight: 700; }
    [data-testid="stSidebar"] hr { border-color: #CBD5E0; }
    [data-testid="stSidebar"] .stMarkdown { color: #2D3748 !important; }

    .stTabs [data-baseweb="tab-list"] {
        background-color: #FFFFFF; border-radius: 10px; padding: 4px 6px;
        gap: 2px; box-shadow: 0 1px 4px rgba(0,0,0,0.08); border: 1px solid #E2E8F0;
    }
    .stTabs [data-baseweb="tab"] {
        color: #4A5568; font-size: 13px; font-weight: 600;
        border-radius: 7px; padding: 7px 14px; background: transparent;
    }
    .stTabs [aria-selected="true"] { background-color: #DCE9F7 !important; color: #1B2A4A !important; }
    .stTabs [aria-selected="true"] p { color: #1B2A4A !important; }

    h1, h2, h3, h4 { color: #1B2A4A !important; }
    p, li { color: #2D3748; }

    .kpi-card {
        background: #FFFFFF; border: 1px solid #E2E8F0;
        border-top: 4px solid #2D6BC4; border-radius: 10px;
        padding: 18px 14px; text-align: center;
        box-shadow: 0 2px 8px rgba(0,0,0,0.06);
    }
    .kpi-icon { font-size: 1.5rem; }
    .kpi-value { font-size: 1.8rem; font-weight: 800; color: #1B2A4A; margin: 4px 0; }
    .kpi-label { font-size: 0.82rem; color: #718096; font-weight: 500; }

    .section-title {
        font-size: 1.05rem; font-weight: 700; color: #1B2A4A;
        border-left: 4px solid #2D6BC4; padding: 8px 12px;
        margin: 20px 0 12px 0; background: #EBF4FF;
        border-radius: 0 6px 6px 0;
    }
    .insight-box {
        background: #EBF8FF; border-left: 4px solid #2D6BC4;
        border-radius: 0 8px 8px 0; padding: 12px 16px;
        margin: 8px 0; font-size: 0.9rem; color: #1B2A4A;
    }
    .warning-box {
        background: #FFFBEB; border-left: 4px solid #F5A623;
        border-radius: 0 8px 8px 0; padding: 12px 16px;
        margin: 8px 0; font-size: 0.9rem; color: #744210;
    }
    .danger-box {
        background: #FFF5F5; border-left: 4px solid #E53E3E;
        border-radius: 0 8px 8px 0; padding: 12px 16px;
        margin: 8px 0; font-size: 0.9rem; color: #742A2A;
    }
    .success-box {
        background: #F0FFF4; border-left: 4px solid #27AE60;
        border-radius: 0 8px 8px 0; padding: 12px 16px;
        margin: 8px 0; font-size: 0.9rem; color: #1A4731;
    }
    [data-testid="stDataFrame"] { border-radius: 8px; border: 1px solid #E2E8F0; }
    .stTextInput input { background-color: #FFFFFF; color: #1B2A4A; border: 1.5px solid #CBD5E0; border-radius: 7px; }
    .stSelectbox > div > div { background-color: #FFFFFF; border: 1.5px solid #CBD5E0; border-radius: 7px; }
    /* [디자인 3] 필터 multiselect 선택 태그 무채색 계열 */
    .stMultiSelect [data-baseweb="tag"] {
        background-color: #64748B !important;
        color: #FFFFFF !important;
        border-radius: 5px !important;
    }
    .stMultiSelect [data-baseweb="tag"] span { color: #FFFFFF !important; }
    .stMultiSelect [data-baseweb="tag"] svg { fill: #E2E8F0 !important; }
    [data-testid="stSidebar"] .stMultiSelect [data-baseweb="tag"] {
        background-color: #64748B !important;
    }
    hr { border-color: #E2E8F0; }
    footer { color: #A0AEC0 !important; }
</style>
""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════════
# [v2.6.2] 접근 제한 — 공유 비밀번호 관문(사업부 구성원 전용, 옵션 C)
#   · Streamlit Cloud → Manage app → Settings → Secrets 에 app_password 설정
#   · 세션당 1회 입력 후 통과(재입력 불필요). 미설정 시 경고 후 개방(잠금 해제).
# ════════════════════════════════════════════════════════════════════
def _require_login():
    if st.session_state.get("_pcid_authed"):
        return
    try:
        _pw = str(st.secrets["app_password"]) if "app_password" in st.secrets else ""
    except Exception:
        _pw = ""
    if not _pw:
        st.session_state["_pcid_authed"] = True
        st.warning("⚠ 접근 비밀번호(app_password)가 설정되지 않아 접근 제한이 비활성화되어 있습니다. "
                   "Streamlit Cloud → Manage app → Settings → Secrets 에 app_password 를 추가하세요.")
        return
    _l, _c, _r = st.columns([1, 1.5, 1])
    with _c:
        st.markdown("<div style='height:7vh'></div>", unsafe_allow_html=True)
        st.markdown("<div style='text-align:center;font-weight:800;font-size:1.5rem;color:#1B2A4A'>"
                    "PASCUCCI Consumer Intelligence Dashboard</div>", unsafe_allow_html=True)
        st.markdown("<div style='text-align:center;color:#71808F;margin:6px 0 18px'>"
                    "🔒 사업부 구성원 전용 — 접근 비밀번호를 입력하세요</div>", unsafe_allow_html=True)
        with st.form("_pcid_login", clear_on_submit=False):
            _pwin = st.text_input("비밀번호", type="password",
                                  label_visibility="collapsed", placeholder="접근 비밀번호")
            _ok = st.form_submit_button("입장", use_container_width=True)
        if _ok:
            if _pwin == _pw:
                st.session_state["_pcid_authed"] = True
                st.rerun()
            else:
                st.error("비밀번호가 올바르지 않습니다. 다시 시도해 주세요.")
        st.caption("문의: 사업부 대시보드 운영자")
    st.stop()


_require_login()

# ══════════════════════════════════════════════════════
# 공통 상수
# ══════════════════════════════════════════════════════
BRAND_COLORS = {
    "스타벅스":    "#E74C3C",
    "투썸플레이스": "#8E44AD",
    "메가커피":    "#27AE60",
    "컴포즈커피":  "#16A085",
    "이디야":      "#1D6A3A",
    "폴바셋":      "#0F6E56",
    "테라로사":    "#2980B9",
    "빽다방":      "#E67E22",
    "블루보틀":    "#D35400",
    "할리스":      "#7F8C8D",
    "파스쿠찌":    "#C0392B",
}
TARGET_BRANDS = ["파스쿠찌","스타벅스","투썸플레이스","메가커피","컴포즈커피","이디야","폴바셋","테라로사","빽다방","블루보틀","할리스"]
SENT_COLORS = {"긍정": "#27AE60", "중립": "#F39C12", "부정": "#E74C3C"}
TOPIC_LABELS = ["브랜드선택·대안","디카페인·건강","맛·원두품질","카페공간·작업","가격·가성비","불매·이슈","선물·추천"]
TOPIC_COLORS = ["#2980B9","#16A085","#27AE60","#E67E22","#8E44AD","#E74C3C","#F39C12"]

STOPWORDS = set(
    "이 가 을 를 은 는 에 의 도 로 고 그 하다 있다 되다 않다 이다 것 수 나 저 우리 그리고 하지만 "
    "정말 진짜 너무 좀 더 제 여기 거기 그냥 있어요 없어요 해요 했어요 같아요 입니다 합니다 그게 "
    "이게 에서 한테 보다 있고 하고 하면 그래서 아니고 아니라 있는데 많이 제가 지금 먹고 마시고 "
    "가서 먹는 먹어요 사람 사람이 다른 하나 이렇게 그렇게 요즘 분들 내가 아주 엄청 무슨 이거 "
    "저거 저도 나도 저는 나는 우리 댓글 같이 항상 어디 얼마 있는 하는 그런 이런 같은 근데 이제 "
    "그래도 해서 있어서 없는 없고 것도 사람들이 사람들 누가 보면 얼마나 하는데 없어서 가끔 정도 "
    "다시 바로 오늘 시간 감사 아니".split()
)

# [텍스트마이닝 2] 의미 추론이 어려운 부사·수식어 — 토큰엔 남기되
# 대응일치분석(CA)·토픽모델링(LDA)·키워드 집계 등 '지표 집계'에서만 제외.
# 이들 단어만으로는 사용 맥락을 알 수 없어 해석 신뢰도를 떨어뜨리기 때문.
MEANINGLESS_MODIFIERS = set(
    "많이 아니면 거의 조금 약간 매우 굉장히 엄청 되게 너무 진짜 정말 완전 그냥 좀 더 덜 "
    "아주 워낙 훨씬 가장 제일 무척 상당히 꽤 별로 그다지 다소 어느정도 대체로 대부분 "
    "보통 아마 혹시 만약 일단 우선 결국 그냥 막 살짝 슬쩍 딱 그저 단지 오직 마치 "
    "이런 저런 그런 어떤 무슨 이렇게 저렇게 그렇게 어떻게 왜 언제 어디 누가 "
    "그리고 그러나 하지만 그래서 그래도 또한 또 즉 따라서 게다가 다만 한편 "
    "있다 없다 하다 되다 같다 이다 보다 들다 나다 오다 가다 주다 받다".split()
)

ATTR_DICT = {
    "맛·품질":    ["맛있","맛나","달콤","진한","부드럽","식감","쓴맛","산미","향","원두","에스프레소"],
    "가격·가성비": ["가격","비싸","저렴","가성비","합리","인상","할인","쿠폰"],
    "공간·분위기": ["카페","분위기","인테리어","좌석","조용","넓은","공부","작업","콘센트","카공"],
    "서비스":     ["직원","서비스","친절","속도","주문","대기","청결","응대"],
    "접근성":     ["위치","매장","점포","테이크아웃","드라이브","배달","픽업"],
    "브랜드이미지": ["이미지","브랜드","선호","느낌","감성","프리미엄","불매"],
    "프로모션·혜택":["이벤트","프로모션","혜택","한정","쿠폰","멤버십","상품권","환불"],
    "MD·굿즈":     ["굿즈","md","엠디","다이어리","텀블러","머그","키링","리유저블","콜라보","컬래버","한정판","피규어","스티커"],
    "디카페인":   ["디카페인","디카페","카페인"],
}

RISK_PATTERNS = {
    "🔥 정치·이슈":  ["불매","일베","정용진","이마트","신세계","극우","광주","멸공"],
    "💳 환불·쿠폰":  ["환불","상품권","쿠폰","잔액","기프티콘","카드"],
    "💸 가격":      ["비싸","인상","가격","원가","올린"],
    "😞 품질 실망": ["맛없","실망","별로","최악","부족","아쉽"],
    "🚶 접근성":    ["없어서","품절","재고","대체","매장에"],
}

OCCASION_DICT = {
    "공부·업무": ["공부","작업","노트북","콘센트","카공","업무","혼자"],
    "만남·미팅": ["친구","약속","만남","미팅","모임","데이트","커플","같이","함께"],
    "출근·이동": ["출근","이동","테이크아웃","픽업","드라이브"],
    "선물·기념": ["선물","기념","생일","기프트","굿즈"],
    "배달·홈카페":["배달","집에서","홈카페","캡슐","드립백","인스턴트"],
}

# ── 제품 사전 (커피·음료·디저트·푸드 망라) [항목 4] ──
PRODUCT_DICT = {
    # 커피·음료
    "아메리카노":["아메리카노","아아","뜨아"], "카페라떼":["카페라떼","라떼"],
    "에스프레소":["에스프레소"], "콜드브루":["콜드브루","콜브"],
    "카푸치노":["카푸치노"], "디카페인":["디카페인","디카페"],
    "프라푸치노":["프라푸치노","프라푸"], "에이드":["에이드"],
    "스무디":["스무디"], "밀크티":["밀크티"], "말차라떼":["말차"],
    "녹차·티":["녹차","얼그레이","루이보스","호지차"], "초콜릿음료":["핫초코","초콜릿라떼"],
    # 디저트
    "케이크":["케이크","케익"], "마카롱":["마카롱"], "티라미수":["티라미수"],
    "타르트":["타르트","에그타르트"], "크로플":["크로플"], "쿠키":["쿠키"],
    "빙수":["빙수"], "푸딩":["푸딩"],
    "아이스크림":["아이스크림"], "젤라또":["젤라또","젤라토"],
    # 트렌디·시그니처 [v2.3.1]
    "두바이초콜릿":["두바이초콜릿","두바이","쪼꼬바"], "두쫀쿠":["두쫀쿠","두바이쿠키","쪼득쿠키"],
    "우베":["우베","우베라떼","자색고구마"], "피스타치오":["피스타치오","피스타치오라떼"],
    "흑임자":["흑임자","흑임자라떼"], "말차디저트":["말차케이크","말차티라미수"],
    # 베이커리·푸드
    "소금빵":["소금빵"], "크로와상":["크로와상","크로아상"],
    "베이글":["베이글"], "스콘":["스콘"], "머핀":["머핀"],
    "빵·베이커리":["베이커리","식빵","바게트"],
    "샌드위치":["샌드위치","파니니"], "샐러드":["샐러드"],
}

# ── 프로모션 사전 [항목 5] ──
PROMO_DICT = {
    "통신·카드할인": ["통신할인","카드할인","제휴할인","KT","SKT","LG유플","삼성카드","현대카드","멤버십할인"],
    "이벤트":       ["이벤트","응모","참여","추첨","경품"],
    "프로모션·행사": ["프로모션","행사","페스타","페어"],
    "1+1·증정":     ["1+1","원플러스원","증정","사은품"],
    "굿즈·MD":      ["굿즈","MD","엠디","다이어리","텀블러","프리퀀시","스티커","리유저블","머그","키링"],
    "기프티콘·상품권":["기프티콘","기프트콘","모바일상품권","모바일쿠폰","교환권","e쿠폰","상품권"],
    "페이 행사":     ["네이버페이","카카오페이","토스페이","페이백","간편결제","네이버페이","제로페이"],
    "쿠폰·적립":    ["쿠폰","적립","스탬프","포인트"],
    "신메뉴·출시":  ["신메뉴","출시","새로나온","리뉴얼","한정"],
}

# 촉진 반응 강도 [항목 5]
PROMO_REACTION = {
    "강한 긍정": (["개꿀","꿀이","필참","머스트","잇템","갓","혜자","대박","무조건","사야","존맛","꼭사"], 2.0),
    "긍정":      (["좋아","괜찮","유용","이득","득템","쏠쏠","갖고싶","사고싶","혹하"], 1.0),
    "부정":      (["별로","아쉽","쓸데","꽝","손해","실망","비싸기만","글쎄"], -1.0),
}

# ── 트렌드 카테고리 분류 [항목 10] ──
TREND_CATEGORY = {
    "커피·음료": ["디카페인","말차·녹차","콜드브루","오트밀크","흑당","유자·시트러스"],
    "푸드류":    ["소금빵","크로플","바스크치즈","티라미수","두바이초콜릿"],
    "식재료":    ["크림치즈","피스타치오","흑임자·전통"],
}

def chart_style(fig, height=420, showlegend=True):
    fig.update_layout(
        plot_bgcolor="#FFFFFF", paper_bgcolor="#F7F9FC",
        font=dict(color="#2D3748", size=12), height=height,
        showlegend=showlegend,
        margin=dict(l=50, r=30, t=50, b=45),
        title_font=dict(color="#1B2A4A", size=14, family="Arial"),
        legend=dict(bgcolor="#FFFFFF", bordercolor="#E2E8F0", borderwidth=1,
                    font=dict(size=11, color="#2D3748")),
        xaxis=dict(gridcolor="#EDF2F7", zerolinecolor="#CBD5E0",
                   tickfont=dict(color="#4A5568"), title_font=dict(color="#4A5568")),
        yaxis=dict(gridcolor="#EDF2F7", zerolinecolor="#CBD5E0",
                   tickfont=dict(color="#4A5568"), title_font=dict(color="#4A5568")),
    )
    return fig

# ══════════════════════════════════════════════════════
# 데이터 로딩 및 전처리
# ══════════════════════════════════════════════════════
@st.cache_data(show_spinner="데이터 전처리 중...")
def load_and_preprocess(uploaded_file):
    _name = getattr(uploaded_file, "name", str(uploaded_file))  # [v2.4.1] 경로 문자열 허용
    if _name.endswith((".csv", ".csv.gz", ".gz")):
        df = pd.read_csv(uploaded_file, encoding="utf-8-sig",
                         compression=("gzip" if _name.endswith(".gz") else "infer"))
    else:
        df = pd.read_excel(uploaded_file)

    df = df[df["doc_id"].str.match(r"^[a-f0-9]{12}$", na=False)].copy()

    coffee_kw = ["커피","카페","아메리카노","라떼","에스프레소","스타벅스","스벅","투썸","이디야",
                 "메가커피","컴포즈","빽다방","할리스","폴바셋","파스쿠찌","블루보틀","테라로사",
                 "더벤티","드립","콜드브루","디카페인","원두","핸드드립"]

    def get_comments(j):
        try:
            return " ".join(c.get("body","") for c in json.loads(j)) if pd.notna(j) else ""
        except:
            return ""

    df["comment_text"] = df["comments_json"].apply(get_comments)
    df["full_text"] = (df["title"].fillna("") + " " +
                       df["body"].fillna("") + " " +
                       df["comment_text"])

    df["is_rel"] = (
        df["full_text"].apply(lambda t: any(k in t for k in coffee_kw)) |
        (df["brand_mentions"].notna() & df["brand_mentions"].str.strip().ne(""))
    )
    v = df[df["is_rel"]].copy()

    brand_map = {
        "스타벅스":    ["스타벅스","스벅"],
        "투썸플레이스": ["투썸플레이스","투썸"],
        "메가커피":    ["메가커피"],
        "컴포즈커피":  ["컴포즈커피","컴포즈"],
        "이디야":      ["이디야"],
        "빽다방":      ["빽다방"],
        "할리스":      ["할리스"],
        "파스쿠찌":    ["파스쿠찌"],
        "폴바셋":      ["폴바셋"],
        "블루보틀":    ["블루보틀"],
        "테라로사":    ["테라로사"],
        "더벤티":      ["더벤티"],
    }

    pos_kw = ["맛있","좋아","좋다","추천","최고","만족","친절","합리","가성비","저렴","달콤","훌륭","깔끔","즐겨","단골","완벽"]
    neg_kw = ["맛없","별로","실망","불만","최악","비싸","불친절","느리","아쉽","안좋","싫","불매","후회","환불","절대","다시는"]

    def detect_brands(t):
        return [b for b, als in brand_map.items() if any(a in str(t) for a in als)]

    def detect_primary_brand(t):
        """본문에서 가장 많이 언급된 '주요 브랜드' 판별 (리스크 필터 정합성용)"""
        text = str(t)
        counts = {}
        for b, als in brand_map.items():
            c = sum(text.count(a) for a in als)
            if c > 0:
                counts[b] = c
        if not counts:
            return ""
        return max(counts, key=counts.get)

    def get_sentiment(t):
        p = sum(1 for k in pos_kw if k in str(t))
        n = sum(1 for k in neg_kw if k in str(t))
        return "긍정" if p > n else ("부정" if n > p else "중립")

    def tokenize(t):
        return [w for w in re.findall(r"[가-힣]{2,6}", str(t)) if w not in STOPWORDS]

    def get_occasion(t):
        for oc, kws in OCCASION_DICT.items():
            if any(k in str(t) for k in kws):
                return oc
        return "기타"

    v["brands"]     = v["full_text"].apply(detect_brands)
    v["primary_brand"] = v["full_text"].apply(detect_primary_brand)
    v["sentiment"]  = v["full_text"].apply(get_sentiment)
    v["tokens"]     = v["full_text"].apply(tokenize)
    v["tokens_str"] = v["tokens"].apply(" ".join)
    v["occasion"]   = v["full_text"].apply(get_occasion)
    v["created_at"] = pd.to_datetime(v["created_at"], errors="coerce")
    v["year_month"] = v["created_at"].dt.to_period("M").astype(str)

    return v

# ══════════════════════════════════════════════════════
# 분석 함수 모음
# ══════════════════════════════════════════════════════
FNB_CATS = {
    "디저트": ["디저트","케이크","마카롱","다쿠아즈","쿠키","타르트","푸딩","마들렌","휘낭시에","스콘","크로플","와플","도넛","빙수","젤라또","아이스크림","초코","초콜릿","티라미수","에그타르트","슈크림","약과","까눌레"],
    "베이커리": ["베이커리","크루아상","바게트","소금빵","식빵","단팥빵","앙버터","프레첼","페이스트리","크림빵","빵집","모닝빵","마늘빵"],
    "카페푸드": ["샌드위치","파니니","샐러드","베이글","브런치","토스트","핫도그","키슈","그릭요거트","에그마요"],
    "식재료·플레이버": ["피스타치오","두바이","흑임자","말차","녹차","쑥","인절미","생크림","오트밀","그래놀라","크림치즈","연유","단호박","고구마","얼그레이","바닐라","캐러멜","헤이즐넛","딸기","블루베리","흑당"],
}
_FNB_POS = ["맛있","좋아","좋다","추천","최고","만족","친절","합리","가성비","저렴","달콤","훌륭","깔끔","즐겨","단골","완벽","부드럽","촉촉","고소"]
_FNB_NEG = ["맛없","별로","실망","불만","최악","비싸","불친절","느리","아쉽","안좋","싫","불매","후회","환불","절대","다시는","퍽퍽","느끼"]
_COFFEE_KW = ["커피","카페","아메리카노","라떼","에스프레소","스타벅스","스벅","투썸","이디야","메가커피","컴포즈","빽다방","할리스","폴바셋","파스쿠찌","블루보틀","테라로사","더벤티","드립","콜드브루","디카페인","원두","핸드드립"]


@st.cache_data(show_spinner="인접 식음료 트렌드 준비 중...")
def build_adjacent_food_voc(uploaded_file):
    """커피·브랜드 무관 글 중 디저트·베이커리·카페푸드·식재료 VoC만 추출(트렌드 센싱용).
    브랜드 분석(v)과 완전 분리 — NSS·SOV 등에 절대 반영되지 않음."""
    if hasattr(uploaded_file, "seek"):
        try: uploaded_file.seek(0)
        except Exception: pass
    _name = getattr(uploaded_file, "name", str(uploaded_file))
    try:
        if _name.endswith((".csv", ".csv.gz", ".gz")):
            df = pd.read_csv(uploaded_file, encoding="utf-8-sig",
                             compression=("gzip" if _name.endswith(".gz") else "infer"))
        else:
            df = pd.read_excel(uploaded_file)
    except Exception:
        return pd.DataFrame()
    if "doc_id" not in df.columns:
        return pd.DataFrame()
    df = df[df["doc_id"].astype(str).str.match(r"^[a-f0-9]{12}$", na=False)].copy()

    def _gc(j):
        try:
            return " ".join(c.get("body", "") for c in json.loads(j)) if pd.notna(j) else ""
        except Exception:
            return ""
    df["comment_text"] = df["comments_json"].apply(_gc) if "comments_json" in df.columns else ""
    df["full_text"] = (df.get("title", pd.Series([""] * len(df))).fillna("").astype(str) + " " +
                       df.get("body", pd.Series([""] * len(df))).fillna("").astype(str) + " " +
                       df["comment_text"].astype(str))
    is_coffee = df["full_text"].apply(lambda t: any(k in t for k in _COFFEE_KW))
    if "brand_mentions" in df.columns:
        is_coffee = is_coffee | (df["brand_mentions"].notna() & df["brand_mentions"].astype(str).str.strip().ne(""))
    adj = df[~is_coffee].copy()
    adj["food_cats"] = adj["full_text"].apply(lambda t: [c for c, kws in FNB_CATS.items() if any(k in t for k in kws)])
    adj = adj[adj["food_cats"].apply(len) > 0].copy()

    def _senti(t):
        p = sum(1 for k in _FNB_POS if k in t); n = sum(1 for k in _FNB_NEG if k in t)
        return "긍정" if p > n else ("부정" if n > p else "중립")
    adj["sentiment"] = adj["full_text"].apply(_senti)
    adj["created_at"] = pd.to_datetime(adj["created_at"], errors="coerce")
    adj["year_month"] = adj["created_at"].dt.to_period("M").astype(str)
    return adj


@st.cache_data(show_spinner=False)
def compute_fnb_trends(v_adj):
    if v_adj is None or len(v_adj) == 0:
        return pd.DataFrame(), pd.DataFrame()
    months = sorted([m for m in v_adj["year_month"].dropna().unique() if m and m != "NaT"])
    recent = months[-1] if months else None
    base = months[:-1] if len(months) > 1 else months
    rows = []
    for cat, kws in FNB_CATS.items():
        for kw in kws:
            hit = v_adj["full_text"].str.contains(kw, regex=False, na=False)
            tot = int(hit.sum())
            if tot < 5:
                continue
            sub = v_adj[hit]
            rc = int((sub["year_month"] == recent).sum()) if recent else 0
            bc = sub[sub["year_month"].isin(base)].groupby("year_month").size() if base else pd.Series(dtype=int)
            base_avg = float(bc.mean()) if len(bc) else 0.0
            mom = round(rc / base_avg, 2) if base_avg > 0 else float(rc)
            pos = (sub["sentiment"] == "긍정").mean() * 100
            neg = (sub["sentiment"] == "부정").mean() * 100
            rows.append({"키워드": kw, "카테고리": cat, "총언급": tot, "최근월": rc,
                         "기준평균": round(base_avg, 1), "모멘텀": mom, "NSS": round(pos - neg, 1)})
    kw_df = (pd.DataFrame(rows).sort_values("모멘텀", ascending=False).reset_index(drop=True)
             if rows else pd.DataFrame())
    ex = v_adj.explode("food_cats")
    crows = []
    for cat, g in ex.groupby("food_cats"):
        pos = (g["sentiment"] == "긍정").mean() * 100
        neg = (g["sentiment"] == "부정").mean() * 100
        crows.append({"카테고리": cat, "건수": len(g), "긍정%": round(pos, 0),
                      "부정%": round(neg, 0), "NSS": round(pos - neg, 1)})
    cat_df = (pd.DataFrame(crows).sort_values("건수", ascending=False).reset_index(drop=True)
              if crows else pd.DataFrame())
    return kw_df, cat_df


@st.cache_data(show_spinner="브랜드 지표 계산 중...")
def compute_brand_stats(_v):
    total_docs = len(_v)
    pos_total  = (_v["sentiment"] == "긍정").sum()
    neg_total  = (_v["sentiment"] == "부정").sum()
    rows = []
    for b in TARGET_BRANDS:
        docs = _v[_v["brands"].apply(lambda x: b in x)]
        n = len(docs)
        if n < 3:
            continue
        pos = (docs["sentiment"] == "긍정").sum()
        neg = (docs["sentiment"] == "부정").sum()
        rows.append({
            "브랜드":      b,
            "언급량":      n,
            "긍정":        int(pos),
            "부정":        int(neg),
            "중립":        int(n - pos - neg),
            "긍정률":      round(pos / n * 100, 1),
            "부정률":      round(neg / n * 100, 1),
            "NSS":        round((pos - neg) / n * 100, 1),
            "SOV_전체":   round(n / total_docs * 100, 1),
            "SOV_긍정":   round(pos / pos_total * 100, 1) if pos_total > 0 else 0,
            "SOV_부정":   round(neg / neg_total * 100, 1) if neg_total > 0 else 0,
        })
    return pd.DataFrame(rows)

@st.cache_data(show_spinner="월별 NSS 계산 중...")
def compute_monthly_nss(_v):
    months = sorted(
        [m for m in _v["year_month"].dropna().unique()
         if "NaT" not in m and ("2025" in m or "2026" in m)]
    )
    main6 = ["스타벅스","투썸플레이스","메가커피","컴포즈커피","이디야","폴바셋"]
    rows = []
    for b in main6:
        bd = _v[_v["brands"].apply(lambda x: b in x)]
        for m in months:
            sub = bd[bd["year_month"] == m]
            if len(sub) < 2:
                continue
            p = (sub["sentiment"] == "긍정").sum()
            n = (sub["sentiment"] == "부정").sum()
            t = len(sub)
            rows.append({"브랜드": b, "월": m, "NSS": round((p - n) / t * 100, 1), "문서수": t})
    return pd.DataFrame(rows), months

@st.cache_data(show_spinner="속성 드라이버 분석 중...")
def compute_absa(_v):
    pos_kw = ["맛있","좋아","좋다","추천","최고","만족","친절","합리","가성비","저렴","달콤","훌륭"]
    neg_kw = ["맛없","별로","실망","불만","최악","비싸","불친절","느리","아쉽","안좋","싫","불매","후회","환불"]
    rows = []
    for attr, kws in ATTR_DICT.items():
        pat  = "|".join(kws)
        mask = _v["full_text"].str.contains(pat, na=False)
        sub  = _v[mask]
        if len(sub) < 3:
            continue
        s   = sub["sentiment"].value_counts()
        tot = len(sub)
        pos = s.get("긍정", 0)
        neg = s.get("부정", 0)
        rows.append({
            "속성":    attr, "문서수": tot,
            "긍정":   int(pos), "부정": int(neg),
            "긍정률": round(pos / tot * 100, 1),
            "부정률": round(neg / tot * 100, 1),
            "NSS":    round((pos - neg) / tot * 100, 1),
        })
    return pd.DataFrame(rows).sort_values("NSS", ascending=False)

@st.cache_data(show_spinner="LDA 토픽 모델링 중...")
def compute_lda(_v):
    docs = _v[_v["tokens_str"].str.len() > 5]["tokens_str"].tolist()
    if len(docs) < 20:
        return None, None, None, None, None
    tfidf = TfidfVectorizer(max_features=500, min_df=2, max_df=0.85,
                             token_pattern=r"[가-힣]{2,6}",
                             stop_words=list(MEANINGLESS_MODIFIERS))
    X     = tfidf.fit_transform(docs)
    vocab = tfidf.get_feature_names_out()
    lda   = LatentDirichletAllocation(n_components=7, random_state=42,
                                       max_iter=30, learning_method="batch")
    lda.fit(X)
    doc_t    = lda.transform(X)
    topic_dist = doc_t.mean(axis=0)
    return lda, tfidf, vocab, doc_t, topic_dist

CA_STOP = {
    "때문에","보고","있습니다","같습니다","합니다","입니다","였습니다","됩니다","했습니다","드립니다",
    "그리고","하지만","그래서","그런데","근데","해서","하고","에서","으로","까지","부터","보면","보니",
    "대한","대해","관련","위해","통해","따라","정도","경우","생각","약간","정말","너무","완전","많이","조금",
    "이라고","라고","처럼","만큼","에게","우리","저는","제가","그냥","이거","저거","요즘","내가",
    "커피를","커피는","커피가","맛이","마시","마셔","먹고","가서","있는","없는","같이","위해서",
    "라떼","아메리카노","음료","메뉴","그리고요","근데요","했는데","하는데","인데","그게","이게","저게"}
_CA_END = ("습니다", "됩니다", "입니다", "였습니다", "때문에", "이라고", "라고요",
           "하는데", "있는데", "면서", "거나", "지만", "어서요", "아서요")


def _ca_meaningful(w):
    if (w in STOPWORDS or w in MEANINGLESS_MODIFIERS or w in CA_STOP or len(w) < 2):
        return False
    if w.endswith(_CA_END):
        return False
    return True


@st.cache_data(show_spinner="대응일치분석 계산 중...")
def compute_ca(_v):
    excl = {"커피","카페","스벅","스타벅스","있는","하는","그런","이런","같은","근데","이제",
            "맛이","커피를","커피는","커피가","카페에서","라떼","있고","없는","없고","것도","해서","있어서"}
    all_toks = [t for ts in _v["tokens"] for t in ts]
    top_kws  = [w for w, _ in Counter(all_toks).most_common(260)
                if w not in excl and _ca_meaningful(w)][:40]
    ca_brands = [b for b in TARGET_BRANDS
                 if _v["brands"].apply(lambda x: b in x).sum() >= 5]
    if len(ca_brands) < 2:
        return None
    cont = pd.DataFrame(0, index=ca_brands, columns=top_kws)
    for b in ca_brands:
        bd  = _v[_v["brands"].apply(lambda x: b in x)]
        cnt = Counter([t for ts in bd["tokens"] for t in ts])
        for kw in top_kws:
            cont.loc[b, kw] = cnt.get(kw, 0)
    N_ca = cont.values.astype(float) + 0.5
    r_ca = N_ca.sum(axis=1); c_ca = N_ca.sum(axis=0); tot = N_ca.sum()
    P    = N_ca / tot; rv = r_ca / tot; cv = c_ca / tot
    Dr   = np.diag(1 / np.sqrt(rv)); Dc = np.diag(1 / np.sqrt(cv))
    S    = Dr @ (P - np.outer(rv, cv)) @ Dc
    U, sigma, Vt = la.svd(S, full_matrices=False)
    row_c = Dr @ U[:, :2] * sigma[:2]
    col_c = Dc @ Vt[:2, :].T
    expl  = sigma[:2] ** 2 / (sigma ** 2).sum() * 100
    return {"brands": ca_brands, "row_x": row_c[:, 0], "row_y": row_c[:, 1],
            "col_x": col_c[:, 0], "col_y": col_c[:, 1],
            "kws": top_kws, "expl": expl, "cont": cont}

@st.cache_data(show_spinner="Burst 탐지 중...")
def compute_burst(_v):
    months = sorted([m for m in _v["year_month"].dropna().unique()
                     if "NaT" not in m and ("2025" in m or "2026" in m)])
    recent_m = [m for m in months if "2026" in m][-2:]
    prev_m   = [m for m in months if "2025" in m][-3:]
    r_cnt, p_cnt = Counter(), Counter()
    for m in recent_m:
        sub = _v[_v["year_month"] == m]
        r_cnt.update([t for ts in sub["tokens"] for t in ts])
    for m in prev_m:
        sub = _v[_v["year_month"] == m]
        p_cnt.update([t for ts in sub["tokens"] for t in ts])
    rows = []
    for kw, cnt in r_cnt.most_common(500):
        if kw in STOPWORDS or len(kw) < 2:
            continue
        prev  = p_cnt.get(kw, 0)
        ratio = round(cnt / (prev + 1), 1)
        if cnt >= 5:
            rows.append({"키워드": kw, "최근": cnt, "이전": prev, "증가율(배)": ratio})
    return pd.DataFrame(rows).sort_values("증가율(배)", ascending=False).head(20)

@st.cache_data(show_spinner="PMI 계산 중...")
def compute_pmi(_v):
    N  = len(_v)
    bd = {b: set(_v[_v["brands"].apply(lambda x: b in x)].index) for b in TARGET_BRANDS}
    rows = []
    for i, b1 in enumerate(TARGET_BRANDS):
        for b2 in TARGET_BRANDS[i + 1:]:
            co  = len(bd.get(b1, set()) & bd.get(b2, set()))
            p1  = len(bd.get(b1, set())) / N
            p2  = len(bd.get(b2, set())) / N
            pco = co / N
            pmi = round(np.log2(pco / (p1 * p2) + 1e-9), 2) if pco > 0 else -99
            rows.append({"브랜드1": b1, "브랜드2": b2, "공출현수": co, "PMI": pmi})
    return pd.DataFrame(rows).sort_values("PMI", ascending=False)

# ══════════════════════════════════════════════════════
# 4대 고도화 분석 함수
# ══════════════════════════════════════════════════════

# ── 고도화 A: 브랜드 로열티(찐팬) 지수 ──────────────────
LOYALTY_SIGNALS = [
    ("독점 충성", ["여기 아니면", "무조건", "만 먹", "만 마셔", "만 가", "아니면 안", "밖에 안"], 3.0),
    ("강한 애착", ["최애", "푹 빠졌", "단골", "사랑", "진심", "믿고", "인생", "내 최애"], 2.0),
    ("반복 소비", ["자주", "매일", "또 갔", "또 가", "재방문", "늘 먹", "맨날", "항상"], 1.0),
]

@st.cache_data(show_spinner="브랜드 로열티 지수 계산 중...")
def compute_loyalty(_v):
    def loyalty_score_of_doc(text):
        t = str(text)
        for name, kws, w in LOYALTY_SIGNALS:
            if any(k in t for k in kws):
                return w
        return 0.0
    def loyalty_tier_of_doc(text):
        t = str(text)
        for name, kws, w in LOYALTY_SIGNALS:
            if any(k in t for k in kws):
                return name
        return "일반"
    rows = []
    for b in TARGET_BRANDS:  # [v2.5.1] 파스쿠찌 중복 제거(이미 포함)
        bd = _v[_v["brands"].apply(lambda x: b in x)]
        n = len(bd)
        if n < 3:
            continue
        scores = bd["full_text"].apply(loyalty_score_of_doc)
        tiers  = bd["full_text"].apply(loyalty_tier_of_doc)
        fan_docs = int((scores > 0).sum())
        loyalty_idx = round(scores.mean() / 3 * 100, 1)
        fan_ratio   = round(fan_docs / n * 100, 1)
        tier_counts = tiers.value_counts()
        # NSS도 함께
        pos = (bd["sentiment"] == "긍정").sum()
        neg = (bd["sentiment"] == "부정").sum()
        nss = round((pos - neg) / n * 100, 1)
        rows.append({
            "브랜드": b, "언급량": n,
            "로열티지수": loyalty_idx,
            "찐팬비율": fan_ratio,
            "찐팬문서수": fan_docs,
            "독점충성": int(tier_counts.get("독점 충성", 0)),
            "강한애착": int(tier_counts.get("강한 애착", 0)),
            "반복소비": int(tier_counts.get("반복 소비", 0)),
            "NSS": nss,
        })
    return pd.DataFrame(rows).drop_duplicates("브랜드").sort_values("로열티지수", ascending=False)


# ── 고도화 B: 소비맥락 세분화 (15유형) ──────────────────
OCCASION_DETAILED = {
    "집중 카공":     ["카공", "공부", "노트북", "콘센트", "집중", "과제"],
    "그룹 스터디":   ["스터디", "팀플", "모임공부", "같이공부"],
    "업무 미팅":     ["미팅", "회의", "업무", "비즈니스", "외근"],
    "친구 수다":     ["친구", "수다", "만나", "웃긴", "친구들"],
    "데이트":        ["데이트", "커플", "남친", "여친", "기념일"],
    "가족 나들이":   ["가족", "부모님", "아이", "엄마", "아빠", "애들"],
    "혼카페 힐링":   ["혼자", "혼카", "힐링", "멍때", "여유"],
    "책읽기·사색":   ["책", "독서", "읽기"],
    "아침·출근":     ["출근", "아침", "모닝", "등굣"],
    "테이크아웃":    ["테이크아웃", "포장", "픽업", "들고"],
    "드라이브":      ["드라이브", "차에서", "드스루"],
    "디저트 타임":   ["디저트", "케이크", "달달", "당충전"],
    "선물·기프티콘": ["선물", "기프티콘", "상품권", "기프트"],
    "배달·홈카페":   ["배달", "집에서", "홈카페", "캡슐"],
    "여행·나들이":   ["여행", "나들이", "바다", "놀러"],
}

@st.cache_data(show_spinner="소비맥락 세분화 분석 중...")
def compute_occasion_detailed(_v):
    rows = []
    for oc, kws in OCCASION_DETAILED.items():
        pat = "|".join(kws)
        mask = _v["full_text"].str.contains(pat, na=False)
        sub  = _v[mask]
        n = len(sub)
        if n < 2:
            continue
        pos = (sub["sentiment"] == "긍정").sum()
        neg = (sub["sentiment"] == "부정").sum()
        nss = round((pos - neg) / n * 100, 1)
        rows.append({"소비맥락": oc, "언급량": int(n), "긍정률": round(pos/n*100,1),
                     "NSS": nss})
    return pd.DataFrame(rows).sort_values("언급량", ascending=False)

@st.cache_data(show_spinner="맥락 × 브랜드 교차 분석 중...")
def compute_occasion_brand_matrix(_v):
    occ_list = list(OCCASION_DETAILED.keys())
    mat = pd.DataFrame(0, index=occ_list, columns=TARGET_BRANDS)
    for oc, kws in OCCASION_DETAILED.items():
        pat = "|".join(kws)
        oc_docs = _v[_v["full_text"].str.contains(pat, na=False)]
        for b in TARGET_BRANDS:
            mat.loc[oc, b] = int(oc_docs["brands"].apply(lambda x: b in x).sum())
    return mat


# ── 고도화 C: LDA 토픽 클러스터링 (토픽×브랜드 중첩) ─────
@st.cache_data(show_spinner="토픽 클러스터링 중...")
def compute_topic_clustering(_v):
    sub = _v[_v["tokens_str"].str.len() > 5].copy().reset_index(drop=True)
    if len(sub) < 20:
        return None
    tfidf = TfidfVectorizer(max_features=300, min_df=2, max_df=0.85,
                             token_pattern=r"[가-힣]{2,6}",
                             stop_words=list(MEANINGLESS_MODIFIERS))
    X = tfidf.fit_transform(sub["tokens_str"])
    vocab = tfidf.get_feature_names_out()
    n_topics = 6  # [v2.3.1] 4~5개 → 6개 세분화
    lda = LatentDirichletAllocation(n_components=n_topics, random_state=42, max_iter=20)
    doc_topics = lda.fit_transform(X)
    sub["main_topic"] = doc_topics.argmax(axis=1)

    topic_words = []
    for t in range(n_topics):
        top_idx = lda.components_[t].argsort()[-8:][::-1]
        topic_words.append([vocab[i] for i in top_idx])

    # 토픽 × 브랜드 매트릭스
    tb_mat = pd.DataFrame(0, index=[f"T{i}" for i in range(n_topics)], columns=TARGET_BRANDS)
    for t in range(n_topics):
        t_docs = sub[sub["main_topic"] == t]
        for b in TARGET_BRANDS:
            tb_mat.loc[f"T{t}", b] = int(t_docs["brands"].apply(lambda x: b in x).sum())

    # 토픽 × 감성
    ts_mat = pd.DataFrame(0, index=[f"T{i}" for i in range(n_topics)], columns=["긍정","중립","부정"])
    for t in range(n_topics):
        t_docs = sub[sub["main_topic"] == t]
        for s in ["긍정","중립","부정"]:
            ts_mat.loc[f"T{t}", s] = int((t_docs["sentiment"] == s).sum())

    # [항목 6] 토픽 의미 자동 정의 (키워드 기반 규칙)
    def define_topic_meaning(words):
        ws = " ".join(words)
        if any(k in ws for k in ["디카페인","디카페","카페인","건강","임산부"]):
            return "건강·디카페인 니즈"
        if any(k in ws for k in ["가격","비싸","불매","환불","쿠폰","할인"]):
            return "가격·이슈 담론"
        if any(k in ws for k in ["맛있","맛","원두","에스프레소","향","진한"]):
            return "맛·품질 평가"
        if any(k in ws for k in ["공부","카공","노트북","작업","공간","분위기"]):
            return "카페 공간·작업"
        if any(k in ws for k in ["추천","후기","어디","비교","고민"]):
            return "브랜드 탐색·비교"
        if any(k in ws for k in ["케이크","디저트","빵","마카롱","베이커리"]):
            return "디저트·베이커리"
        if any(k in ws for k in ["선물","기프티콘","생일","굿즈"]):
            return "선물·기프트"
        return "일상 카페 소비"

    # [항목 6] 토픽별 주요 소비자 유형(소비맥락 기반)
    topic_meanings = [define_topic_meaning(topic_words[t]) for t in range(n_topics)]
    topic_occasions = []
    for t in range(n_topics):
        t_docs = sub[sub["main_topic"] == t]
        if "occasion" in t_docs.columns and len(t_docs) > 0:
            occ_top = t_docs["occasion"].value_counts().head(2).index.tolist()
        else:
            occ_top = []
        topic_occasions.append(occ_top)

    # 토픽별 대표 브랜드
    topic_brands = []
    for t in range(n_topics):
        t_docs = sub[sub["main_topic"] == t]
        tb = Counter(b for bs in t_docs["brands"] for b in bs)
        topic_brands.append([b for b, _ in tb.most_common(3)])

    return {
        "n_topics": n_topics, "topic_words": topic_words,
        "doc_topics": doc_topics, "sub": sub,
        "tb_mat": tb_mat, "ts_mat": ts_mat,
        "topic_sizes": [int((sub["main_topic"]==t).sum()) for t in range(n_topics)],
        "topic_meanings": topic_meanings,
        "topic_occasions": topic_occasions,
        "topic_brands": topic_brands,
    }


# ── 고도화 D: 상대평가 트렌드 조기경보 ──────────────────
BASELINE_PRODUCTS = ["아메리카노", "라떼", "아이스아메리카노", "아아"]
EMERGING_CANDIDATES = {
    "디카페인": ["디카페인", "디카페"],
    "말차·녹차": ["말차", "녹차라떼"],
    "크림치즈": ["크림치즈"],
    "소금빵": ["소금빵"],
    "피스타치오": ["피스타치오", "피스타"],
    "오트밀크": ["오트밀크", "오트라떼", "귀리"],
    "흑임자·전통": ["흑임자", "미숫가루", "콩가루"],
    "콜드브루": ["콜드브루", "콜브"],
    "티라미수": ["티라미수"],
    "두바이초콜릿": ["두바이", "카다이프"],
    "크로플": ["크로플"],
    "유자·시트러스": ["유자", "레몬"],
    "바스크치즈": ["바스크"],
    "흑당": ["흑당", "버블티"],
}

@st.cache_data(show_spinner="트렌드 조기경보 계산 중...")
def compute_trend_earlywarning(_v):
    months = sorted([m for m in _v["year_month"].dropna().unique()
                     if "NaT" not in str(m) and ("2025" in str(m) or "2026" in str(m))])
    if len(months) < 3:
        return None

    # 기준선: 안정 제품군 월평균 언급량
    baseline_monthly = []
    for m in months:
        sub = _v[_v["year_month"] == m]
        cnt = sub["full_text"].apply(lambda t: any(b in str(t) for b in BASELINE_PRODUCTS)).sum()
        baseline_monthly.append(int(cnt))
    baseline_avg = np.mean(baseline_monthly) if baseline_monthly else 1
    if baseline_avg < 1:
        baseline_avg = 1

    rows = []
    for prod, kws in EMERGING_CANDIDATES.items():
        pat = "|".join(kws)
        monthly = []
        for m in months:
            sub = _v[_v["year_month"] == m]
            monthly.append(int(sub["full_text"].str.contains(pat, na=False).sum()))
        total = sum(monthly)
        if total < 2:
            continue
        recent = np.mean(monthly[-2:]) if len(monthly) >= 2 else 0
        rel_ratio = round(recent / baseline_avg, 2)
        # 가속도 (최근 4개월 기울기)
        if len(monthly) >= 3:
            tail = monthly[-4:] if len(monthly) >= 4 else monthly
            x = np.arange(len(tail))
            slope = round(float(np.polyfit(x, tail, 1)[0]), 2) if len(set(tail)) > 1 else 0.0
        else:
            slope = 0.0
        # 브랜드 확산도
        prod_docs = _v[_v["full_text"].str.contains(pat, na=False)]
        brand_spread = len(set(b for bs in prod_docs["brands"] for b in bs))
        # 감성
        pos = (prod_docs["sentiment"] == "긍정").sum()
        neg = (prod_docs["sentiment"] == "부정").sum()
        nss = round((pos - neg) / max(len(prod_docs),1) * 100, 1)
        # Emerging Signal Score
        es = round(rel_ratio * 10 + slope * 5 + brand_spread * 2, 1)
        # 경보 등급
        if rel_ratio >= 0.8 and slope > 0:
            alert = "🚨 강한 신호"
        elif slope > 0.5:
            alert = "📈 성장 조짐"
        elif slope < -0.5:
            alert = "📉 둔화"
        else:
            alert = "➡️ 안정"
        rows.append({
            "제품·식재료": prod, "총언급": total,
            "상대배수": rel_ratio, "성장기울기": slope,
            "브랜드확산": brand_spread, "NSS": nss,
            "ES점수": es, "경보": alert, "_monthly": monthly,
        })
    result_df = pd.DataFrame(rows).sort_values("ES점수", ascending=False)
    # 트렌드 카테고리 매핑 [항목 10]
    prod_to_cat = {}
    for cat, items in TREND_CATEGORY.items():
        for it in items:
            prod_to_cat[it] = cat
    result_df["카테고리"] = result_df["제품·식재료"].map(prod_to_cat).fillna("기타")
    return {"df": result_df, "months": months, "baseline_avg": round(baseline_avg,1),
            "baseline_monthly": baseline_monthly}


# ── 고도화: 브랜드별 제품 언급·감성 분석 [항목 4] ──────
@st.cache_data(show_spinner="브랜드별 제품 분석 중...")
def compute_brand_products(_v, brand):
    bdocs = _v[_v["brands"].apply(lambda x: brand in x)]
    if len(bdocs) < 2:
        return pd.DataFrame()
    rows = []
    for prod, kws in PRODUCT_DICT.items():
        pat = "|".join(kws)
        sub = bdocs[bdocs["full_text"].str.contains(pat, na=False)]
        n = len(sub)
        if n < 1:
            continue
        pos = (sub["sentiment"] == "긍정").sum()
        neg = (sub["sentiment"] == "부정").sum()
        nss = round((pos - neg) / n * 100, 1)
        # 카테고리 분류
        if prod in ["아메리카노","카페라떼","에스프레소","콜드브루","카푸치노","디카페인",
                    "프라푸치노","에이드","스무디","밀크티","말차라떼","녹차·티","초콜릿음료"]:
            cat = "커피·음료"
        elif prod in ["케이크","마카롱","티라미수","타르트","크로플","쿠키","빙수","푸딩","아이스크림","젤라또"]:
            cat = "디저트"
        else:
            cat = "베이커리·푸드"
        rows.append({"제품": prod, "카테고리": cat, "언급량": int(n),
                     "긍정": int(pos), "부정": int(neg), "NSS": nss})
    return pd.DataFrame(rows).sort_values("언급량", ascending=False)

@st.cache_data(show_spinner="제품 월별 추이 분석 중...")
def compute_product_trend(_v, brand, product):
    bdocs = _v[_v["brands"].apply(lambda x: brand in x)]
    kws = PRODUCT_DICT.get(product, [product])
    pat = "|".join(kws)
    sub = bdocs[bdocs["full_text"].str.contains(pat, na=False)]
    months = sorted([m for m in _v["year_month"].dropna().unique() if "NaT" not in str(m)])
    rows = []
    for m in months:
        msub = sub[sub["year_month"] == m]
        n = len(msub)
        pos = (msub["sentiment"] == "긍정").sum()
        neg = (msub["sentiment"] == "부정").sum()
        nss = round((pos - neg) / n * 100, 1) if n > 0 else 0
        rows.append({"월": m, "언급량": int(n), "NSS": nss})
    return pd.DataFrame(rows)


# ── 고도화: 프로모션 모니터링 [항목 5] ─────────────────
@st.cache_data(show_spinner="프로모션 모니터링 중...")
def compute_promotion(_v):
    # 프로모션 유형별 × 브랜드 언급
    promo_brand = pd.DataFrame(0, index=list(PROMO_DICT.keys()), columns=TARGET_BRANDS)
    promo_total = {}
    for promo, kws in PROMO_DICT.items():
        pat = "|".join(kws)
        pdocs = _v[_v["full_text"].str.contains(pat, na=False, case=False)]
        promo_total[promo] = int(len(pdocs))
        for b in TARGET_BRANDS:
            promo_brand.loc[promo, b] = int(pdocs["brands"].apply(lambda x: b in x).sum())
    promo_total_df = pd.DataFrame(
        [{"프로모션": k, "언급량": v} for k, v in promo_total.items()]
    ).sort_values("언급량", ascending=False)
    return promo_brand, promo_total_df

@st.cache_data(show_spinner="최근 프로모션 추세 분석 중...")
def compute_promo_recent(_v):
    """[기능 4] 최근(마지막 2개월) 가장 많이 언급된 프로모션 + 가장 열심히 하는 브랜드"""
    months = sorted([m for m in _v["year_month"].dropna().unique() if "NaT" not in str(m)])
    if not months:
        return pd.DataFrame(), pd.DataFrame()
    recent_months = months[-2:] if len(months) >= 2 else months
    recent = _v[_v["year_month"].isin(recent_months)]
    rows = []
    for promo, kws in PROMO_DICT.items():
        pat = "|".join(kws)
        rdocs = recent[recent["full_text"].str.contains(pat, na=False, case=False)]
        n = len(rdocs)
        if n < 1:
            continue
        # 가장 열심히 하는(많이 언급된) 브랜드
        bc = Counter()
        for _, r in rdocs.iterrows():
            for b in r["brands"]:
                bc[b] += 1
        top_brand = bc.most_common(1)[0][0] if bc else "—"
        top_brand_cnt = bc.most_common(1)[0][1] if bc else 0
        rows.append({"프로모션": promo, "최근언급": int(n),
                     "주력브랜드": top_brand, "주력브랜드_언급": int(top_brand_cnt)})
    recent_df = pd.DataFrame(rows).sort_values("최근언급", ascending=False)

    # 브랜드별 프로모션 적극도 (전체 기간)
    brand_promo_rows = []
    for b in TARGET_BRANDS:
        bdocs = _v[_v["brands"].apply(lambda x: b in x)]
        if len(bdocs) < 1:
            continue
        all_promo_kws = [kw for kws in PROMO_DICT.values() for kw in kws]
        ppat = "|".join(all_promo_kws)
        promo_cnt = int(bdocs["full_text"].str.contains(ppat, na=False, case=False).sum())
        ratio = round(promo_cnt / len(bdocs) * 100, 1)
        brand_promo_rows.append({"브랜드": b, "프로모션언급": promo_cnt,
                                  "전체언급": int(len(bdocs)), "프로모션비중": ratio})
    brand_promo_df = pd.DataFrame(brand_promo_rows).sort_values("프로모션언급", ascending=False)
    return recent_df, brand_promo_df


@st.cache_data(show_spinner="급증 프로모션 탐지 중...")
def compute_promo_surge(_v):
    """[v2.3.1] 최근 급증 프로모션 + 프로모션×브랜드 cross 매트릭스"""
    months = sorted([m for m in _v["year_month"].dropna().unique() if "NaT" not in str(m)])
    if len(months) < 2:
        return pd.DataFrame(), pd.DataFrame()
    recent_months = months[-1:]
    prior_months = months[-4:-1] if len(months) >= 4 else months[:-1]
    recent = _v[_v["year_month"].isin(recent_months)]
    prior = _v[_v["year_month"].isin(prior_months)]
    rows = []
    for promo, kws in PROMO_DICT.items():
        pat = "|".join(kws)
        r = int(recent["full_text"].str.contains(pat, na=False, case=False).sum())
        p_total = int(prior["full_text"].str.contains(pat, na=False, case=False).sum())
        p_avg = p_total / max(len(prior_months), 1)
        growth = (r / p_avg - 1) * 100 if p_avg > 0 else (100.0 if r > 0 else 0.0)
        rows.append({"프로모션": promo, "최근월": r, "직전월평균": round(p_avg, 1),
                     "증감률": round(growth, 1), "급증": "🔥" if growth >= 50 and r >= 2 else ""})
    surge_df = pd.DataFrame(rows).sort_values("증감률", ascending=False)
    recent2 = _v[_v["year_month"].isin(months[-2:])]
    cross = pd.DataFrame(0, index=list(PROMO_DICT.keys()), columns=TARGET_BRANDS)
    for promo, kws in PROMO_DICT.items():
        pat = "|".join(kws)
        pdocs = recent2[recent2["full_text"].str.contains(pat, na=False, case=False)]
        for _, rr in pdocs.iterrows():
            for b in rr["brands"]:
                if b in cross.columns:
                    cross.loc[promo, b] += 1
    return surge_df, cross

@st.cache_data(show_spinner="촉진 반응강도 측정 중...")
def compute_promo_reaction(_v):
    """[기능 4] 반응 강도를 0~100 지표로 정규화하여 관리"""
    all_promo_kws = [kw for kws in PROMO_DICT.values() for kw in kws]
    promo_pat = "|".join(all_promo_kws)
    promo_docs = _v[_v["full_text"].str.contains(promo_pat, na=False, case=False)]
    rows = []
    for reaction, (kws, weight) in PROMO_REACTION.items():
        pat = "|".join(kws)
        cnt = int(promo_docs["full_text"].str.contains(pat, na=False).sum())
        rows.append({"반응": reaction, "언급량": cnt, "가중치": weight})
    react_df = pd.DataFrame(rows)
    # 프로모션 유형별 반응 강도 — 원점수 + 0~100 정규화 지표
    type_reaction = []
    for promo, pkws in PROMO_DICT.items():
        ppat = "|".join(pkws)
        pdocs = _v[_v["full_text"].str.contains(ppat, na=False, case=False)]
        if len(pdocs) < 1:
            continue
        pos_strong = pdocs["full_text"].str.contains("|".join(PROMO_REACTION["강한 긍정"][0]), na=False).sum()
        pos = pdocs["full_text"].str.contains("|".join(PROMO_REACTION["긍정"][0]), na=False).sum()
        neg = pdocs["full_text"].str.contains("|".join(PROMO_REACTION["부정"][0]), na=False).sum()
        raw_score = (pos_strong * 2 + pos * 1 - neg * 1) / len(pdocs)
        type_reaction.append({"프로모션": promo, "반응강도_raw": round(raw_score, 2),
                               "강한긍정": int(pos_strong), "긍정": int(pos), "부정": int(neg),
                               "언급량": int(len(pdocs))})
    type_react_df = pd.DataFrame(type_reaction)
    # 0~100 정규화 (반응강도 지표화)
    if len(type_react_df) > 0:
        rmin, rmax = type_react_df["반응강도_raw"].min(), type_react_df["반응강도_raw"].max()
        if rmax > rmin:
            type_react_df["반응강도지표"] = ((type_react_df["반응강도_raw"] - rmin) / (rmax - rmin) * 100).round(1)
        else:
            type_react_df["반응강도지표"] = 50.0
        type_react_df = type_react_df.sort_values("반응강도지표", ascending=False)
    return react_df, type_react_df, int(len(promo_docs))

@st.cache_data(show_spinner="프로모션 원문 추출 중...")
def _kw_sentence(text, kws):
    """[v2.5.2] 키워드가 포함된 문장만 발췌 — 전체 글 대신 관련 부분만 노출(공용)."""
    t = str(text)
    for s in re.split(r"(?<=[.!?。…])\s+|\n+", t):
        if any(k in s for k in kws):
            s2 = re.sub(r"\s+", " ", s).strip()
            if len(s2) >= 6:
                return s2[:180]
    for k in kws:
        i = t.find(k)
        if i >= 0:
            return re.sub(r"\s+", " ", t[max(0, i - 40):i + 110]).strip()[:180]
    return ""


def get_promo_detail_docs(_v, promo_type):
    """[기능 5/v2.5.2] 프로모션 유형 원문 — 키워드가 실제 든 '문장'을 발췌하고
    무관 글은 제외, 관련도(키워드 빈도+제목 포함) 순 정렬."""
    kws = PROMO_DICT.get(promo_type, [])
    if not kws:
        return _v.iloc[0:0].copy()
    pat = "|".join(map(re.escape, kws))
    pdocs = _v[_v["full_text"].str.contains(pat, na=False, case=False)].copy()
    if len(pdocs) == 0:
        return pdocs
    pdocs["발췌"] = pdocs["full_text"].apply(lambda t: _kw_sentence(t, kws))
    pdocs = pdocs[pdocs["발췌"].str.len() > 0].copy()
    if len(pdocs) == 0:
        return pdocs

    def _rel(r):
        ft = str(r["full_text"]); ti = str(r.get("title", ""))
        return sum(ft.count(k) for k in kws) + (3 if any(k in ti for k in kws) else 0)
    pdocs["_rel"] = pdocs.apply(_rel, axis=1)
    return pdocs.sort_values("_rel", ascending=False)


# ── 고도화: 고객 클러스터링 & 페르소나 [항목 7] ─────────
# [v2.3.1] 라이프스타일·관심도 추론 사전 / 커피 외 동시언급 브랜드
LIFESTYLE_DICT = {
    "맛집·미식 탐방": ["맛집","미식","존맛","빵지순례","디저트투어","베이커리투어"],
    "여행·나들이":   ["여행","나들이","드라이브","호캉스","카페투어","제주","바다"],
    "예술·문화 관여": ["전시","공연","영화","뮤지컬","미술관","독서","콘서트","연극"],
    "운동·건강":     ["운동","헬스","러닝","요가","필라테스","다이어트","등산"],
    "육아·가족":     ["육아","아기","엄마","가족","키즈","아이랑"],
    "재택·디지털":   ["재택","카공","작업","노트북","유튜브","넷플릭스"],
    "패션·뷰티":     ["패션","화장품","뷰티","올리브영","무신사"],
}
COMPANION_BRANDS = ["배스킨라빈스","파리바게뜨","뚜레쥬르","맥도날드","버거킹","롯데리아",
                    "올리브영","다이소","무신사","CGV","넷플릭스","쿠팡","이마트","코스트코","애플"]
@st.cache_data(show_spinner="고객 클러스터링 중...")
def compute_customer_clustering(_v, k=8):
    sub = _v[_v["tokens_str"].str.len() > 5].copy().reset_index(drop=True)
    if len(sub) < 30:
        return None
    # 피처: TF-IDF 텍스트
    tfidf = TfidfVectorizer(max_features=200, min_df=3, max_df=0.85,
                             token_pattern=r"[가-힣]{2,6}",
                             stop_words=list(MEANINGLESS_MODIFIERS))
    X_text = tfidf.fit_transform(sub["tokens_str"]).toarray()
    # 행동 피처: 감성, 소비상황, 브랜드 수
    sent_map = {"긍정": 1, "중립": 0, "부정": -1}
    feat_sent = sub["sentiment"].map(sent_map).fillna(0).values.reshape(-1, 1)
    feat_nbrand = sub["brands"].apply(len).values.reshape(-1, 1)
    feat_loyalty = sub["full_text"].apply(
        lambda t: 1 if any(k in str(t) for k in ["단골","최애","자주","매일","항상"]) else 0
    ).values.reshape(-1, 1)
    feat_decaf = sub["full_text"].apply(
        lambda t: 1 if any(k in str(t) for k in ["디카페인","디카페","카페인"]) else 0
    ).values.reshape(-1, 1)
    # 결합 (텍스트 70% + 행동 30%)
    behav = np.hstack([feat_sent, feat_nbrand, feat_loyalty, feat_decaf]).astype(float)
    behav = StandardScaler().fit_transform(behav)
    X = np.hstack([X_text * 0.7, behav * 0.3])

    km = KMeans(n_clusters=k, random_state=42, n_init=10)
    sub["cluster"] = km.fit_predict(X)
    # PCA 2D
    pca = PCA(n_components=2, random_state=42)
    coords = pca.fit_transform(X)
    sub["pca_x"] = coords[:, 0]
    sub["pca_y"] = coords[:, 1]

    # 클러스터별 프로파일
    profiles = []
    vocab = tfidf.get_feature_names_out()
    for c in range(k):
        cdocs = sub[sub["cluster"] == c]
        n = len(cdocs)
        if n == 0:
            continue
        # 대표 키워드
        ctoks = [t for ts in cdocs["tokens"] for t in ts]
        top_kw = [w for w, _ in Counter(ctoks).most_common(8)]
        # 주요 브랜드
        cbrands = Counter(b for bs in cdocs["brands"] for b in bs)
        top_brands = [b for b, _ in cbrands.most_common(3)]
        # 감성
        pos_r = round((cdocs["sentiment"] == "긍정").mean() * 100, 1)
        neg_r = round((cdocs["sentiment"] == "부정").mean() * 100, 1)
        nss = round(pos_r - neg_r, 1)
        # 충성도
        loyalty_r = round(cdocs["full_text"].apply(
            lambda t: any(k in str(t) for k in ["단골","최애","자주","매일"])).mean() * 100, 1)
        # 디카페인 관심
        decaf_r = round(cdocs["full_text"].apply(
            lambda t: any(k in str(t) for k in ["디카페인","디카페"])).mean() * 100, 1)
        # 주요 소비상황
        occ_top = cdocs["occasion"].value_counts().head(2).index.tolist() if "occasion" in cdocs.columns else []

        # [v2.3.1] 라이프스타일·관심도·동행 브랜드 추론
        ftext = " ".join(cdocs["full_text"].astype(str))
        life_scored = [(cat, sum(ftext.count(k) for k in kws)) for cat, kws in LIFESTYLE_DICT.items()]
        lifestyle = [cat for cat, sc in sorted(life_scored, key=lambda x: -x[1]) if sc > 0][:3]
        co_scored = [(b, ftext.count(b)) for b in COMPANION_BRANDS]
        co_brands = [b for b, sc in sorted(co_scored, key=lambda x: -x[1]) if sc > 0][:4]
        profiles.append({
            "cluster": c, "n": n, "top_kw": top_kw, "top_brands": top_brands,
            "pos_r": pos_r, "neg_r": neg_r, "nss": nss,
            "loyalty_r": loyalty_r, "decaf_r": decaf_r, "occ_top": occ_top,
            "lifestyle": lifestyle, "co_brands": co_brands,
        })
    # [v2.3.1] 브랜드 PCA 센트로이드 — 세그먼트 인접 브랜드 포지셔닝용
    brand_centroids = {}
    for b in TARGET_BRANDS:
        bd = sub[sub["brands"].apply(lambda x: b in x)]
        if len(bd) >= 5:
            brand_centroids[b] = (float(bd["pca_x"].mean()), float(bd["pca_y"].mean()))
    return {"sub": sub, "profiles": profiles, "k": k, "brand_centroids": brand_centroids,
            "pca_var": [round(x*100,1) for x in pca.explained_variance_ratio_]}


KW_NAV_STOP = {
    "다음","이전","목록","처음","마지막","검색","클릭","링크","답글","댓글","대댓글","원문","출처",
    "사진","이미지","첨부","파일","수정","삭제","신고","공유","스크랩","조회","추천수","비추","페이지",
    "로그인","회원","가입","커피를","커피가","커피는","커피도","커피랑","커피만",
    "때문에","그리고","하지만","그래서","근데","해서","에서","으로","까지","부터","정도","경우","생각",
    "약간","정말","너무","완전","많이","조금","대한","관련","위해","보고","같이"}


def _kw_clean(words, n=12):
    """[v2.5.2] 관심 키워드에서 조사·종결어미·네비게이션 등 의미 없는 토큰 제거(공용)."""
    out = []
    for w in words:
        if w in KW_NAV_STOP or w in STOPWORDS or len(w) < 2:
            continue
        if w.endswith(("습니다", "입니다", "됩니다", "이라고", "하는데", "있는데")):
            continue
        if w not in out:
            out.append(w)
        if len(out) >= n:
            break
    return out


def make_persona(profile):
    """[v2.5.2] 클러스터 프로파일 → 직관적·간결한 페르소나 네이밍 + 키워드 정제"""
    kw = _kw_clean(profile["top_kw"], 12)
    brands = profile["top_brands"]
    nss = profile["nss"]
    loyalty = profile["loyalty_r"]
    decaf = profile["decaf_r"]
    lifestyle = profile.get("lifestyle", [])
    life0 = lifestyle[0] if lifestyle else ""

    # 라이프스타일 → (짧은 한 단어 태그, 맞춤 마케팅 액티비티)
    LIFE_TAG = {
        "맛집·미식 탐방": ("미식파", "신메뉴 시식단·한정 디저트 선출시"),
        "여행·나들이":   ("여행파", "여행지 팝업스토어·시즌 MD"),
        "예술·문화 관여": ("컬처파", "전시·공연 제휴 쿠폰·아트 콜라보"),
        "운동·건강":     ("웰니스파", "프로틴·저당 메뉴·러닝 크루 제휴"),
        "육아·가족":     ("패밀리", "키즈존·가족 세트·주말 패밀리데이"),
        "재택·디지털":   ("카공파", "카공 패키지·콘센트 좌석·롱블랙 무한리필"),
        "패션·뷰티":     ("트렌드파", "뷰티 브랜드 콜라보 굿즈·포토존"),
    }
    tag, life_mkt = LIFE_TAG.get(life0, ("", ""))
    _suffix = f" ({tag})" if tag else ""

    if decaf > 25:
        base = "디카페인 웰니스"
        name = f"오후의 디카페인족{_suffix}"
        desc = "오후·저녁에 디카페인·논카페인을 찾는 건강 지향 고객. 카페인 부담 없이 즐기고 싶어 함"
        ua = "오후·저녁 디카페인 / 임산부·수유부 가능 / 저당·저자극 선호"
        mkt = f"디카페인 라인 전면 노출 + {life_mkt or '오후 타임딜·저녁 논카페인 추천'}"
    elif loyalty > 35:
        base = "충성 단골"
        name = f"한 브랜드 찐단골{_suffix}"
        desc = "특정 브랜드만 반복 방문하는 고충성 고객. 출근길·일상 루틴에 카페가 자리잡음"
        ua = "매일·자주 방문 / 출근길 루틴 / 멤버십·적립 적극 활용"
        mkt = f"멤버십 등급·리워드 강화 + {life_mkt or '단골 전용 시크릿 메뉴·선공개'}"
    elif nss < -10:
        base = "이탈 위험"
        name = f"떠나려는 이탈주의 고객{_suffix}"
        desc = "가격·품질·서비스 불만이 쌓여 브랜드 전환을 저울질하는 고객. 즉각 케어가 필요"
        ua = "가격·품질·서비스 불만 / 대안 적극 탐색 / 불매 동참 가능"
        mkt = "VOC 신속 대응·보상 쿠폰 + 재방문 유도 윈백(win-back) 캠페인"
    elif any(k in ["추천", "후기", "비교", "맛집", "신메뉴"] for k in kw) or life0 == "맛집·미식 탐방":
        base = "탐색·전파"
        name = f"신메뉴 얼리어답터{_suffix}"
        desc = "신메뉴·맛집을 먼저 시도하고 SNS로 후기를 퍼뜨리는 오피니언 리더"
        ua = "신메뉴 빠른 시도 / SNS 인증·후기 / 카페 투어"
        mkt = f"신메뉴 선공개·체험단·해시태그 챌린지 + {life_mkt or 'SNS 인증 리워드'}"
    elif life0:
        base = "라이프스타일"
        name = f"{life0.split('·')[0]} 취향 소비자{_suffix}"
        desc = f"{life0}에 관심이 많아, 카페를 취향·활동의 연장선으로 즐기는 소비자"
        ua = f"{life0} 연계 방문 / 취향 기반 메뉴·공간 선택"
        mkt = life_mkt or "관심사 기반 제휴·콘텐츠 마케팅"
    else:
        base = "일상 소비"
        name = "가성비 데일리족"
        desc = "특별한 선호보다 접근성·편의·가성비로 카페를 이용하는 대중 소비자"
        ua = "테이크아웃·이동 중 소비 / 가성비 중시 / 가까운 매장"
        mkt = "가성비 세트·앱 전용 쿠폰·접근성 기반 입지 마케팅"

    # [v2.4.1] Pain-point · Key Buying Factor 추론
    PAIN_BY_BASE = {
        "디카페인 웰니스": "오후 카페인 부담·당·자극 / 논카페인 선택지 부족",
        "충성 단골": "리워드 체감 부족·대기시간 / 혜택 매너리즘",
        "이탈 위험": "가격 인상·품질 편차·서비스 불만 누적",
        "탐색·전파": "신메뉴 품절·정보 과부하 / 기대 대비 실망",
        "라이프스타일": "취향·활동과 카페 경험의 연결 부족",
        "일상 소비": "가격 민감·접근성 / 차별점 부재",
    }
    KBF_BY_BASE = {
        "디카페인 웰니스": "디카페인·저당 옵션, 건강 신뢰",
        "충성 단골": "멤버십 혜택·적립, 일관된 품질·접근성",
        "이탈 위험": "가성비·문제 해결 속도, 보상",
        "탐색·전파": "신메뉴 화제성·SNS 인증성, 한정성",
        "라이프스타일": "취향 적합도·공간 경험·콜라보",
        "일상 소비": "가격·접근성·속도(편의)",
    }
    pain = PAIN_BY_BASE.get(base, "—")
    kbf = KBF_BY_BASE.get(base, "—")
    return {"name": name, "base": base, "desc": desc, "ua": ua, "marketing": mkt,
            "pain": pain, "kbf": kbf,
            "brands": brands, "nss": nss, "loyalty": loyalty, "decaf": decaf,
            "n": profile["n"], "kw": kw,
            "lifestyle": profile.get("lifestyle", []),
            "co_brands": profile.get("co_brands", []),
            "occ": profile.get("occ_top", []),
            "pos_r": profile.get("pos_r", 0), "neg_r": profile.get("neg_r", 0)}



# ══════════════════════════════════════════════════════
# Market Intelligence 모듈 [v2.3.1]
# ══════════════════════════════════════════════════════
@st.cache_data(show_spinner="시장 데이터 로딩 중...", ttl=3600)
def fetch_market_series():
    """원두 선물(KC=F)·환율 시계열. yfinance 사용, 실패 시 None 반환."""
    try:
        import yfinance as yf
        tickers = {"원두선물": "KC=F", "USD/KRW": "KRW=X", "USD/BRL": "BRL=X",
                   "EUR/KRW": "EURKRW=X", "DXY": "DX-Y.NYB"}
        out = {}
        for name, tk in tickers.items():
            h = yf.Ticker(tk).history(period="6mo")["Close"].dropna()
            if len(h) > 5:
                out[name] = h
        return out if out else None
    except Exception:
        return None


def _series_metrics(s):
    cur = float(s.iloc[-1])
    ma20 = float(s.tail(20).mean())
    ma60 = float(s.tail(60).mean())
    base = s[s.index <= (s.index[-1] - pd.Timedelta(days=30))]
    mom = (cur / float(base.iloc[-1]) - 1) * 100 if len(base) else float("nan")
    hi = float(s.tail(63).max())
    lo = float(s.tail(63).min())
    return {"cur": cur, "ma20": ma20, "ma60": ma60, "mom": mom, "hi": hi, "lo": lo}


# ── [v2.3.3] News Monitoring 자동 연동 ─────────────────
NEWS_FILE_GLOB = "news_master_*.csv"   # 크롤러(naver_news_crawler.py) 산출물 패턴

def _find_latest_news_file():
    """대시보드 폴더 / ./news / 현재 경로에서 최신 news_master_*.csv 탐색."""
    try:
        base = os.path.dirname(os.path.abspath(__file__))
    except Exception:
        base = "."
    cands = []
    for d in (base, os.path.join(base, "news"), os.path.join(base, ".news"), ".", "news", ".news"):
        try:
            cands += glob.glob(os.path.join(d, NEWS_FILE_GLOB))
        except Exception:
            pass
    cands = list(set(cands))
    if not cands:
        return None
    return max(cands, key=os.path.getmtime)  # 가장 최근 수정본 = 최신본


# [v2.4.1] 데이터 자동 로딩 — GitHub 리포지토리 .data 폴더 최신 CSV
DATA_FILE_GLOBS = ("*.csv", "*.csv.gz", "*.gz", "*.xlsx")

def _data_search_dirs():
    try:
        base = os.path.dirname(os.path.abspath(__file__))
    except Exception:
        base = "."
    return [os.path.join(base, "data"), os.path.join(base, ".data"),
            "data", ".data", base, "."]


def _find_latest_data_file():
    """data / .data 폴더(및 앱 폴더·리포 루트)에서 최신 데이터 파일 자동 탐색."""
    cands = []
    for d in _data_search_dirs():
        for g in DATA_FILE_GLOBS:
            try:
                cands += glob.glob(os.path.join(d, g))
            except Exception:
                pass
    cands = [c for c in set(cands)
             if os.path.isfile(c) and not os.path.basename(c).startswith("news_master")]
    if not cands:
        return None
    import re as _re
    def _key(p):
        m = _re.search(r"(20\d{6})", os.path.basename(p))
        return (int(m.group(1)) if m else 0, os.path.getmtime(p))
    return max(cands, key=_key)


@st.cache_data(ttl=120, show_spinner=False)
def load_news_csv(path, _mtime):
    """뉴스 CSV 로드(+표준 컬럼 보정). _mtime은 파일 변경 시 캐시 무효화용."""
    df = pd.read_csv(path)
    for col in ["일자", "브랜드", "유형", "제목", "언론사", "링크", "요약"]:
        if col not in df.columns:
            df[col] = ""
    df["일자"] = df["일자"].astype(str)
    return df


TYPE_EMOJI = {"신제품 출시": "🆕", "혜택/프로모션/이벤트": "🎁", "굿즈/MD": "🛍️", "기타 이슈": "📌"}


def news_today_stats(today_df, latest_day):
    """[v2.3.6 1.1] 날짜·브랜드·유형별 기사 수를 이모지로 시인성 있게 (HTML 반환)."""
    if len(today_df) == 0:
        return "<div class='insight-box'>오늘자 뉴스가 없습니다.</div>"
    n = len(today_df)
    by_brand = today_df["브랜드"].value_counts()
    by_type = today_df["유형"].value_counts()
    brand_lines = "　".join(f"☕ {b} <b>{c}건</b>" for b, c in by_brand.items())
    type_lines = "<br>".join(f"{TYPE_EMOJI.get(t, '📰')} {t} <b>{c}건</b>" for t, c in by_type.items())
    return (
        "<div style='background:#F7F9FC;border:1px solid #E2E8F0;border-radius:8px;"
        "padding:12px 16px;font-size:0.9rem;color:#2D3748;line-height:1.7;'>"
        f"📅 <b>{latest_day}</b> · 오늘 수집 뉴스 총 <b>{n}건</b><br><br>"
        f"<b>🏷️ 브랜드별</b><br>{brand_lines}<br><br>"
        f"<b>🗂️ 유형별</b><br>{type_lines}"
        "</div>")


def make_news_briefing(today_df, latest_day, llm_only=False):
    """[v2.3.6 1.2] 요약(G칼럼) 내용을 다시 종합 → 1000자 이내 브리핑 기사.
    ANTHROPIC_API_KEY 있으면 Claude 자연어 요약, 없으면 규칙기반 합성(기본)."""
    if len(today_df) == 0:
        return "" if llm_only else "오늘자 뉴스가 없습니다. 크롤러를 실행해 최신 뉴스를 수집하세요."
    key = os.getenv("ANTHROPIC_API_KEY")
    if key:
        try:
            import anthropic
            clips = "\n".join(
                f"- [{r['브랜드']}/{r['유형']}] {r['제목']} :: {str(r.get('요약', ''))[:120]}"
                for _, r in today_df.head(60).iterrows())
            model = os.getenv("ANTHROPIC_MODEL", "claude-3-5-sonnet-latest")
            cli = anthropic.Anthropic(api_key=key)
            res = cli.messages.create(
                model=model, max_tokens=700,
                messages=[{"role": "user", "content":
                           "다음은 오늘 국내 커피 업계 뉴스의 제목과 요약이다. 이를 종합하여 한국어로 "
                           "1000자 이내의 '오늘의 업계 종합 브리핑' 기사를 자연스러운 문장으로 작성하라. "
                           "브랜드별 핵심 움직임과 시사점을 포함하라.\n\n" + clips}])
            txt = "".join(getattr(b, "text", "") for b in res.content)
            if txt.strip():
                return txt.strip()[:1000]
        except Exception:
            pass
    if llm_only:
        return ""  # AI 요약 전용 호출 시, 키 없거나 실패하면 빈 문자열
    # 규칙기반 합성(기본) — 요약(G칼럼) 기반 유형별 종합 기사
    n = len(today_df)
    type_order = today_df["유형"].value_counts().index.tolist()
    lead = f"{latest_day} 국내 커피 업계에서는 총 {n}건의 뉴스가 포착됐다. "
    sents = []
    for typ in type_order:
        grp = today_df[today_df["유형"] == typ]
        brands = ", ".join(pd.unique(grp["브랜드"])[:4])
        descs = [re.sub(r"\s+", " ", str(d)).strip()
                 for d in grp["요약"] if str(d).strip() and str(d) != "nan"]
        rep = descs[0] if descs else re.sub(r"\s+", " ", str(grp["제목"].iloc[0]))
        rep = rep[:70]
        sents.append(f"{typ} 분야에서는 {brands} 등을 중심으로 {len(grp)}건이 다뤄졌으며, "
                     f"대표적으로 '{rep}…' 등이 보도됐다.")
    return (lead + " ".join(sents))[:1000]


def _title_sim(a, b):
    """제목 토큰 자카드 유사도 — 중복 기사 제거용."""
    ta = set(re.findall(r"[가-힣A-Za-z]{2,}", str(a)))
    tb = set(re.findall(r"[가-힣A-Za-z]{2,}", str(b)))
    if not ta or not tb:
        return 0.0
    return len(ta & tb) / len(ta | tb)


def news_reps_by_type(today_df, k=3):
    """[v2.4.1] 유형별로 유사·반복되는 대표기사 2~3건 선별 (중복 제거).
    반환: {유형: [(제목, 요약, 링크, 브랜드), ...]}"""
    res = {}
    if len(today_df) == 0:
        return res
    for typ, grp in today_df.groupby("유형"):
        all_tok = []
        for t in grp["제목"].astype(str):
            all_tok += re.findall(r"[가-힣A-Za-z]{2,}", t)
        freq = Counter(all_tok)
        scored, seen = [], []
        for _, r in grp.iterrows():
            title = str(r.get("제목", "")).strip()
            if not title or any(_title_sim(title, s) >= 0.6 for s in seen):
                continue
            seen.append(title)
            toks = re.findall(r"[가-힣A-Za-z]{2,}", title)
            score = sum(freq[w] for w in toks)  # 반복 테마일수록 高
            summary = re.sub(r"\s+", " ", str(r.get("요약", "") or "")).strip()
            scored.append((score, title, summary, str(r.get("링크", "") or ""),
                           str(r.get("브랜드", "") or "")))
        scored.sort(key=lambda x: -x[0])
        res[typ] = [(t, s, u, b) for _, t, s, u, b in scored[:k]]
    return res


def render_market_intelligence(view, v, brand_df):
    if view == "mi_raw":
        st.markdown("<div class='section-title'>🌐 원두 선물 · 국제 환율</div>",
                    unsafe_allow_html=True)
        series = fetch_market_series()
        if series is None or "원두선물" not in series:
            st.warning("⚠️ 라이브 시장 데이터(yfinance) 연결에 실패했습니다. "
                       "`pip install yfinance` 및 네트워크(야후 파이낸스) 접근을 확인하세요.")
            c1, c2 = st.columns(2)
            c1.link_button("ICE Coffee C Futures", "https://www.ice.com/products/15/coffee-c-futures/data?marketId=7455479")
            c2.link_button("investing.com 아라비카 커피", "https://kr.investing.com/commodities/arabica-coffee-4-5")
            st.caption(
                "ⓘ **ICE Coffee C Futures** — 미국 ICE(인터컨티넨탈 거래소)의 아라비카 원두 국제 선물 가격. "
                "전 세계 원두 거래의 기준 시세로, 원가 동향 파악의 1차 지표입니다.  　"
                "ⓘ **investing.com 아라비카 커피** — 같은 아라비카 커피 선물의 실시간 시세·차트를 무료로 제공하는 "
                "금융정보 사이트로, ICE 시세를 보조 확인하는 용도입니다.")
            return
        # [상단] 핵심 지표 5종
        st.markdown("**핵심 지표**")
        order = ["원두선물", "USD/KRW", "USD/BRL", "EUR/KRW", "DXY"]
        cols = st.columns(5)
        metr = {}
        for i, name in enumerate(order):
            if name in series:
                m = _series_metrics(series[name])
                metr[name] = m
                cols[i].metric(name, f"{m['cur']:,.2f}",
                               f"{m['mom']:+.1f}% (전월대비)" if m['mom'] == m['mom'] else "—")
        # [v2.4.1] 핵심 지표 5종 풋노트
        st.caption(
            "ⓘ **원두선물(KC=F)** 미국 ICE 아라비카 커피 국제 선물(¢/lb) — 원가의 기준 시세　|　"
            "**USD/KRW** 원/달러 환율 — 수입 원가의 원화 부담　|　"
            "**USD/BRL** 달러/헤알 — 최대 생산국 브라질 통화로 산지 공급가에 영향　|　"
            "**EUR/KRW** 유로/원 — 유럽산 장비·원부자재 수입가 참고　|　"
            "**DXY** 달러인덱스 — 달러 강세 정도로 원자재 가격 전반의 방향성 가늠")
        # [v2.6.1] 좌: 원두선물(20/60 MA), 우: EUR/KRW
        cL, cR = st.columns(2)
        with cL:
            s = series["원두선물"]
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=s.index, y=s.values, name="원두선물(KC=F)", line=dict(color="#6F4E37", width=2)))
            fig.add_trace(go.Scatter(x=s.index, y=s.rolling(20).mean(), name="20일 MA", line=dict(color="#E67E22", width=1, dash="dot")))
            fig.add_trace(go.Scatter(x=s.index, y=s.rolling(60).mean(), name="60일 MA", line=dict(color="#2980B9", width=1, dash="dash")))
            fig.update_layout(title="원두 선물 (Coffee C, ¢/lb)")
            chart_style(fig, height=340)
            st.plotly_chart(fig, use_container_width=True)
            m = metr.get("원두선물", {})
            if m:
                st.caption(f"최근 3개월 고점 {m['hi']:,.1f} / 저점 {m['lo']:,.1f}")
        with cR:
            if "EUR/KRW" in series:
                s2 = series["EUR/KRW"]
                fig2 = go.Figure()
                fig2.add_trace(go.Scatter(x=s2.index, y=s2.values, name="EUR/KRW", line=dict(color="#1B2A4A", width=2)))
                fig2.add_trace(go.Scatter(x=s2.index, y=s2.rolling(20).mean(), name="20일 MA", line=dict(color="#E67E22", width=1, dash="dot")))
                fig2.update_layout(title="원/유로 환율 (EUR/KRW)")
                chart_style(fig2, height=340)
                st.plotly_chart(fig2, use_container_width=True)
                m2 = metr.get("EUR/KRW", {})
                if m2:
                    st.caption(f"최근 3개월 고점 {m2['hi']:,.1f} / 저점 {m2['lo']:,.1f}")
        # [v2.3.5] 그래프 하단 외부 시세 사이트 접근 버튼 (복원)
        st.markdown("---")
        bcol1, bcol2 = st.columns(2)
        bcol1.link_button("ICE Coffee C Futures",
                          "https://www.ice.com/products/15/coffee-c-futures/data?marketId=7455479",
                          use_container_width=True)
        bcol2.link_button("investing.com 아라비카 커피",
                          "https://kr.investing.com/commodities/arabica-coffee-4-5",
                          use_container_width=True)
        st.caption("ⓘ **ICE Coffee C Futures** — 미국 ICE(인터컨티넨탈 거래소)의 아라비카 원두 국제 선물 가격(전 세계 기준 시세).　"
                   "ⓘ **investing.com 아라비카 커피** — 동일 선물의 실시간 시세·차트를 무료 제공(보조 확인용).")

    elif view == "mi_news":
        st.markdown("<div class='section-title'>📰 PAS 뉴스룸 (국내 언론)</div>",
                    unsafe_allow_html=True)
        npath = _find_latest_news_file()
        if npath is None:
            st.warning("최신 뉴스 CSV를 찾지 못했습니다.")
            st.caption(f"크롤러(naver_news_crawler.py)가 생성한 `{NEWS_FILE_GLOB}` 파일을 "
                       "대시보드와 같은 폴더(또는 ./news 폴더)에 두면 업로드 없이 자동으로 최신본을 읽어옵니다.")
            st.link_button("📧 뉴스 사내 이메일로 공유",
                           "https://hubmail.spc.co.kr/Mail/Message/CubeMessageMaster.aspx")
            return
        import datetime as _dt
        mtime = os.path.getmtime(npath)
        ndf = load_news_csv(npath, mtime)
        updated = _dt.datetime.fromtimestamp(mtime).strftime("%Y-%m-%d %H:%M")
        cap_l, cap_r = st.columns([4, 1])
        cap_l.caption(f"🔗 자동 연동: {os.path.basename(npath)} · 총 {len(ndf):,}건 · 갱신 {updated} "
                      "(크롤러 실행 후 새로고침하면 최신본 자동 반영)")
        if cap_r.button("🔄 새로고침", key="news_reload", use_container_width=True):
            st.cache_data.clear(); st.rerun()

        # 오늘의 업계 뉴스 Keywords — 최신 일자 제목 빈도 상위
        latest_day = ndf["일자"].max()
        day_titles = " ".join(ndf[ndf["일자"] == latest_day]["제목"].astype(str))
        toks = [w for w in re.findall(r"[가-힣A-Za-z]{2,}", day_titles)
                if w not in STOPWORDS and w not in TARGET_BRANDS]
        top_kw = [w for w, _ in Counter(toks).most_common(6)]
        st.markdown(f"**오늘의 업계 뉴스 Keywords**  ({latest_day})")
        if "news_kw" not in st.session_state:
            st.session_state["news_kw"] = None
        kwcols = st.columns(max(len(top_kw), 1) + 1)
        for i, kw in enumerate(top_kw):
            if kwcols[i].button(f"#{kw}", key=f"newskw_{kw}", use_container_width=True):
                st.session_state["news_kw"] = kw
        if kwcols[-1].button("전체", key="newskw_all", use_container_width=True):
            st.session_state["news_kw"] = None

        # 필터 (데이터에 실제 존재하는 값으로)
        brands_in = sorted([b for b in ndf["브랜드"].dropna().unique() if str(b).strip()])
        types_in = sorted([t for t in ndf["유형"].dropna().unique() if str(t).strip()])
        c1, c2 = st.columns(2)
        with c1:
            sel_b = st.multiselect("브랜드 필터", brands_in, default=brands_in, key="news_brand_sel")
        with c2:
            sel_t = st.multiselect("뉴스 유형", types_in, default=types_in, key="news_type_sel")

        fdf = ndf[ndf["브랜드"].isin(sel_b) & ndf["유형"].isin(sel_t)].copy()
        active_kw = st.session_state.get("news_kw")
        if active_kw:
            fdf = fdf[fdf["제목"].astype(str).str.contains(active_kw, na=False)]
            st.caption(f"🔎 키워드 필터 적용: #{active_kw}  ('전체' 버튼으로 해제)")
        fdf = fdf.sort_values("일자", ascending=False)

        # [v2.4.1] 오늘의 뉴스 브리핑 = (1.1) 이모지 통계 + (선택)AI요약 + (1.2) 유형별 대표기사
        st.markdown("<div class='section-title'>📋 오늘의 뉴스 브리핑</div>", unsafe_allow_html=True)
        _today = ndf[ndf["일자"] == latest_day]
        st.markdown(news_today_stats(_today, latest_day), unsafe_allow_html=True)
        st.markdown("")  # 줄바꿈(가독성)
        _ai = make_news_briefing(_today, latest_day, llm_only=True)
        if _ai:
            st.markdown(f"<div class='insight-box'><b>🤖 AI 종합 브리핑</b><br>{_ai}</div>",
                        unsafe_allow_html=True)
        st.markdown("<div class='section-title'>📰 오늘의 핵심 뉴스 — 유형별 대표기사</div>",
                    unsafe_allow_html=True)
        st.caption("유형별로 유사·반복되는 대표기사 2~3건을 추려 요약과 함께 정리했습니다. (제목 옆 '원문 보기'로 기사 확인)")
        _reps = news_reps_by_type(_today)
        if not _reps:
            st.caption("오늘자 대표기사가 없습니다.")
        for _typ, _items in _reps.items():
            with st.expander(f"{TYPE_EMOJI.get(_typ, '📰')} {_typ} · 대표기사 {len(_items)}건", expanded=True):
                for _title, _summary, _url, _brand in _items:
                    _link = f"  [원문 보기]({_url})" if str(_url).startswith("http") else ""
                    st.markdown(f"**[{_brand}] {_title}**{_link}")
                    if _summary and _summary != "nan":
                        st.caption(_summary)

        st.markdown(f"**뉴스 목록** — {len(fdf):,}건")
        show = fdf[["일자", "브랜드", "유형", "제목", "언론사", "링크"]]
        st.dataframe(
            show, use_container_width=True, hide_index=True, height=460,
            column_config={
                "링크": st.column_config.LinkColumn("링크", display_text="기사보기"),
                "제목": st.column_config.TextColumn("제목", width="large"),
            })
        st.link_button("📧 뉴스 사내 이메일로 공유",
                       "https://hubmail.spc.co.kr/Mail/Message/CubeMessageMaster.aspx")

    elif view == "mi_voc":
        st.markdown("<div class='section-title'>🗣️ 소비자 VoC 모니터링 시스템</div>",
                    unsafe_allow_html=True)
        st.warning("🚧 **To Be Developed** — 소비자 불만 데이터를 VoC 서버 API와 연동하여 "
                   "불만 유형/지역별/점포별로 자동 분석·시각화하는 기능입니다. "
                   "현재 API 연동 이슈로 개발 대기 중입니다.")
        c1, c2, c3 = st.columns(3)
        c1.markdown("**불만 유형 분석**\n\n(TBD)")
        c2.markdown("**지역별 분포**\n\n(TBD)")
        c3.markdown("**점포별 분포**\n\n(TBD)")


# ══════════════════════════════════════════════════════
# 사이드바
# ══════════════════════════════════════════════════════
with st.sidebar:
    # [디자인 v2.3.1] 타이틀 +15% 확대, 부제 삭제
    st.markdown(
        "<div style='font-size:1.5rem;font-weight:800;color:#1B2A4A;line-height:1.22;'>"
        "☕ PASCUCCI<br><span style='font-size:1.17rem;font-weight:700;'>"
        "Consumer Intelligence Dashboard</span></div>",
        unsafe_allow_html=True)
    # [디자인 1] 버전 표시
    st.markdown(
        f"<div style='color:#1B2A4A;font-weight:700;font-size:0.82rem;margin-top:2px;'>"
        f"{APP_VERSION} <span style='color:#718096;font-weight:400;'>({APP_DATE})</span></div>",
        unsafe_allow_html=True)
    # [디자인 2] Copyright 작은 글씨
    st.markdown(
        f"<div style='color:#A0AEC0;font-size:0.62rem;line-height:1.3;margin-top:3px;'>"
        f"{COPYRIGHT}</div>",
        unsafe_allow_html=True)
    st.markdown("---")
    # [v2.3.4] Consumer Intelligence Dashboard — MI에서 분리, 데이터 업로드 위로 이동
    if "view" not in st.session_state:
        st.session_state["view"] = "dashboard"
    if st.button("📊 컨슈머 인텔리전스 대시보드", use_container_width=True, key="nav_dash"):
        st.session_state["view"] = "dashboard"
    # [v2.5.3] 데이터 파일 업로드 위젯 제거 — 배포 버전은 데이터를 자동 인식하므로 불필요
    # (로컬 업로드가 필요한 버전에서만 재도입; 다운스트림 호환 위해 uploaded=None 유지)
    uploaded = None
    st.markdown("---")
    # [v2.3.4] Market Intelligence 네비게이션 (CID 분리)
    st.markdown("**🧭 Market Intelligence**")
    if st.button("🌐 원두선물·환율 모니터링", use_container_width=True, key="nav_raw"):
        st.session_state["view"] = "mi_raw"
    if st.button("📰 PAS 뉴스룸", use_container_width=True, key="nav_news"):
        st.session_state["view"] = "mi_news"
    if st.button("🗣️ 소비자 VoC (개발 중)", use_container_width=True, key="nav_voc"):
        st.session_state["view"] = "mi_voc"
    st.markdown("---")
    # [v2.3.1→2] 분석 리포트 다운로드 — 네비 아래 배치(데이터 로드 후 채움)
    report_dl_slot = st.empty()
    st.markdown("---")

# [v2.4.1] 업로드가 없으면 GitHub .data 폴더의 최신 파일을 자동 사용
auto_data_path = None
if uploaded is None:
    auto_data_path = _find_latest_data_file()

if uploaded is None and auto_data_path is None:
    st.markdown("""
    <div style='text-align:center; padding:60px 20px;'>
        <div style='font-size:4rem'>☕</div>
        <h1 style='color:#1B2A4A; margin:20px 0 10px;'>PASCUCCI Consumer Intelligence Dashboard</h1>
        <h3 style='color:#2D6BC4; font-weight:400;'>Consumer Data-driven Marketing for PASCUCCI</h3>
    </div>""", unsafe_allow_html=True)
    cols = st.columns(4)
    features = [
        ("📊", "Overview",       "KPI · 감성분포 · 언급량"),
        ("📈", "NSS 평판",       "Net Sentiment 월별 추이"),
        ("📣", "SOV 분석",       "담론 점유율 매트릭스"),
        ("🔍", "드라이버 분석",  "ABSA 속성별 긍/부정"),
        ("🗺️", "브랜드 이미지맵",  "브랜드 이미지 맵 - 대응일치분석(CA)"),
        ("🎯", "LDA 토픽",      "7-토픽 모델링 결과"),
        ("⚡", "Burst 탐지",    "급증 키워드 Burst Detection"),
        ("🔗", "PMI 분석",      "브랜드 공출현 강도"),
        ("📍", "포지셔닝맵",    "버블차트 · 레이더"),
        ("☁️", "키워드 버블",   "브랜드별 연관어"),
        ("🎯", "소비 맥락",     "Occasion × 브랜드"),
        ("🚨", "리스크 탐지",   "부정 드라이버 맵"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        with cols[i % 4]:
            st.markdown(f"""
            <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-top:3px solid #2D6BC4;
                        border-radius:10px;padding:14px;text-align:center;margin-bottom:10px;
                        box-shadow:0 2px 6px rgba(0,0,0,0.05)'>
                <div style='font-size:1.6rem'>{icon}</div>
                <div style='color:#1B2A4A;font-weight:700;margin:5px 0;font-size:0.9rem'>{title}</div>
                <div style='color:#718096;font-size:0.75rem'>{desc}</div>
            </div>""", unsafe_allow_html=True)
    st.info("ℹ️ 데이터를 자동으로 찾지 못했습니다. 사이드바에서 직접 업로드하거나, "
            "GitHub 리포지토리에 **data/ 폴더(점 없이)**를 만들고 그 안에 크롤링 CSV를 커밋하세요. "
            "(`.gitignore`에 `*.csv` 또는 `data/`가 포함되면 깃에 올라가지 않으니 확인 필요)")
    with st.expander("자동 탐색 경로 확인"):
        for _d in _data_search_dirs():
            st.caption(f"· {_d}")
    st.stop()

# ══════════════════════════════════════════════════════
# 데이터 로딩 — 업로드 우선, 없으면 data/.data 폴더 최신 파일 [v2.4.2]
# ══════════════════════════════════════════════════════
@st.cache_data(show_spinner=False)
def _tag_source_quality(v):
    """소스(블로그/커뮤니티) + 콘텐츠 유형(리뷰/여론/혜택·홍보성/보도·뉴스) 태깅. [v2.6.4]
    현재 본문이 짧아 휴리스틱 기반(풀텍스트 재수집 후 정밀화 예정)."""
    import re as _re
    v = v.copy()
    _src = v["source_type"].astype(str) if "source_type" in v.columns \
        else __import__("pandas").Series(["community"] * len(v), index=v.index)
    v["_src"] = _src.map({"blog": "블로그", "community": "커뮤니티"}).fillna("기타")
    if "full_text" in v.columns:
        _txt = v["full_text"].astype(str)
    else:
        _txt = (v.get("title", "").fillna("").astype(str) + " " + v.get("body", "").fillna("").astype(str))
    _PROMO = r"배민|배달의민족|배달의 민족|쿠팡이츠|요기요|쿠폰|할인 정보|할인정보|할인코드|이벤트 정리|픽업 할인|적립 혜택|혜택 정리|혜택정리|프로모션 정리|페스타|천원의|기프티콘 이벤트"
    _PR = r"보도자료|\[알림\]|연합뉴스|뉴시스|뉴스1|조선비즈|머니투데이|이데일리|기자\]|배포일|출처\s*=|제공\s*=|보도 자료"
    _promo = _txt.str.contains(_PROMO, regex=True, na=False)
    _pr = _txt.str.contains(_PR, regex=True, na=False)
    import pandas as _pd
    _ct = _pd.Series("블로그 리뷰", index=v.index)
    _ct[v["_src"] == "커뮤니티"] = "커뮤니티 여론"
    _ct[(v["_src"] == "블로그") & _promo] = "혜택·홍보성"
    _ct[(v["_src"] == "블로그") & _pr] = "보도·뉴스"
    v["_ctype"] = _ct
    return v


_data_source = uploaded if uploaded is not None else auto_data_path
try:
    v = load_and_preprocess(_data_source)
except Exception as _e:
    import traceback as _tb
    st.error(f"데이터 로딩 중 오류: {type(_e).__name__}: {_e}")
    st.code(_tb.format_exc())
    st.info("업로드한 CSV가 82Cook 크롤링 표준 스키마인지, 파일이 손상되지 않았는지 확인하세요.")
    st.stop()
if uploaded is None and auto_data_path is not None:
    st.sidebar.caption(f"🔗 자동 로딩: {os.path.basename(auto_data_path)}")

v = _tag_source_quality(v)

with st.sidebar:
    st.markdown("**🔧 필터**")
    brand_filter = st.multiselect(
        "브랜드",
        sorted(TARGET_BRANDS),
        default=sorted(TARGET_BRANDS)
    )
    sent_filter = st.multiselect(
        "감성",
        ["긍정", "중립", "부정"],
        default=["긍정", "중립", "부정"]
    )
    st.markdown("---")
    st.markdown("**🧭 데이터 소스 (1단계)**")
    _src_all = [x for x in ["블로그", "커뮤니티"] if x in set(v["_src"])]
    src_sel = st.multiselect("소스 선택", _src_all, default=_src_all,
                             help="블로그=노출·후기 / 커뮤니티=여론·리스크. 두 소스는 성격이 달라 분리해 봐야 합니다.")
    st.markdown("**🎯 콘텐츠 품질 (2단계)**")
    qual_sel = st.radio("필터", ["전체", "VoC만 (리뷰·여론)", "홍보·보도성만"], index=0,
                        help="VoC만=혜택·홍보·보도성 글 제외(리얼 여론). 홍보·보도성만=역으로 점검용.")
    with st.expander("📊 소스별 특성 비교", expanded=False):
        _rows = []
        for _sname, _g in v.groupby("_src"):
            _pos = (_g["sentiment"] == "긍정").mean() * 100
            _neg = (_g["sentiment"] == "부정").mean() * 100
            _rows.append({"소스": _sname, "건수": len(_g), "긍정%": round(_pos, 1),
                          "부정%": round(_neg, 1), "NSS": round(_pos - _neg, 1)})
        if _rows:
            st.dataframe(pd.DataFrame(_rows).set_index("소스"), use_container_width=True)
            st.caption("같은 시장이라도 소스에 따라 감성·구성이 크게 달라집니다.")
    st.markdown("---")
    st.markdown(f"<small style='color:#99AACC'>총 {len(v):,}건 로드됨</small>",
                unsafe_allow_html=True)

# ── 소스/품질 필터 적용: 선택한 렌즈로 분석 범위 재설정(모든 분석에 반영) ──
_mask = v["_src"].isin(src_sel)
if qual_sel == "VoC만 (리뷰·여론)":
    _mask &= v["_ctype"].isin(["커뮤니티 여론", "블로그 리뷰"])
elif qual_sel == "홍보·보도성만":
    _mask &= v["_ctype"].isin(["혜택·홍보성", "보도·뉴스"])
v = v[_mask].copy()
if len(v) == 0:
    st.warning("선택한 소스/품질 조건에 맞는 데이터가 없습니다. 필터를 완화하세요.")
    st.stop()
st.sidebar.caption(f"🔍 현재 분석 범위: {' + '.join(src_sel) if src_sel else '없음'} · {qual_sel} · {len(v):,}건")

fv = v[v["sentiment"].isin(sent_filter)].copy()
if len(fv) == 0:
    st.warning("필터 조건에 맞는 데이터가 없습니다.")
    st.stop()

# 분석 실행
brand_df = compute_brand_stats(v)
monthly_nss_df, all_months = compute_monthly_nss(v)
absa_df   = compute_absa(v)
lda_model, tfidf_vec, vocab, doc_topics, topic_dist = compute_lda(v)
ca_result = compute_ca(v)
burst_df  = compute_burst(v)
pmi_df    = compute_pmi(v)
loyalty_df = compute_loyalty(v)
occ_detail_df = compute_occasion_detailed(v)
occ_brand_mat = compute_occasion_brand_matrix(v)
topic_cluster = compute_topic_clustering(v)
trend_ew  = compute_trend_earlywarning(v)
v_adj = build_adjacent_food_voc(_data_source)
fnb_trends = compute_fnb_trends(v_adj)
promo_brand_mat, promo_total_df = compute_promotion(v)
promo_react_df, promo_type_react_df, promo_doc_count = compute_promo_reaction(v)
promo_recent_df, brand_promo_df = compute_promo_recent(v)
promo_surge_df, promo_cross_df = compute_promo_surge(v)
customer_cluster = compute_customer_clustering(v)

# ══════════════════════════════════════════════════════
# [v2.5.5] 외부 근거 — 2025 서울시민 먹거리조사 참조 연동
# ══════════════════════════════════════════════════════
@st.cache_data(show_spinner=False)
def load_survey_ref():
    """reference/seoul_food_survey_ref.json 자동 탐색·로드 → dict (없으면 {})."""
    cands = []
    for d in _data_search_dirs():
        cands.append(os.path.join(d, "reference", "seoul_food_survey_ref.json"))
        cands.append(os.path.join(d, "seoul_food_survey_ref.json"))
    for p in cands:
        try:
            if os.path.exists(p):
                with open(p, encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            continue
    return {}


SURVEY_REF = load_survey_ref()

# 리포트 섹션 → 조사 토픽 매핑
REPORT_REF_MAP = {
    "intro":   ["price_perception", "food_cost_burden", "diet_lifestyle"],
    "nss":     ["diet_lifestyle", "nutrition_label"],
    "loyalty": ["info_channels"],
    "trend":   ["coffee", "zero_drinks", "zero_foods", "sweet_drinks"],
    "promo":   ["food_cost_burden", "price_perception"],
    "persona": ["dining_out", "diet_lifestyle"],
    "risk":    ["diet_lifestyle", "wellbeing"],
}


def survey_meta_line():
    m = (SURVEY_REF or {}).get("meta", {})
    if not m:
        return "2025 서울시민 먹거리조사"
    return f"{m.get('title', '2025 서울시민 먹거리조사')} ({m.get('sample', '')})"


def survey_comment(section_key, max_topics=2, max_findings=2):
    """섹션 키 → 해당 조사 토픽의 근거 코멘트 리스트(없으면 [])."""
    topics = (SURVEY_REF or {}).get("topics", {})
    if not topics:
        return []
    out = []
    for tk in REPORT_REF_MAP.get(section_key, [])[:max_topics]:
        t = topics.get(tk)
        if not t:
            continue
        fnd = " / ".join(t.get("findings", [])[:max_findings])
        use = t.get("pcid_use", "")
        line = f"[{t.get('label', tk)}] {fnd}"
        if use:
            line += f" → {use}"
        out.append(line)
    return out


# ══════════════════════════════════════════════════════
# [기능 3] 분석 결과 요약·인사이트 docx 생성
# ══════════════════════════════════════════════════════
def build_summary_docx():
    """[v2.6.1] 사전저작 코멘트(report_commentary.json) 기반 워드 리포트 → bytes.
    PPTX와 동일 스킴: 표제/Executive Summary/5개 챕터(결과·의미·함의·시나리오·시장근거)/마무리.
    차트는 matplotlib 이미지(가능 시)·표(폴백), 폰트 맑은 고딕, 데이터 지문 가드 포함."""
    import io, tempfile
    from docx import Document
    from docx.shared import Pt as DPt, Inches as DInches, RGBColor as DRGB
    from docx.oxml.ns import qn as _dqn
    from docx.enum.text import WD_ALIGN_PARAGRAPH as WDA

    NAVY = DRGB(0x1B, 0x2A, 0x4A); RED = DRGB(0xC0, 0x39, 0x2B)
    GRAY = DRGB(0x71, 0x80, 0x96); GOLD = DRGB(0xB4, 0x8A, 0x3C); DARK = DRGB(0x2D, 0x37, 0x48)

    com = load_report_commentary()
    REG = {}
    for nm in ("brand_df", "absa_df", "loyalty_df", "occ_detail_df", "promo_type_react_df"):
        if nm in globals():
            REG[nm] = globals()[nm]

    PLT = None; KFP = None
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as _plt
        from matplotlib import font_manager as _fm
        fpath = _kfont_path()
        if fpath:
            _fm.fontManager.addfont(fpath); KFP = _fm.FontProperties(fname=fpath)
        _plt.rcParams["axes.unicode_minus"] = False
        PLT = _plt
    except Exception:
        PLT = None

    doc = Document()
    # 기본 폰트 맑은 고딕
    try:
        st_normal = doc.styles["Normal"]
        st_normal.font.name = "맑은 고딕"; st_normal.font.size = DPt(11)
        st_normal.element.rPr.rFonts.set(_dqn('w:eastAsia'), "맑은 고딕")
    except Exception:
        pass

    def kfont(run, size=None, bold=None, color=None):
        run.font.name = "맑은 고딕"
        rPr = run._element.get_or_add_rPr()
        rF = rPr.find(_dqn('w:rFonts'))
        if rF is None:
            rF = rPr.makeelement(_dqn('w:rFonts'), {}); rPr.append(rF)
        rF.set(_dqn('w:eastAsia'), "맑은 고딕")
        if size is not None: run.font.size = DPt(size)
        if bold is not None: run.font.bold = bold
        if color is not None: run.font.color.rgb = color

    def para(text="", size=11, bold=False, color=None, align=None, after=4, bold_token="파스쿠찌"):
        p = doc.add_paragraph(); p.paragraph_format.space_after = DPt(after)
        if align is not None: p.alignment = align
        import re
        for seg in (re.split(r'(' + bold_token + r')', text) if bold_token else [text]):
            if seg == "": continue
            r = p.add_run(seg); kfont(r, size, bold or seg == bold_token, color)
        return p

    def label_body(label, body, lcolor=NAVY):
        if not body: return
        p = doc.add_paragraph(); p.paragraph_format.space_after = DPt(2)
        r = p.add_run("▪ " + label + "  "); kfont(r, 11, True, lcolor)
        import re
        for seg in re.split(r'(파스쿠찌)', body):
            if seg == "": continue
            r2 = p.add_run(seg); kfont(r2, 11, seg == "파스쿠찌", DARK)

    def render_bar_png(cats, vals, title=None):
        cats = [str(c) for c in cats]; vals = [float(x) for x in vals]
        o = sorted(range(len(vals)), key=lambda i: vals[i])
        cats = [cats[i] for i in o]; vals = [vals[i] for i in o]
        fig, ax = PLT.subplots(figsize=(6.6, 3.6), dpi=180)
        CN = "#1B2A4A"; CR = "#C03929"; CD = "#2D3748"
        for i, (c, val) in enumerate(zip(cats, vals)):
            if c == "파스쿠찌":
                ax.barh(i, val, color=CR, edgecolor=CR, height=0.62, zorder=3)
            elif val < 0:
                ax.barh(i, val, facecolor="white", edgecolor=CN, hatch="////", linewidth=0.6, height=0.62, zorder=3)
            else:
                ax.barh(i, val, color=CN, edgecolor=CN, height=0.62, zorder=3)
        ax.set_yticks(range(len(vals))); ax.set_yticklabels(cats, fontproperties=KFP, fontsize=9)
        for t, c in zip(ax.get_yticklabels(), cats):
            if c == "파스쿠찌": t.set_fontweight("bold"); t.set_color(CR)
        for lab in ax.get_xticklabels():
            if KFP: lab.set_fontproperties(KFP)
            lab.set_fontsize(8)
        span = (max(vals) - min(vals)) or 1
        for i, val in enumerate(vals):
            off = span * 0.013; c = cats[i]
            ax.text(val + (off if val >= 0 else -off), i, f"{val:.1f}", va="center",
                    ha=("left" if val >= 0 else "right"), fontsize=9, fontproperties=KFP,
                    color=(CR if c == "파스쿠찌" else CD), fontweight=("bold" if c == "파스쿠찌" else "normal"), zorder=4)
        ax.axvline(0, color="#888888", linewidth=0.6, zorder=2)
        ax.grid(axis="x", color="#CCCCCC", linewidth=0.4, zorder=0); ax.set_axisbelow(True)
        for sp in ("top", "right", "left"): ax.spines[sp].set_visible(False)
        ax.spines["bottom"].set_color("#BBBBBB"); ax.spines["bottom"].set_linewidth(0.5)
        ax.set_xlim(min(0, min(vals)) * 1.15 - 2, max(0, max(vals)) * 1.18 + 4); ax.margins(y=0.02)
        if title: ax.set_title(title, fontproperties=KFP, fontsize=11, fontweight="bold", color=CD)
        fig.tight_layout(pad=0.4)
        fn = tempfile.NamedTemporaryFile(suffix=".png", delete=False).name
        fig.savefig(fn, dpi=180, transparent=False, bbox_inches="tight"); PLT.close(fig)
        return fn

    def add_chart(sec):
        ch = sec.get("chart", {}) or {}
        if ch.get("type") == "barh":
            src = REG.get(ch.get("source")); cat = ch.get("cat"); val = ch.get("val")
            if src is not None and cat in src.columns and val in src.columns:
                d = src[[cat, val]].dropna()
                if PLT is not None:
                    try:
                        fn = render_bar_png(d[cat].tolist(), d[val].tolist(), ch.get("title"))
                        doc.add_picture(fn, width=DInches(6.3)); return
                    except Exception:
                        pass
                # 표 폴백
                try:
                    dd = d.sort_values(val, ascending=False)
                    tb = doc.add_table(rows=1, cols=2); tb.style = "Light Grid Accent 1"
                    h = tb.rows[0].cells; h[0].text = str(cat); h[1].text = str(val)
                    for _, rrow in dd.iterrows():
                        c = tb.add_row().cells
                        c[0].text = str(rrow[cat]); c[1].text = f"{float(rrow[val]):.1f}"
                    return
                except Exception:
                    pass
        # text형: 데이터 요약 bullet
        for b in ch.get("bullets", []):
            p = doc.add_paragraph(style="List Bullet"); r = p.add_run(b); kfont(r, 10.5, False, DARK)

    # ===== 표제 =====
    es = com.get("executive_summary", {}); meta = com.get("meta", {}); fp = meta.get("fingerprint", {}) or {}
    live_docs = len(v) if "v" in globals() else None
    stale = (fp.get("docs") is not None and live_docs is not None and fp.get("docs") != live_docs)
    para("PASCUCCI 브랜드 컨슈머 인텔리전스 리포트", 20, True, NAVY, WDA.CENTER, after=2)
    _period = APP_DATE[:7].replace("-", ".") if isinstance(APP_DATE, str) and len(APP_DATE) >= 7 else "2026.06"
    para(f"{_period}  ·  {APP_VERSION}  ·  분석 문서 {live_docs:,}건", 10, False, GRAY, WDA.CENTER, after=10)
    if stale:
        para(f"⚠ 데이터가 갱신됨(현재 {live_docs}건 ≠ 코멘트 {fp.get('docs')}건) — 코멘트 재생성 필요", 10, True, RED, after=8)

    # ===== Executive Summary =====
    if es:
        para("Executive Summary", 15, True, NAVY, after=4)
        para(es.get("headline", ""), 12, True, DARK, after=8)
        para("핵심 인사이트", 12, True, NAVY, after=2)
        for c in es.get("key_insights", []):
            p = doc.add_paragraph(style="List Bullet"); p.paragraph_format.space_after = DPt(2)
            r = p.add_run(f"[{c.get('area','')}] "); kfont(r, 11, True, GOLD)
            import re
            for seg in re.split(r'(파스쿠찌)', c.get("text", "")):
                if seg == "": continue
                rr = p.add_run(seg); kfont(rr, 11, seg == "파스쿠찌", DARK)
        para("", 4)
        para("Key Recommendations", 12, True, RED, after=2)
        for it in es.get("recommendations", []):
            p = doc.add_paragraph(style="List Bullet"); p.paragraph_format.space_after = DPt(2)
            import re
            for seg in re.split(r'(파스쿠찌)', it):
                if seg == "": continue
                rr = p.add_run(seg); kfont(rr, 11, seg == "파스쿠찌", DARK)

    # ===== 챕터별 섹션 =====
    chno = 0; cur = None
    for sec in com.get("sections", []):
        try:
            ch = sec.get("chapter", "")
            if ch != cur:
                cur = ch; chno += 1
                doc.add_page_break()
                para(f"Chapter {chno}.  {ch}", 16, True, NAVY, after=6)
            para(sec.get("tag", ch), 10, True, GOLD, after=1)
            para(sec.get("headline", ""), 13, True, DARK, after=6)
            add_chart(sec)
            para("", 2)
            label_body("의미 (So what)", sec.get("meaning", ""))
            label_body("함의 (Implication)", sec.get("implication", ""))
            scs = sec.get("scenarios", [])
            if scs:
                p = doc.add_paragraph(); p.paragraph_format.space_after = DPt(1)
                r = p.add_run("▪ 실행 시나리오 (Action)"); kfont(r, 11, True, RED)
                for s2 in scs:
                    pp = doc.add_paragraph(style="List Bullet"); pp.paragraph_format.space_after = DPt(1)
                    import re
                    for seg in re.split(r'(파스쿠찌)', s2):
                        if seg == "": continue
                        rr = pp.add_run(seg); kfont(rr, 11, seg == "파스쿠찌", DARK)
            mv = sec.get("market_view", "")
            if mv:
                p = doc.add_paragraph(); p.paragraph_format.space_after = DPt(8)
                r = p.add_run("시장 근거 · 서울 먹거리조사  "); kfont(r, 9.5, True, GOLD)
                r2 = p.add_run(mv); kfont(r2, 9.5, False, GRAY)
        except Exception:
            continue

    if not com.get("sections"):
        para("reference/report_commentary.json 을 리포에 올리면 분석 서술이 반영됩니다.", 11, False, GRAY)

    para("", 6)
    foot = para(COPYRIGHT if "COPYRIGHT" in globals() else "© PAS DIVISION", 9, False, GRAY, WDA.CENTER)

    buf = io.BytesIO(); doc.save(buf); buf.seek(0)
    return buf.getvalue()


def _report_template_path():
    """리포트 템플릿(insight_report_templ*.pptx) 자동 탐색 — 대소문자·하위폴더·철자 무시."""
    bases = []
    try:
        bases.append(os.path.dirname(os.path.abspath(__file__)))
    except Exception:
        pass
    bases += [os.getcwd(), "."]
    names = ["insight_report_templete.pptx", "insight_report_template.pptx"]
    # 1) 흔한 위치 직접 확인
    for b in bases:
        for sub in ("reference", "templates", "data", ".data", ""):
            for nm in names:
                p = os.path.join(b, sub, nm)
                if os.path.isfile(p):
                    return p
    # 2) 대소문자·철자·폴더 무시 재귀 탐색
    seen = set()
    for b in bases:
        try:
            b = os.path.abspath(b)
        except Exception:
            continue
        if b in seen:
            continue
        seen.add(b)
        for root, dirs, files in os.walk(b):
            dirs[:] = [d for d in dirs if d not in (".git", "__pycache__", ".venv", "venv", "node_modules")]
            for f in files:
                fl = f.lower()
                if fl.endswith(".pptx") and fl.startswith("insight_report_templ"):
                    return os.path.join(root, f)
    return None


def _list_repo_pptx():
    """진단용: 작업 경로 하위 모든 .pptx 상대경로 목록."""
    out = []
    try:
        base = os.getcwd()
        for root, dirs, files in os.walk(base):
            dirs[:] = [d for d in dirs if d not in (".git", "__pycache__", ".venv", "venv", "node_modules")]
            for f in files:
                if f.lower().endswith(".pptx"):
                    out.append(os.path.relpath(os.path.join(root, f), base))
    except Exception:
        pass
    return out


def _kfont_path():
    """matplotlib용 한글 폰트 경로 탐색(Noto CJK·나눔 등)."""
    import glob
    cands = ["/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
             "/usr/share/fonts/truetype/nanum/NanumGothic.ttf"]
    for pat in ("/usr/share/fonts/**/*NotoSansCJK*", "/usr/share/fonts/**/*NotoSansKR*",
                "/usr/share/fonts/**/*Nanum*"):
        cands += glob.glob(pat, recursive=True)
    for p in cands:
        if os.path.exists(p):
            return p
    return None


@st.cache_data(show_spinner=False)
def load_report_commentary():
    """reference/report_commentary.json 자동 탐색·로드 → dict(없으면 {})."""
    for d in _data_search_dirs():
        for p in (os.path.join(d, "reference", "report_commentary.json"),
                  os.path.join(d, "report_commentary.json")):
            try:
                if os.path.isfile(p):
                    with open(p, encoding="utf-8") as f:
                        return json.load(f)
            except Exception:
                continue
    return {}


def build_summary_pptx():
    """[v2.6.0] 브랜드 템플릿 + 사전저작 코멘트(report_commentary.json) 기반
    컨설팅형 분석 리포트 → bytes. 표지/Exec Summary/5개 챕터(분석별 결과·의미·함의·시나리오)
    /마지막 장. 차트는 라이브 데이터 이미지, 서술은 코멘트 파일에서 반영(지문 동기화 가드 포함).
    폰트 맑은 고딕(자간 좁게), 챕터 10·헤드라인 18·본문 12·범례·수치 9pt."""
    import io, copy, math, tempfile
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn as _qn

    MALGUN = "맑은 고딕"
    NAVY = PRGB(0x1B, 0x2A, 0x4A); RED = PRGB(0xC0, 0x39, 0x2B); GRAY = PRGB(0x71, 0x80, 0x96)
    DARK = PRGB(0x2D, 0x37, 0x48); PANEL = PRGB(0xF4, 0xF6, 0xF9); LINE = PRGB(0xE3, 0xE8, 0xEF)
    GOLD = PRGB(0xB4, 0x8A, 0x3C); BLACK = PRGB(0x11, 0x11, 0x11)
    CARD = PRGB(0xF7, 0xF8, 0xFA); CARDLINE = PRGB(0xDD, 0xE3, 0xEC)

    tpl = _report_template_path()
    if not tpl:
        prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9.0), Inches(5.5)).text_frame
        tf.word_wrap = True
        _found = _list_repo_pptx()
        for i, ln in enumerate([
            "리포트 템플릿을 찾지 못했습니다. (진단 정보)",
            f"· 작업 경로(cwd): {os.getcwd()}",
            f"· 리포 내 .pptx 파일: {', '.join(_found) if _found else '(하나도 없음)'}",
            "→ reference/insight_report_templete.pptx 가 main 브랜치에 올라가 있는지 확인하세요."]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = ln; r.font.size = PPt(13)
        buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf.getvalue()

    com = load_report_commentary()

    # ---- matplotlib 차트 렌더 준비 ----
    PLT = None; KFP = None
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as _plt
        from matplotlib import font_manager as _fm
        fpath = _kfont_path()
        if fpath:
            _fm.fontManager.addfont(fpath); KFP = _fm.FontProperties(fname=fpath)
        _plt.rcParams["axes.unicode_minus"] = False
        PLT = _plt
    except Exception:
        PLT = None

    REG = {}
    for nm in ("brand_df", "absa_df", "loyalty_df", "occ_detail_df", "promo_type_react_df"):
        if nm in globals():
            REG[nm] = globals()[nm]

    prs = Presentation(tpl)
    SW = prs.slide_width / 914400

    def _setface(rPr):
        rPr.set('spc', '-80')
        for t in ("a:latin", "a:ea", "a:cs"):
            e = rPr.find(_qn(t))
            if e is None: e = rPr.makeelement(_qn(t), {}); rPr.append(e)
            e.set("typeface", MALGUN)

    def kf(run, size, bold=False, color=None):
        run.font.size = PPt(size); run.font.bold = bold
        if color is not None: run.font.color.rgb = color
        _setface(run._r.get_or_add_rPr())

    def addrich(p, text, size, bold, color):
        import re
        for seg in re.split(r'(파스쿠찌)', text):
            if seg == "": continue
            r = p.add_run(); r.text = seg; kf(r, size, bold or seg == "파스쿠찌", color)

    def txt(sl, x, y, w, h, lines, size, bold=False, color=DARK, sp=2):
        tf = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
        first = True
        for ln in (lines if isinstance(lines, list) else [lines]):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_after = PPt(sp); addrich(p, ln, size, bold, color); first = False
        return tf

    def setpara(p, text):
        rs = p.runs
        if rs:
            rs[0].text = text
            for r in rs[1:]: r._r.getparent().remove(r._r)
            _setface(rs[0]._r.get_or_add_rPr())
        else:
            r = p.add_run(); r.text = text; _setface(r._r.get_or_add_rPr())

    def resize_logo(sl, scale=0.9):
        for sh in sl.shapes:
            if sh.shape_type == 13 and sh.width:
                w = sh.width; sh.width = int(w * scale); sh.height = int(sh.height * scale)
                sh.left = sh.left + (w - sh.width)
                return sh.left / 914400
        return SW - 1.5

    def dup(i):
        src = prs.slides[i]; dest = prs.slides.add_slide(src.slide_layout)
        for sh in list(dest.shapes): sh._element.getparent().remove(sh._element)
        for sh in src.shapes: dest.shapes._spTree.append(copy.deepcopy(sh._element))
        for _r, rel in src.part.rels.items():
            if "image" in rel.reltype: dest.part.rels.get_or_add(rel.reltype, rel._target)
        dest._logo_left = resize_logo(dest, 0.9)
        return dest

    def reorder(order):
        lst = prs.slides._sldIdLst; by = {}
        for s in list(lst):
            by[id(prs.part.related_part(s.get(_qn('r:id'))))] = s; lst.remove(s)
        for sl in order: lst.append(by[id(sl.part)])

    def topline(sl, y=0.36):
        x2 = min(getattr(sl, "_logo_left", SW - 1.5) - 0.15, SW - 0.6)
        ln = sl.shapes.add_connector(2, Inches(0.6), Inches(y), Inches(x2), Inches(y))
        ln.line.color.rgb = BLACK; ln.line.width = PPt(0.5); ln.shadow.inherit = False

    def head(sl, tag, headline):
        topline(sl)
        txt(sl, 0.6, 0.45, SW - 1.2, 0.30, tag, 10, bold=True, color=GOLD)
        txt(sl, 0.6, 0.74, SW - 1.2, 0.95, headline, 18, bold=True, color=NAVY)

    def panel(sl, x, y, w, h, fill=PANEL, line=LINE, rounded=False):
        shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        sp = sl.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.color.rgb = line
        sp.line.width = PPt(0.75); sp.shadow.inherit = False
        if rounded:
            try: sp.adjustments[0] = 0.06
            except Exception: pass
        return sp

    def rightcol(sl, sec, x=6.15, y=1.82, w=4.15, h=4.55):
        panel(sl, x, y, w, h); ix = x + 0.16; iw = w - 0.32; cy = y + 0.16
        for lbl, key, col in [("의미 (So what)", "meaning", NAVY), ("함의 (Implication)", "implication", NAVY)]:
            body = sec.get(key, "")
            if not body: continue
            txt(sl, ix, cy, iw, 0.3, "▪ " + lbl, 13, bold=True, color=col, sp=1); cy += 0.32
            txt(sl, ix + 0.02, cy, iw - 0.02, 1.0, body, 12, color=DARK, sp=1)
            cy += 0.205 * max(1, math.ceil(len(body) / 22)) + 0.10
        scs = sec.get("scenarios", [])
        if scs:
            txt(sl, ix, cy, iw, 0.3, "▪ 실행 시나리오 (Action)", 13, bold=True, color=RED, sp=1); cy += 0.32
            for sc in scs:
                txt(sl, ix + 0.02, cy, iw - 0.02, 0.6, "· " + sc, 12, color=DARK, sp=1)
                cy += 0.205 * max(1, math.ceil(len(sc) / 21)) + 0.06

    def market_band(sl, text, y=6.5):
        if not text: return
        sp = panel(sl, 0.55, y, SW - 1.1, 0.64, fill=PRGB(0xF7, 0xF2, 0xE8), line=PRGB(0xE6, 0xD8, 0xBE))
        tf = sp.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]; r = p.add_run(); r.text = "시장 근거 · 서울 먹거리조사   "; kf(r, 10, True, GOLD)
        addrich(p, text, 10, False, GRAY)

    def textbox_left(sl, title, bullets, x=0.55, y=1.82, w=5.35, h=4.55):
        panel(sl, x, y, w, h, fill=PRGB(0xFB, 0xFC, 0xFD), line=LINE)
        if title:
            txt(sl, x + 0.18, y + 0.16, w - 0.36, 0.32, title, 12, bold=True, color=NAVY)
        cy = y + 0.62
        for b in bullets:
            txt(sl, x + 0.18, cy, w - 0.36, 0.6, ("· " + b) if b[:1] not in ("·", "▪", "→", "(") else b, 12, color=DARK, sp=1)
            cy += 0.205 * max(1, math.ceil(len(b) / 30)) + 0.10

    def render_bar_png(cats, vals, title=None, w=5.5, h=4.3):
        cats = [str(c) for c in cats]; vals = [float(x) for x in vals]
        o = sorted(range(len(vals)), key=lambda i: vals[i])
        cats = [cats[i] for i in o]; vals = [vals[i] for i in o]
        fig, ax = PLT.subplots(figsize=(w, h), dpi=200)
        CN = "#1B2A4A"; CR = "#C03929"; CD = "#2D3748"
        for i, (c, val) in enumerate(zip(cats, vals)):
            if c == "파스쿠찌":
                ax.barh(i, val, color=CR, edgecolor=CR, height=0.62, zorder=3)
            elif val < 0:
                ax.barh(i, val, facecolor="white", edgecolor=CN, hatch="////", linewidth=0.6, height=0.62, zorder=3)
            else:
                ax.barh(i, val, color=CN, edgecolor=CN, height=0.62, zorder=3)
        ax.set_yticks(range(len(vals)))
        ax.set_yticklabels(cats, fontproperties=KFP, fontsize=9)
        for t, c in zip(ax.get_yticklabels(), cats):
            if c == "파스쿠찌": t.set_fontweight("bold"); t.set_color(CR)
        for lab in ax.get_xticklabels():
            if KFP: lab.set_fontproperties(KFP)
            lab.set_fontsize(8)
        span = (max(vals) - min(vals)) or 1
        for i, val in enumerate(vals):
            off = span * 0.013; c = cats[i]
            ax.text(val + (off if val >= 0 else -off), i, f"{val:.1f}",
                    va="center", ha=("left" if val >= 0 else "right"), fontsize=9,
                    fontproperties=KFP, color=(CR if c == "파스쿠찌" else CD),
                    fontweight=("bold" if c == "파스쿠찌" else "normal"), zorder=4)
        ax.axvline(0, color="#888888", linewidth=0.6, zorder=2)
        ax.grid(axis="x", color="#CCCCCC", linewidth=0.4, zorder=0); ax.set_axisbelow(True)
        for sp in ("top", "right", "left"): ax.spines[sp].set_visible(False)
        ax.spines["bottom"].set_color("#BBBBBB"); ax.spines["bottom"].set_linewidth(0.5)
        ax.set_xlim(min(0, min(vals)) * 1.15 - 2, max(0, max(vals)) * 1.18 + 4); ax.margins(y=0.02)
        if title:
            ax.set_title(title, fontproperties=KFP, fontsize=11, fontweight="bold", color=CD)
        fig.tight_layout(pad=0.4)
        fn = tempfile.NamedTemporaryFile(suffix=".png", delete=False).name
        fig.savefig(fn, dpi=200, transparent=True, bbox_inches="tight"); PLT.close(fig)
        return fn

    def _native_bar(sl, cats, vals, title):
        """matplotlib 미사용 폴백 — 네이티브 PPT 가로막대(외부 의존성 없음, PowerPoint에서 항상 렌더)."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        pairs = sorted(zip([str(c) for c in cats], [float(x) for x in vals]), key=lambda t: t[1])
        cs = [p[0] for p in pairs]; vs = [p[1] for p in pairs]
        cd = CategoryChartData(); cd.categories = cs; cd.add_series(title or "값", vs)
        gf = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.85), Inches(5.6), Inches(4.3), cd)
        chart = gf.chart; chart.has_legend = False
        try:
            chart.font.size = PPt(9); chart.font.name = MALGUN
            txPr = chart._chartSpace.find(_qn('c:txPr'))
            if txPr is not None:
                for d in txPr.iter(_qn('a:defRPr')):
                    for t in ('a:ea', 'a:cs'):
                        e = d.find(_qn(t))
                        if e is None: e = d.makeelement(_qn(t), {}); d.append(e)
                        e.set('typeface', MALGUN)
        except Exception: pass
        try:
            ser = chart.series[0]
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = NAVY
            if "파스쿠찌" in cs:
                pt = ser.points[cs.index("파스쿠찌")]
                pt.format.fill.solid(); pt.format.fill.fore_color.rgb = RED
        except Exception: pass
        try:
            plot = chart.plots[0]; plot.has_data_labels = True
            plot.data_labels.font.size = PPt(9); plot.data_labels.font.name = MALGUN
            plot.data_labels.font.color.rgb = DARK
            plot.data_labels.number_format = '0.0'; plot.data_labels.number_format_is_linked = False
        except Exception: pass
        try:
            va = chart.value_axis; va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = PRGB(0xCC, 0xCC, 0xCC)
            va.major_gridlines.format.line.width = PPt(0.4)
            chart.category_axis.has_major_gridlines = False
        except Exception: pass

    def chart_left(sl, sec):
        ch = sec.get("chart", {}) or {}
        ctype = ch.get("type", "text")
        if ctype == "barh":
            src = REG.get(ch.get("source"))
            cat = ch.get("cat"); val = ch.get("val")
            if src is not None and cat in src.columns and val in src.columns:
                d = src[[cat, val]].dropna()
                cats = d[cat].tolist(); vals = d[val].tolist()
                if PLT is not None:
                    try:
                        fn = render_bar_png(cats, vals, ch.get("title"))
                        sl.shapes.add_picture(fn, Inches(0.5), Inches(1.78), width=Inches(5.6))
                        return
                    except Exception:
                        pass
                try:
                    _native_bar(sl, cats, vals, ch.get("title")); return
                except Exception:
                    pass
            textbox_left(sl, ch.get("title"), ["(차트 데이터를 불러올 수 없어 요약으로 대체합니다.)"])
        else:
            textbox_left(sl, ch.get("title"), ch.get("bullets", []))

    # ===== 표지 =====
    cover = prs.slides[0]; closing = prs.slides[2]; body = []
    es = com.get("executive_summary", {})
    meta = com.get("meta", {}); fp = meta.get("fingerprint", {}) or {}
    live_docs = len(v) if "v" in globals() else None
    stale = (fp.get("docs") is not None and live_docs is not None and fp.get("docs") != live_docs)
    _period = APP_DATE[:7].replace("-", ".") if isinstance(APP_DATE, str) and len(APP_DATE) >= 7 else "2026.06"
    cover_map = {"분석 주제를 입력하세요": "PASCUCCI 브랜드 컨슈머 인텔리전스 리포트",
                 "YYYY.MM": _period, "1.0": APP_VERSION, "Draft / Final": "Final", "Team / Owner": "PAS DIVISION"}
    for sh in cover.shapes:
        if not sh.has_text_frame: continue
        for p in sh.text_frame.paragraphs:
            if p.text.strip() in cover_map: setpara(p, cover_map[p.text.strip()])

    # ===== Executive Summary =====
    if es:
        s = dup(1); head(s, "EXECUTIVE SUMMARY", es.get("headline", "핵심 요약"))
        if stale:
            txt(s, 0.6, 1.55, SW - 1.2, 0.26, f"⚠ 데이터가 갱신됨(현재 {live_docs}건 ≠ 코멘트 {fp.get('docs')}건) — 코멘트 재생성 필요", 10, bold=True, color=RED)
        ky = 1.9 if not stale else 2.02
        txt(s, 0.6, ky, SW - 1.2, 0.34, "핵심 인사이트", 13, bold=True, color=NAVY)
        cards = es.get("key_insights", []); n = max(1, len(cards)); gap = 0.25
        cw = (SW - 1.2 - gap * (n - 1)) / n; cx = 0.6; cyy = ky + 0.42; chh = 1.9
        for c in cards:
            panel(s, cx, cyy, cw, chh, fill=CARD, line=CARDLINE, rounded=True)
            txt(s, cx + 0.14, cyy + 0.13, cw - 0.28, 0.3, c.get("area", ""), 10, bold=True, color=GOLD, sp=1)
            txt(s, cx + 0.14, cyy + 0.5, cw - 0.28, chh - 0.55, c.get("text", ""), 12, color=DARK, sp=1)
            cx += cw + gap
        ry0 = cyy + chh + 0.25
        txt(s, 0.6, ry0, SW - 1.2, 0.34, "Key Recommendations", 13, bold=True, color=RED)
        ry = ry0 + 0.43
        for it in es.get("recommendations", []):
            txt(s, 0.62, ry, SW - 1.25, 0.5, "· " + it, 12, color=DARK, sp=1)
            ry += 0.205 * max(1, math.ceil(len(it) / 56)) + 0.10
        body.append(s)

    # ===== 챕터별 섹션 =====
    chno = 0; cur = None
    for sec in com.get("sections", []):
        try:
            ch = sec.get("chapter", "")
            if ch != cur:
                cur = ch; chno += 1
                d = dup(1)
                txt(d, 0.6, 2.85, SW - 1.2, 0.5, f"CHAPTER {chno}", 16, bold=True, color=GOLD)
                txt(d, 0.6, 3.35, SW - 1.2, 0.9, ch, 30, bold=True, color=NAVY)
                subs = [x.get("tag", "").split(" · ")[-1] for x in com["sections"] if x.get("chapter") == ch]
                txt(d, 0.6, 4.45, SW - 1.2, 0.5, " · ".join(subs), 13, color=GRAY)
                body.append(d)
            s = dup(1); head(s, sec.get("tag", ch), sec.get("headline", ""))
            chart_left(s, sec)
            rightcol(s, sec)
            market_band(s, sec.get("market_view", ""))
            body.append(s)
        except Exception:
            continue

    if not body:
        s = dup(1); head(s, "NOTICE", "코멘트 파일을 찾지 못했습니다")
        textbox_left(s, "안내", ["reference/report_commentary.json 을 리포에 올리면 분석 서술이 반영됩니다.",
                                 "차트·수치는 라이브 데이터에서 생성됩니다."])
        body.append(s)

    reorder([cover] + body + [closing])
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.getvalue()


with report_dl_slot.container():
    st.markdown("**📄 분석 리포트 다운로드**")
    # [v2.3.4] PPT를 기본 리포트로(상단), 시사점·Lessons 포함
    if PPTX_AVAILABLE:
        try:
            pptx_bytes = build_summary_pptx()
            st.download_button(
                label="📊 분석 리포트 (.pptx · 권장)",
                data=pptx_bytes,
                file_name=f"PASCUCCI_분석리포트_{APP_DATE}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        except Exception as e:
            st.caption(f"pptx 생성 오류: {e}")
    else:
        st.caption("📊 PPT 리포트(권장): `python-pptx` 미설치 상태입니다. "
                   "`pip install python-pptx` 후 앱을 재시작하세요 "
                   "(requirements.txt 포함 — `pip install -r requirements.txt`).")
    # docx는 보조 포맷
    if DOCX_AVAILABLE:
        try:
            docx_bytes = build_summary_docx()
            st.download_button(
                label="📄 분석 리포트 (.docx · 보조)",
                data=docx_bytes,
                file_name=f"PASCUCCI_분석리포트_{APP_DATE}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
            )
        except Exception as e:
            st.caption(f"docx 생성 오류: {e}")

# [v2.3.1] MI 뷰 라우팅 — 대시보드가 아니면 MI 화면 표시 후 종료
if st.session_state.get("view", "dashboard") != "dashboard":
    render_market_intelligence(st.session_state["view"], v, brand_df)
    st.markdown("---")
    st.markdown(f"<p style='text-align:center;color:#A0AEC0;font-size:11px;'>{COPYRIGHT}</p>",
                unsafe_allow_html=True)
    st.stop()

# ══════════════════════════════════════════════════════
# 탭 구성 — 12개
# ══════════════════════════════════════════════════════
# [v2.5.1] 탭은 카테고리별 렌더 함수 + 하단 디스패치(_CATS)로 구성

# ════════════════════════════════════════════════════
# TAB 1 — Overview
# ════════════════════════════════════════════════════
def _rt_ov():
    # [v2.5.2] 파스쿠찌 중심 스냅샷 — 담당자가 들어오면 가장 먼저 보이는 패널
    _pb = brand_df[brand_df["브랜드"] == "파스쿠찌"]
    _tot = brand_df["브랜드"].nunique()
    if len(_pb):
        _row = _pb.iloc[0]
        _ment = int(_row["언급량"]); _nss = float(_row["NSS"])
        _sov_col = "SOV_전체" if "SOV_전체" in brand_df.columns else ("SOV" if "SOV" in brand_df.columns else None)
        _sov = float(_row[_sov_col]) if _sov_col else float("nan")
        _ment_rank = int((brand_df["언급량"] > _ment).sum()) + 1
        _nss_rank = int((brand_df["NSS"] > _nss).sum()) + 1
        _loy = None; _loy_rank = None; _loy_tot = 0
        try:
            if loyalty_df is not None and "파스쿠찌" in set(loyalty_df["브랜드"]):
                _lr = loyalty_df[loyalty_df["브랜드"] == "파스쿠찌"].iloc[0]
                _loy = float(_lr["로열티지수"]); _loy_tot = len(loyalty_df)
                _loy_rank = int((loyalty_df["로열티지수"] > _loy).sum()) + 1
        except Exception:
            pass
        st.markdown("""
        <div style='background:linear-gradient(135deg,#7A1F2B 0%,#A52A35 100%);
                    border-radius:14px;padding:16px 22px;margin-bottom:10px;color:#FFF;
                    box-shadow:0 4px 14px rgba(122,31,43,0.3);'>
            <span style='font-size:1.6rem;vertical-align:middle;'>☕</span>
            <span style='font-weight:800;font-size:1.25rem;letter-spacing:0.5px;vertical-align:middle;'>&nbsp;PASCUCCI Brand Status</span>
            <span style='font-size:0.8rem;opacity:0.88;'>&nbsp;— 핵심 지표 요약</span>
        </div>""", unsafe_allow_html=True)

        def _pcard(col, icon, val, lbl, sub):
            with col:
                st.markdown(f"""
                <div style='background:#FFF;border:1px solid #EBD9DB;border-top:4px solid #A52A35;
                            border-radius:12px;padding:14px 10px;text-align:center;'>
                    <div style='font-size:1.4rem;'>{icon}</div>
                    <div style='font-size:1.5rem;font-weight:800;color:#7A1F2B;line-height:1.2;'>{val}</div>
                    <div style='font-size:0.78rem;color:#4A5568;font-weight:700;'>{lbl}</div>
                    <div style='font-size:0.72rem;color:#A0AEC0;'>{sub}</div>
                </div>""", unsafe_allow_html=True)
        _h1, _h2, _h3, _h4 = st.columns(4)
        _pcard(_h1, "💬", f"{_ment:,}", "언급량", f"{_tot}개 중 {_ment_rank}위")
        _pcard(_h2, "😊", f"{_nss:+.1f}", "NSS(평판)", f"{_tot}개 중 {_nss_rank}위")
        _pcard(_h3, "📣", (f"{_sov:.1f}%" if _sov == _sov else "—"), "SOV(점유율)", "전체 언급 중 비중")
        _pcard(_h4, "💎", (f"{_loy:.1f}" if _loy is not None else "—"), "로열티지수",
               (f"{_loy_tot}개 중 {_loy_rank}위" if _loy_rank else "—"))
        st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)

    st.markdown("<div class='section-title'>📊 종합 현황 (시장 전체)</div>", unsafe_allow_html=True)
    _pas_arr = brand_df[brand_df['브랜드'] == '파스쿠찌']['NSS'].values
    _pas_nss_txt = f"NSS {_pas_arr[0]:+.1f}" if len(_pas_arr) else "NSS —"
    c1, c2, c3, c4, c5 = st.columns(5)
    metrics = [
        ("📝", f"{len(v):,}건",     "분석 문서 수"),
        ("🏢", f"{brand_df['브랜드'].nunique()}개사", "언급 브랜드"),
        ("😊", f"{(v['sentiment']=='긍정').mean()*100:.1f}%", "전체 긍정률"),
        ("😡", f"{(v['sentiment']=='부정').mean()*100:.1f}%", "전체 부정률"),
        ("☕", _pas_nss_txt, "파스쿠찌 NSS"),
    ]
    for col, (icon, val, lbl) in zip([c1,c2,c3,c4,c5], metrics):
        with col:
            st.markdown(f"""
            <div class='kpi-card'>
                <div class='kpi-icon'>{icon}</div>
                <div class='kpi-value'>{val}</div>
                <div class='kpi-label'>{lbl}</div>
            </div>""", unsafe_allow_html=True)

    st.markdown("")
    c1, c2 = st.columns([1.4, 1])
    with c1:
        fig = px.bar(
            brand_df.sort_values("언급량", ascending=False),
            x="브랜드", y="언급량", color="브랜드",
            color_discrete_map=BRAND_COLORS,
            text="언급량", title="브랜드별 언급량",
        )
        fig.update_traces(textposition="outside", textfont_size=11)
        chart_style(fig, showlegend=False)
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        sent_cnt = v["sentiment"].value_counts().reset_index()
        sent_cnt.columns = ["감성", "수"]
        fig2 = px.pie(sent_cnt, names="감성", values="수",
                      color="감성", color_discrete_map=SENT_COLORS,
                      hole=0.52, title="전체 감성 분포")
        fig2.update_traces(textinfo="percent+label", textfont_size=13)
        chart_style(fig2)
        st.plotly_chart(fig2, use_container_width=True)

    st.markdown("<div class='section-title'>브랜드별 월별 언급량 추이</div>", unsafe_allow_html=True)
    # 상위 6개 브랜드 월별 언급량 시계열
    top6_brands = brand_df.sort_values("언급량", ascending=False).head(6)["브랜드"].tolist()
    months_ov = sorted([m for m in v["year_month"].dropna().unique() if "NaT" not in str(m)])
    ts_rows_ov = []
    for b in top6_brands:
        bdocs = v[v["brands"].apply(lambda x: b in x)]
        for m in months_ov:
            cnt = int((bdocs["year_month"] == m).sum())
            ts_rows_ov.append({"브랜드": b, "월": m, "언급량": cnt})
    ts_ov_df = pd.DataFrame(ts_rows_ov)
    fig_ts = px.line(
        ts_ov_df, x="월", y="언급량", color="브랜드", markers=True,
        color_discrete_map=BRAND_COLORS,
        title="상위 6개 브랜드 월별 언급량 추이",
    )
    fig_ts.update_traces(line=dict(width=2.5), marker=dict(size=6))
    chart_style(fig_ts, height=380)
    st.plotly_chart(fig_ts, use_container_width=True)

    st.markdown("<div class='section-title'>브랜드별 감성 누적 비율</div>", unsafe_allow_html=True)
    sp = brand_df.melt(id_vars="브랜드", value_vars=["긍정률","부정률"],
                       var_name="감성", value_name="비율")
    sp["감성"] = sp["감성"].str.replace("률","")
    fig3 = px.bar(
        brand_df.sort_values("NSS"),
        x="브랜드", y=["긍정률","부정률"],
        title="브랜드별 긍정/부정 비율 (%)",
        color_discrete_map={"긍정률": "#27AE60", "부정률": "#E74C3C"},
        labels={"value":"비율(%)","variable":"감성"},
        text_auto=".0f",
    )
    chart_style(fig3, height=350)
    st.plotly_chart(fig3, use_container_width=True)

    st.markdown("<div class='section-title'>브랜드 종합 지표 테이블</div>", unsafe_allow_html=True)
    st.dataframe(
        brand_df[["브랜드","언급량","긍정","부정","중립","긍정률","부정률","NSS","SOV_전체","SOV_긍정","SOV_부정"]]
        .sort_values("NSS", ascending=False),
        use_container_width=True, hide_index=True,
    )

# ════════════════════════════════════════════════════
# TAB 2 — NSS 평판
# ════════════════════════════════════════════════════
def _rt_nss():
    st.markdown("<div class='section-title'>📈 Net Sentiment Score 평판 진단</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    <b>NSS = (긍정 − 부정) / 전체 × 100</b><br>
    단순 긍정률보다 평판의 '방향성'을 민감하게 포착합니다.
    +100에 가까울수록 압도적 긍정, 0 이하면 부정이 긍정을 추월한 상태입니다.
    </div>""", unsafe_allow_html=True)

    c1, c2 = st.columns([1.2, 1])
    with c1:
        sorted_b = brand_df.sort_values("NSS")
        fig = px.bar(
            sorted_b, x="NSS", y="브랜드", orientation="h",
            color="NSS", color_continuous_scale=["#E74C3C","#FFFFFF","#27AE60"],
            color_continuous_midpoint=0,
            title="브랜드별 NSS",
            text=sorted_b["NSS"].apply(lambda x: f"{x:+.1f}"),
        )
        fig.update_traces(textposition="outside")
        fig.add_vline(x=0, line_color="#E53E3E", line_width=2)
        fig.add_vline(x=brand_df["NSS"].mean(), line_color="#718096",
                      line_width=1.5, line_dash="dot",
                      annotation_text=f"평균 {brand_df['NSS'].mean():.1f}",
                      annotation_position="top right")
        chart_style(fig, height=400, showlegend=False)
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        for _, row in brand_df.sort_values("NSS", ascending=False).iterrows():
            nss = row["NSS"]
            color = BRAND_COLORS.get(row["브랜드"], "#2D6BC4")
            grade = "🟢 우수" if nss >= 50 else ("🟡 양호" if nss >= 0 else "🔴 위험")
            st.markdown(f"""
            <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-left:4px solid {color};
                        border-radius:0 8px 8px 0;padding:9px 12px;margin-bottom:6px;
                        display:flex;justify-content:space-between;align-items:center;'>
                <span style='font-weight:700;color:{color};font-size:0.9rem;'>{row['브랜드']}</span>
                <span style='color:#4A5568;font-size:0.8rem;'>긍 {row['긍정률']}% / 부 {row['부정률']}%</span>
                <span style='font-weight:800;color:{color};font-size:1rem;'>NSS {nss:+.1f}</span>
                <span>{grade}</span>
            </div>""", unsafe_allow_html=True)

    st.markdown("<div class='section-title'>월별 NSS 트렌드 (주요 6개 브랜드)</div>", unsafe_allow_html=True)
    if len(monthly_nss_df) > 0:
        fig_t = px.line(
            monthly_nss_df, x="월", y="NSS", color="브랜드",
            color_discrete_map=BRAND_COLORS, markers=True,
            title="브랜드별 월별 NSS 추이",
            labels={"월": "월", "NSS": "Net Sentiment Score"},
        )
        fig_t.add_hline(y=0, line_color="#E53E3E", line_width=1.5, line_dash="dash")
        chart_style(fig_t, height=380)
        st.plotly_chart(fig_t, use_container_width=True)

        st.markdown("<div class='section-title'>5월 스타벅스 NSS 급락 분석</div>", unsafe_allow_html=True)
        col1, col2 = st.columns(2)
        with col1:
            sbux_m = monthly_nss_df[monthly_nss_df["브랜드"] == "스타벅스"]
            if not sbux_m.empty:
                prev_avg = sbux_m[sbux_m["월"] < "2026-05"]["NSS"].mean()
                may_nss  = sbux_m[sbux_m["월"] == "2026-05"]["NSS"].values
                st.markdown(f"""
                <div class='danger-box'>
                <b>스타벅스 2026.05 NSS = {may_nss[0] if len(may_nss) > 0 else 'N/A':+.1f}</b><br>
                이전 8개월 평균 NSS: <b>+{prev_avg:.1f}</b><br>
                → 낙폭 <b>{prev_avg - (may_nss[0] if len(may_nss) > 0 else 0):.1f}p</b> 급전직하
                </div>""", unsafe_allow_html=True)
        with col2:
            # [기능 1] 스타벅스 외 다양한 시각의 코멘트 (파스쿠찌 시사점 삭제)
            # 데이터 기반으로 NSS 상·하위 브랜드를 자동 추출해 다각도 코멘트 생성
            nss_sorted = brand_df.sort_values("NSS", ascending=False)
            top_brand = nss_sorted.iloc[0]["브랜드"]
            top_nss = nss_sorted.iloc[0]["NSS"]
            bottom_brand = nss_sorted.iloc[-1]["브랜드"]
            bottom_nss = nss_sorted.iloc[-1]["NSS"]
            mid_avg = brand_df["NSS"].mean()
            st.markdown(f"""
            <div class='insight-box'>
            <b>시장 전반 관점</b><br>
            • 평판 1위: <b>{top_brand}</b> (NSS {top_nss:+.1f}) — 긍정 담론 주도<br>
            • 평판 최하위: <b>{bottom_brand}</b> (NSS {bottom_nss:+.1f}) — 리스크 관리 필요<br>
            • 전체 브랜드 평균 NSS: <b>{mid_avg:+.1f}</b>
            </div>""", unsafe_allow_html=True)

        # [기능 1] 추가 다각도 코멘트
        st.markdown("<div class='section-title'>다각도 시사점</div>", unsafe_allow_html=True)
        ic1, ic2, ic3 = st.columns(3)
        with ic1:
            st.markdown("""
            <div class='success-box'>
            <b>가성비 브랜드 관점</b><br>
            메가·컴포즈·빽다방 등 저가 브랜드는 가격 민감 소비자의 지지로
            안정적 NSS를 유지하는 경향. 불황기 수혜 가능성.
            </div>""", unsafe_allow_html=True)
        with ic2:
            st.markdown("""
            <div class='warning-box'>
            <b>프리미엄 브랜드 관점</b><br>
            테라로사·폴바셋 등은 품질 담론으로 높은 NSS를 보이나
            언급량(SOV)이 작아 영향력 확대가 과제.
            </div>""", unsafe_allow_html=True)
        with ic3:
            st.markdown("""
            <div class='insight-box'>
            <b>이슈 관리 관점</b><br>
            정치·사회 이슈는 단기 NSS 급락을 유발. 브랜드별 이슈
            모니터링과 선제적 커뮤니케이션 체계가 중요.
            </div>""", unsafe_allow_html=True)

# ════════════════════════════════════════════════════
# TAB 3 — SOV 분석
# ════════════════════════════════════════════════════
def _rt_sov():
    st.markdown("<div class='section-title'>📣 Share of Voice 분석</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    <b>SOV(담론 점유율)</b>: 전체 vs 긍정 vs 부정 SOV를 분리해 브랜드별 담론 질을 진단합니다.<br>
    전체 SOV가 높아도 부정 SOV가 압도적이면 실질 경쟁력이 낮은 상태입니다.
    </div>""", unsafe_allow_html=True)

    c1, c2, c3 = st.columns(3)
    with c1:
        fig = px.pie(brand_df, names="브랜드", values="SOV_전체",
                     color="브랜드", color_discrete_map=BRAND_COLORS,
                     title="전체 SOV", hole=0.45)
        fig.update_traces(textinfo="percent+label", textfont_size=11)
        chart_style(fig, height=340)
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        sorted_sov = brand_df.sort_values("SOV_긍정", ascending=True)
        fig2 = px.bar(sorted_sov, x="SOV_긍정", y="브랜드", orientation="h",
                      color="브랜드", color_discrete_map=BRAND_COLORS,
                      title="긍정 SOV (%)",
                      text=sorted_sov["SOV_긍정"].apply(lambda x: f"{x:.1f}%"))
        fig2.update_traces(textposition="outside")
        chart_style(fig2, height=340, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)
    with c3:
        sorted_sov2 = brand_df.sort_values("SOV_부정", ascending=True)
        fig3 = px.bar(sorted_sov2, x="SOV_부정", y="브랜드", orientation="h",
                      color="브랜드", color_discrete_map=BRAND_COLORS,
                      title="부정 SOV (%)",
                      text=sorted_sov2["SOV_부정"].apply(lambda x: f"{x:.1f}%"))
        fig3.update_traces(textposition="outside")
        chart_style(fig3, height=340, showlegend=False)
        st.plotly_chart(fig3, use_container_width=True)

    st.markdown("<div class='section-title'>SOV 매트릭스 — 전체 SOV × 긍정 SOV (버블=부정 SOV)</div>", unsafe_allow_html=True)
    fig4 = go.Figure()
    for _, row in brand_df.iterrows():
        color = BRAND_COLORS.get(row["브랜드"], "#2D6BC4")
        fig4.add_trace(go.Scatter(
            x=[row["SOV_전체"]], y=[row["SOV_긍정"]],
            mode="markers+text", name=row["브랜드"],
            marker=dict(
                size=max(row["SOV_부정"] * 6, 12),
                color=color, opacity=0.8,
                line=dict(color="white", width=2)
            ),
            text=[row["브랜드"]], textposition="top center",
            textfont=dict(size=11, color=color, family="Arial Black"),
            hovertemplate=(
                f"<b>{row['브랜드']}</b><br>"
                f"전체 SOV: {row['SOV_전체']:.1f}%<br>"
                f"긍정 SOV: {row['SOV_긍정']:.1f}%<br>"
                f"부정 SOV: {row['SOV_부정']:.1f}%<extra></extra>"
            ),
        ))
    fig4.add_hline(y=brand_df["SOV_긍정"].mean(), line_dash="dot", line_color="#A0AEC0")
    fig4.add_vline(x=brand_df["SOV_전체"].mean(), line_dash="dot", line_color="#A0AEC0")
    fig4.update_layout(
        title="SOV 매트릭스 (버블 크기 = 부정 SOV)",
        xaxis_title="전체 담론 점유율 (%)",
        yaxis_title="긍정 담론 점유율 (%)",
        showlegend=False,
    )
    chart_style(fig4, height=450)
    st.plotly_chart(fig4, use_container_width=True)

    col1, col2 = st.columns(2)
    with col1:
        st.markdown("""
        <div class='danger-box'><b>스타벅스 딜레마</b><br>
        전체 SOV 53.3%로 독보적이나, 부정 SOV 77.7%로 부정 담론 독식.
        담론의 양과 질이 극명하게 괴리되어 있음.</div>""", unsafe_allow_html=True)
    with col2:
        st.markdown("""
        <div class='success-box'><b>파스쿠찌 진입 기회</b><br>
        전체 SOV 2~5%, 긍정 SOV 5~7% 구간(테라로사·이디야·폴바셋 수준)을
        목표로 설정하면 작은 투자로 큰 긍정 담론 효과 가능.</div>""", unsafe_allow_html=True)

# ════════════════════════════════════════════════════
# TAB 4 — 드라이버 분석 (ABSA)
# ════════════════════════════════════════════════════
def _rt_absa():
    st.markdown("<div class='section-title'>🔍 평판 드라이버 분석 (Aspect-Based Sentiment)</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    <b>ABSA</b>: 속성(맛·배송·가격 등)별로 긍정/부정 리뷰를 분해합니다.
    '왜 좋은가? 왜 싫은가?'를 구체적인 속성 수준에서 파악합니다.
    </div>""", unsafe_allow_html=True)

    sel_b = st.selectbox("브랜드 선택 (전체 또는 개별)", ["전체"] + sorted(TARGET_BRANDS), key="absa_brand")
    if sel_b == "전체":
        absa_target = v
    else:
        absa_target = v[v["brands"].apply(lambda x: sel_b in x)]

    absa_computed = compute_absa(absa_target)

    c1, c2 = st.columns(2)
    with c1:
        fig = go.Figure()
        df_absa = absa_computed.sort_values("NSS")
        fig.add_trace(go.Bar(
            name="긍정", y=df_absa["속성"], x=df_absa["긍정률"],
            orientation="h", marker_color="#27AE60",
            text=df_absa["긍정률"].apply(lambda x: f"{x:.0f}%"),
            textposition="inside", textfont=dict(color="white", size=11),
        ))
        fig.add_trace(go.Bar(
            name="부정", y=df_absa["속성"], x=-df_absa["부정률"],
            orientation="h", marker_color="#E74C3C",
            text=df_absa["부정률"].apply(lambda x: f"{x:.0f}%"),
            textposition="inside", textfont=dict(color="white", size=11),
        ))
        fig.update_layout(
            barmode="relative",
            title=f"속성별 긍정/부정 비율 — {sel_b}",
            xaxis=dict(
                tickvals=[-100,-75,-50,-25,0,25,50,75,100],
                ticktext=["100%","75%","50%","25%","0","25%","50%","75%","100%"],
                title="비율 (%)"
            ),
        )
        chart_style(fig, height=400)
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        fig2 = px.bar(
            absa_computed.sort_values("NSS"), x="NSS", y="속성",
            orientation="h", color="NSS",
            color_continuous_scale=["#E74C3C","#FFFFFF","#27AE60"],
            color_continuous_midpoint=0,
            title=f"속성별 NSS — {sel_b}",
            text=absa_computed.sort_values("NSS")["NSS"].apply(lambda x: f"{x:+.1f}"),
        )
        fig2.update_traces(textposition="outside")
        fig2.add_vline(x=0, line_color="#4A5568", line_width=1.5)
        chart_style(fig2, height=400, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)

    st.markdown("<div class='section-title'>브랜드별 속성 NSS 히트맵</div>", unsafe_allow_html=True)
    heat = {}
    for b in TARGET_BRANDS:
        bd = v[v["brands"].apply(lambda x: b in x)]
        row = {}
        for attr, kws in ATTR_DICT.items():
            pat  = "|".join(kws)
            sub  = bd[bd["full_text"].str.contains(pat, na=False)]
            if len(sub) < 2:
                row[attr] = 0
                continue
            p = (sub["sentiment"] == "긍정").sum()
            n = (sub["sentiment"] == "부정").sum()
            t = len(sub)
            row[attr] = round((p - n) / t * 100, 1)
        heat[b] = row
    heat_df = pd.DataFrame(heat).T
    fig3 = px.imshow(
        heat_df,
        color_continuous_scale=["#E74C3C","#FFFFFF","#27AE60"],
        color_continuous_midpoint=0,
        title="브랜드 × 속성 NSS 히트맵 (빨강=부정 강함, 초록=긍정 강함)",
        text_auto=".1f", aspect="auto", zmin=-60, zmax=80,
    )
    fig3.update_traces(textfont=dict(color="#1B2A4A", size=11))
    chart_style(fig3, height=360, showlegend=False)
    st.plotly_chart(fig3, use_container_width=True)

    st.markdown("<div class='section-title'>칭찬 vs 불만 드라이버 TOP 3</div>", unsafe_allow_html=True)
    col1, col2 = st.columns(2)
    with col1:
        for _, row in absa_computed.nlargest(3, "NSS").iterrows():
            st.markdown(f"""
            <div style='background:#F0FFF4;border-left:4px solid #27AE60;border-radius:0 8px 8px 0;
                        padding:10px 14px;margin:6px 0'>
                <b style='color:#276749'>{row['속성']}</b>
                <span style='float:right;color:#27AE60;font-weight:700'>NSS {row['NSS']:+.1f} / 긍정 {row['긍정률']}%</span>
                <div style='color:#718096;font-size:0.82rem;margin-top:4px'>
                    분석 {row['문서수']}건 · 긍정 {row['긍정']}건 · 부정 {row['부정']}건
                </div>
            </div>""", unsafe_allow_html=True)
    with col2:
        for _, row in absa_computed.nsmallest(3, "NSS").iterrows():
            st.markdown(f"""
            <div style='background:#FFF5F5;border-left:4px solid #E53E3E;border-radius:0 8px 8px 0;
                        padding:10px 14px;margin:6px 0'>
                <b style='color:#C53030'>{row['속성']}</b>
                <span style='float:right;color:#E74C3C;font-weight:700'>NSS {row['NSS']:+.1f} / 부정 {row['부정률']}%</span>
                <div style='color:#718096;font-size:0.82rem;margin-top:4px'>
                    분석 {row['문서수']}건 · 긍정 {row['긍정']}건 · 부정 {row['부정']}건
                </div>
            </div>""", unsafe_allow_html=True)

    # [v2.4.1] 칭찬/불만 드라이버 Top3 원문 확인
    st.markdown("<div class='section-title'>🔎 드라이버 Top 3 원문 확인</div>", unsafe_allow_html=True)
    st.caption("각 드라이버를 펼치면 해당 속성 키워드가 포함된 문장만 발췌해 보여줍니다(전체 글이 아닌 관련 부분). 칭찬/불만은 그 문장에 감성어가 함께 있는 글을 우선 노출합니다. [v2.5.1 정확도 개선]")

    _DRV_POS = ["맛있", "좋", "추천", "최고", "만족", "친절", "합리", "가성비", "저렴",
                "달콤", "훌륭", "대박", "꿀", "인생", "강추", "최애"]
    _DRV_NEG = ["맛없", "별로", "실망", "불만", "최악", "비싸", "불친절", "느리", "아쉽",
                "안좋", "싫", "불매", "후회", "환불", "논란", "줄였", "항의"]

    def _aspect_snippet(text, kws):
        """키워드가 포함된 문장만 발췌 — 전체 문서 대신 속성 관련 부분만 노출."""
        t = str(text)
        for s in re.split(r"(?<=[.!?。…])\s+|\n+", t):
            if any(kw in s for kw in kws):
                s2 = re.sub(r"\s+", " ", s).strip()
                if len(s2) >= 6:
                    return s2[:170]
        for kw in kws:
            i = t.find(kw)
            if i >= 0:
                return re.sub(r"\s+", " ", t[max(0, i - 40):i + 90]).strip()[:170]
        return ""

    def _driver_docs(attr, sentiment, k=3):
        # [v2.5.1] 정확도 개선: 속성 키워드가 든 문장을 발췌하고, 그 문장에 해당 감성어가
        # 함께 있을 때만 그 속성에 대한 칭찬/불만으로 채택. 관련도 높은 순으로 정렬.
        kws = ATTR_DICT.get(attr, [])
        if not kws:
            return []
        senti_kw = _DRV_POS if sentiment == "긍정" else _DRV_NEG
        strict, loose = [], []
        for _, rr in absa_target.iterrows():
            ft = str(rr.get("full_text", ""))
            if not any(kw in ft for kw in kws):
                continue
            snip = _aspect_snippet(ft, kws)
            if not snip:
                continue
            title = str(rr.get("title", "")).strip() or "(제목 없음)"
            url = str(rr.get("post_url", "") or "")
            relev = sum(ft.count(kw) for kw in kws) + (3 if any(kw in title for kw in kws) else 0)
            item = (relev, title, snip, url)
            (strict if any(sw in snip for sw in senti_kw) else loose).append(item)
        strict.sort(key=lambda x: -x[0])
        loose.sort(key=lambda x: -x[0])
        picked = strict[:k] + loose[:max(0, k - len(strict))]
        return [(t, s, u) for _, t, s, u in picked]

    dcol1, dcol2 = st.columns(2)
    with dcol1:
        st.markdown("**👍 칭찬 드라이버 원문**")
        for _, row in absa_computed.nlargest(3, "NSS").iterrows():
            with st.expander(f"{row['속성']} · NSS {row['NSS']:+.1f}"):
                docs = _driver_docs(row['속성'], "긍정")
                if not docs:
                    st.caption("표시할 원문이 충분치 않습니다.")
                for title, body, url in docs:
                    link = f" [원문 보기]({url})" if url.startswith("http") else ""
                    st.markdown(f"- **{title}** — {body}…{link}")
    with dcol2:
        st.markdown("**👎 불만 드라이버 원문**")
        for _, row in absa_computed.nsmallest(3, "NSS").iterrows():
            with st.expander(f"{row['속성']} · NSS {row['NSS']:+.1f}"):
                docs = _driver_docs(row['속성'], "부정")
                if not docs:
                    st.caption("표시할 원문이 충분치 않습니다.")
                for title, body, url in docs:
                    link = f" [원문 보기]({url})" if url.startswith("http") else ""
                    st.markdown(f"- **{title}** — {body}…{link}")

# ════════════════════════════════════════════════════
# TAB 5 — CA 이미지 맵
# ════════════════════════════════════════════════════
def _rt_ca():
    st.markdown("<div class='section-title'>🗺️ 브랜드 이미지 맵 — 대응일치분석 (CA)</div>", unsafe_allow_html=True)
    if ca_result is None:
        st.info("CA 분석을 위한 데이터가 부족합니다.")
    else:
        n_kw = st.slider("표시 키워드 수", 15, 40, 30, key="ca_kw")
        expl  = ca_result["expl"]
        brands = ca_result["brands"]
        top_kws = ca_result["kws"][:n_kw]
        cont    = ca_result["cont"][top_kws]
        kw_sizes = [cont[kw].sum() for kw in top_kws]
        kw_max   = max(kw_sizes) if kw_sizes else 1

        col_x = ca_result["col_x"][:n_kw]
        col_y = ca_result["col_y"][:n_kw]

        fig = go.Figure()
        fig.add_trace(go.Scatter(
            x=col_x, y=col_y, mode="markers+text",
            marker=dict(
                size=[10 + int(s / kw_max * 18) for s in kw_sizes],
                color="#A0AEC0", opacity=0.65,
                line=dict(color="white", width=1)
            ),
            text=top_kws, textposition="top center",
            textfont=dict(size=10, color="#4A5568"),
            hovertemplate="<b>%{text}</b><extra>키워드</extra>",
            name="키워드",
        ))
        for i, b in enumerate(brands):
            color = BRAND_COLORS.get(b, "#2D6BC4")
            rx, ry = ca_result["row_x"][i], ca_result["row_y"][i]
            dists = [(np.sqrt((col_x[j]-rx)**2+(col_y[j]-ry)**2), j) for j in range(len(top_kws))]
            for _, j in sorted(dists)[:4]:
                fig.add_trace(go.Scatter(
                    x=[rx, col_x[j]], y=[ry, col_y[j]],
                    mode="lines", line=dict(color=color, width=1.2, dash="dot"),
                    opacity=0.35, showlegend=False, hoverinfo="skip",
                ))
            _is_pas = (b == "파스쿠찌")
            fig.add_trace(go.Scatter(
                x=[rx], y=[ry], mode="markers+text", name=b,
                marker=dict(size=36 if _is_pas else 22,
                            symbol="star" if _is_pas else "circle",
                            color=color,
                            line=dict(color="#C0392B" if _is_pas else "white",
                                      width=4 if _is_pas else 2.5)),
                text=[("★ " + b) if _is_pas else b], textposition="top right",
                textfont=dict(size=16 if _is_pas else 12,
                              color="#C0392B" if _is_pas else color, family="Arial Black"),
                hovertemplate=f"<b>{b}</b><extra>브랜드</extra>",
            ))
        fig.add_hline(y=0, line_color="#CBD5E0", line_dash="dot", line_width=1)
        fig.add_vline(x=0, line_color="#CBD5E0", line_dash="dot", line_width=1)
        # [v2.5.2] 축 라벨을 데이터(양 극단 키워드) 기반으로 — CA 차원 부호는 데이터마다 달라짐
        _cx = np.asarray(col_x); _cy = np.asarray(col_y)
        _ox = np.argsort(_cx); _oy = np.argsort(_cy)
        _xl = " · ".join(top_kws[i] for i in _ox[:3])
        _xr = " · ".join(top_kws[i] for i in _ox[-3:][::-1])
        _yb = " · ".join(top_kws[i] for i in _oy[:3])
        _yt = " · ".join(top_kws[i] for i in _oy[-3:][::-1])
        fig.update_layout(
            title=f"브랜드 이미지 맵  (차원1: {expl[0]:.1f}% | 차원2: {expl[1]:.1f}% | 누적: {sum(expl[:2]):.1f}%)",
            xaxis_title=f"차원 1 ({expl[0]:.1f}%):  ← {_xl}　|　{_xr} →",
            yaxis_title=f"차원 2 ({expl[1]:.1f}%):  ↓ {_yb}　|　{_yt} ↑",
        )
        chart_style(fig, height=620)
        st.plotly_chart(fig, use_container_width=True)

        with st.expander("📘 브랜드 이미지 맵 읽는 방법"):
            st.markdown("""
            - **브랜드 ●** 와 **키워드 ●** 가 **가까울수록** 해당 브랜드 담론에서 그 키워드가 자주 등장합니다.
            - **축의 의미는 고정이 아니라 데이터로 결정됩니다.** 각 축 양 끝의 키워드(축 제목에 표시)가 그 방향의 성격이며, 데이터가 바뀌면 축 방향·해석도 달라집니다.
            - 원점(0,0)에서 **멀수록 개성이 뚜렷한** 이미지, 가까울수록 평균적 이미지입니다.
            - 점선 = 브랜드와 가장 가까운 키워드 TOP 4 연결.
            - **★ 파스쿠찌**는 별 마커로 강조됩니다 — 어떤 키워드 쪽에 위치하는지 먼저 확인하세요.
            """)

# ════════════════════════════════════════════════════
# TAB 6 — LDA 토픽 모델링
# ════════════════════════════════════════════════════
def _rt_lda():
    st.markdown("<div class='section-title'>🎯 LDA 토픽 모델링 (7-토픽)</div>", unsafe_allow_html=True)
    if lda_model is None:
        st.warning("LDA 분석을 위한 데이터가 부족합니다.")
    else:
        sel_b6 = st.selectbox("브랜드 선택", ["전체"] + sorted(TARGET_BRANDS), key="lda_brand")
        if sel_b6 == "전체":
            lda_target = v
        else:
            lda_target = v[v["brands"].apply(lambda x: sel_b6 in x)]

        lm, tv, vc, dt, tdist = compute_lda(lda_target)
        if lm is None:
            st.info("선택한 브랜드의 문서가 부족합니다.")
        else:
            c1, c2 = st.columns([1.4, 1])
            with c1:
                fig = make_subplots(
                    rows=4, cols=2,
                    subplot_titles=[f"T{i+1}: {TOPIC_LABELS[i]}" for i in range(7)],
                    vertical_spacing=0.08, horizontal_spacing=0.1,
                )
                for ti in range(7):
                    r, ci = ti // 2 + 1, ti % 2 + 1
                    top_idx = lm.components_[ti].argsort()[-10:][::-1]
                    tw = [vc[j] for j in top_idx]
                    tv_vals = lm.components_[ti][top_idx]
                    tv_vals = tv_vals / tv_vals.sum()
                    fig.add_trace(
                        go.Bar(x=tv_vals, y=tw, orientation="h",
                               marker_color=TOPIC_COLORS[ti], showlegend=False,
                               hovertemplate="%{y}: %{x:.3f}<extra></extra>"),
                        row=r, col=ci,
                    )
                    fig.update_yaxes(autorange="reversed", row=r, col=ci,
                                     tickfont=dict(size=9, color="#2D3748"))
                    fig.update_xaxes(tickfont=dict(size=8, color="#718096"), row=r, col=ci)
                fig.update_layout(
                    height=700, paper_bgcolor="#F7F9FC", plot_bgcolor="#FFFFFF",
                    font=dict(color="#2D3748"),
                    title_text=f"{sel_b6} — 7개 토픽 키워드",
                    title_font=dict(color="#1B2A4A", size=14),
                )
                for ann in fig.layout.annotations:
                    ann.font.color = "#1B2A4A"; ann.font.size = 11
                st.plotly_chart(fig, use_container_width=True)

            with c2:
                dist = dt.mean(axis=0)
                fig2 = go.Figure(go.Pie(
                    labels=[f"T{i+1}: {TOPIC_LABELS[i]}" for i in range(7)],
                    values=dist, hole=0.5,
                    marker=dict(colors=TOPIC_COLORS, line=dict(color="white", width=2)),
                    textfont=dict(color="#1B2A4A", size=10),
                ))
                fig2.update_layout(title=f"{sel_b6} — 토픽 비중")
                chart_style(fig2, height=360)
                st.plotly_chart(fig2, use_container_width=True)

                dom = dt.argmax(axis=1)
                dom_cnt = pd.Series(dom).value_counts().sort_index()
                fig3 = px.bar(
                    x=[TOPIC_LABELS[i] for i in dom_cnt.index],
                    y=dom_cnt.values,
                    color=[TOPIC_LABELS[i] for i in dom_cnt.index],
                    color_discrete_sequence=TOPIC_COLORS,
                    title="지배 토픽별 문서 수",
                    text=dom_cnt.values,
                    labels={"x":"토픽","y":"문서 수"},
                )
                fig3.update_traces(textposition="outside")
                chart_style(fig3, height=300, showlegend=False)
                st.plotly_chart(fig3, use_container_width=True)

        if sel_b6 == "전체":
            st.markdown("<div class='section-title'>브랜드별 토픽 분포 히트맵</div>", unsafe_allow_html=True)
            brand_topic = {}
            for b in TARGET_BRANDS:
                bd = v[v["brands"].apply(lambda x: b in x)]
                bd2 = bd[bd["tokens_str"].str.len() > 5]
                if len(bd2) < 5:
                    continue
                bX = tfidf_vec.transform(bd2["tokens_str"])
                bt = lda_model.transform(bX).mean(axis=0)
                brand_topic[b] = {TOPIC_LABELS[i]: round(float(bt[i])*100, 1) for i in range(7)}
            if brand_topic:
                bt_df = pd.DataFrame(brand_topic).T
                fig4 = px.imshow(
                    bt_df, color_continuous_scale="Blues",
                    title="브랜드 × 토픽 집중도 히트맵 (%)",
                    text_auto=".1f", aspect="auto",
                )
                fig4.update_traces(textfont=dict(color="#1B2A4A", size=11))
                chart_style(fig4, height=380, showlegend=False)
                st.plotly_chart(fig4, use_container_width=True)

# ════════════════════════════════════════════════════
# TAB 7 — Burst 탐지
# ════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("<div class='section-title'>🧩 토픽 클러스터링 — 토픽 × 브랜드 × 감성 중첩 분석</div>", unsafe_allow_html=True)
    st.caption("※ LDA 토픽 모델링과 토픽 클러스터링을 한 화면에서 함께 분석 [v2.3.1 merge]")
    st.markdown("""
    <div class='insight-box'>
    LDA 토픽을 나열식으로 보지 않고, <b>어떤 토픽이 어떤 브랜드·고객 감성과 중첩되는지</b>를 분석합니다.<br>
    토픽 × 브랜드 히트맵으로 "어떤 브랜드가 어떤 주제로 이야기되는가"를 파악합니다.
    </div>""", unsafe_allow_html=True)

    if topic_cluster is None:
        st.info("토픽 클러스터링을 위한 데이터가 부족합니다.")
    else:
        n_t = topic_cluster["n_topics"]
        topic_words = topic_cluster["topic_words"]
        topic_meanings = topic_cluster.get("topic_meanings", [""]*n_t)
        topic_occasions = topic_cluster.get("topic_occasions", [[]]*n_t)
        topic_brands_list = topic_cluster.get("topic_brands", [[]]*n_t)
        topic_labels_cl = [f"T{i}: {topic_meanings[i]}" for i in range(n_t)]

        st.markdown("<div class='section-title'>토픽별 의미 · 핵심 키워드 · 소비자 유형</div>", unsafe_allow_html=True)
        tcols = st.columns(n_t)
        cluster_colors = ["#2980B9", "#16A085", "#E67E22", "#8E44AD", "#E74C3C",
                          "#D4A017", "#C0392B", "#117A65"]  # [v2.3.2] 6~8토픽 대응
        for i, col in enumerate(tcols):
            with col:
                kw_html = "".join([
                    f"<span style='display:inline-block;background:{cluster_colors[i % len(cluster_colors)]}22;"
                    f"color:{cluster_colors[i % len(cluster_colors)]};border-radius:10px;padding:2px 8px;"
                    f"margin:2px;font-size:11px;'>{w}</span>"
                    for w in topic_words[i][:6]])
                occ_str = ", ".join(topic_occasions[i]) if topic_occasions[i] else "—"
                brand_str = ", ".join(topic_brands_list[i]) if topic_brands_list[i] else "—"
                st.markdown(f"""
                <div style='background:#FFFFFF;border-top:3px solid {cluster_colors[i % len(cluster_colors)]};
                            border:1px solid #E2E8F0;border-radius:8px;padding:10px;'>
                    <div style='font-weight:800;color:{cluster_colors[i % len(cluster_colors)]};font-size:0.9rem;'>토픽 {i}</div>
                    <div style='color:#1B2A4A;font-weight:700;font-size:0.82rem;margin:3px 0;'>{topic_meanings[i]}</div>
                    <div style='color:#718096;font-size:11px;'>{topic_cluster['topic_sizes'][i]}건</div>
                    <div style='margin-top:6px;'>{kw_html}</div>
                    <div style='margin-top:8px;padding-top:6px;border-top:1px dashed #E2E8F0;font-size:0.72rem;color:#4A5568;'>
                        <b>주요 소비상황</b>: {occ_str}<br>
                        <b>주요 브랜드</b>: {brand_str}
                    </div>
                </div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>토픽 × 브랜드 중첩 히트맵</div>", unsafe_allow_html=True)
        tb = topic_cluster["tb_mat"].copy()
        tb.index = topic_labels_cl
        fig = px.imshow(
            tb, color_continuous_scale="Blues",
            title="어떤 브랜드가 어떤 토픽으로 이야기되는가",
            text_auto=True, aspect="auto",
        )
        fig.update_traces(textfont=dict(color="#1B2A4A", size=11))
        chart_style(fig, height=380, showlegend=False)
        st.plotly_chart(fig, use_container_width=True)

        c1, c2 = st.columns(2)
        with c1:
            st.markdown("<div class='section-title'>토픽 × 감성 분포</div>", unsafe_allow_html=True)
            ts = topic_cluster["ts_mat"].copy()
            ts.index = [f"T{i}" for i in range(n_t)]
            fig2 = px.bar(
                ts, barmode="stack",
                color_discrete_map={"긍정":"#27AE60","중립":"#F39C12","부정":"#E74C3C"},
                title="토픽별 감성 구성",
                labels={"value":"문서 수","index":"토픽","variable":"감성"},
            )
            chart_style(fig2, height=340)
            st.plotly_chart(fig2, use_container_width=True)
        with c2:
            st.markdown("<div class='section-title'>토픽 규모 비중</div>", unsafe_allow_html=True)
            fig3 = go.Figure(go.Pie(
                labels=topic_labels_cl, values=topic_cluster["topic_sizes"],
                hole=0.5,
                marker=dict(colors=cluster_colors[:n_t], line=dict(color="white", width=2)),
                textfont=dict(size=10),
            ))
            fig3.update_layout(title="토픽 규모 분포")
            chart_style(fig3, height=340)
            st.plotly_chart(fig3, use_container_width=True)

        st.markdown("""
        <div class='success-box'>
        <b>활용</b>: 특정 토픽이 경쟁사에 집중되어 있다면 그 토픽은 '경쟁사 강점 영역',
        파스쿠찌 언급이 적은 토픽은 '담론 공백 → 선점 기회'로 해석합니다.
        </div>""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════
# TAB — 소비맥락 세분화 (3축 결합) [고도화 B]
# ════════════════════════════════════════════════════

def _rt_burst():
    st.markdown("<div class='section-title'>⚡ Burst Detection — 급증 키워드 탐지</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='warning-box'>
    <b>Burst Detection</b>: 최근 2개월(2026.04-05) vs 이전 3개월(2025.10-12) 키워드 증가율을 비교합니다.
    급증 키워드는 브랜드 위기 또는 새로운 소비 담론의 선행 신호입니다.
    </div>""", unsafe_allow_html=True)

    c1, c2 = st.columns([1.3, 1])
    with c1:
        top15 = burst_df.head(15)
        colors_burst = ["#E74C3C" if any(k in kw for k in ["환불","일베","이마트","정용진","극우","신세계","멸공"])
                        else "#F39C12" if any(k in kw for k in ["불매","텀블러","유니클로"])
                        else "#2980B9"
                        for kw in top15["키워드"]]
        fig = px.bar(
            top15.sort_values("증가율(배)"),
            x="증가율(배)", y="키워드", orientation="h",
            title="급증 키워드 TOP 15 (증가율 배수)",
            text=top15.sort_values("증가율(배)")["증가율(배)"].apply(lambda x: f"{x:.1f}x"),
            color="키워드",
            color_discrete_sequence=colors_burst,
        )
        fig.update_traces(textposition="outside")
        chart_style(fig, height=460, showlegend=False)
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        decaf_months = [m for m in all_months if m >= "2025-09"]
        decaf_docs = v[v["full_text"].str.contains("디카페인|디카페", na=False)]
        decaf_cnt  = {m: int(decaf_docs[decaf_docs["year_month"] == m].shape[0]) for m in decaf_months}
        fig2 = go.Figure()
        fig2.add_trace(go.Scatter(
            x=list(decaf_cnt.keys()),
            y=list(decaf_cnt.values()),
            mode="lines+markers+text",
            line=dict(color="#27AE60", width=2.5),
            marker=dict(size=8, color="#27AE60"),
            fill="tozeroy", fillcolor="rgba(39,174,96,0.12)",
            text=[str(v) for v in decaf_cnt.values()],
            textposition="top center",
            textfont=dict(size=11, color="#1A4731"),
            name="디카페인 언급",
        ))
        fig2.update_layout(title="디카페인 언급 월별 추이",
                           xaxis_title="월", yaxis_title="문서 수")
        chart_style(fig2, height=240, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)

        st.markdown("""
        <div class='success-box'><b>디카페인 — 구조적 성장 신호</b><br>
        2025.09 6건 → 2026.05 10건. NSS +56.5로 전 속성 최고 긍정률.
        파스쿠찌의 선점 가능한 1순위 카테고리.</div>""", unsafe_allow_html=True)
        st.markdown("""
        <div class='danger-box'><b>5월 급증 = 구조적 위기</b><br>
        환불·일베·이마트·정용진·극우 등 정치·사회 이슈 키워드 신규 폭발적 등장.
        스타벅스 이탈은 단기 회복 어려운 구조적 전환.</div>""", unsafe_allow_html=True)

    st.markdown("<div class='section-title'>급증 키워드 상세 데이터</div>", unsafe_allow_html=True)
    st.dataframe(burst_df, use_container_width=True, hide_index=True)

# ════════════════════════════════════════════════════
# TAB 8 — PMI 분석
# ════════════════════════════════════════════════════
def _rt_pmi():
    st.markdown("<div class='section-title'>🔗 PMI — 브랜드 공출현 강도 분석</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    <b>PMI (Pointwise Mutual Information)</b>: 두 브랜드가 같은 문서에서 함께 언급될 때
    우연 대비 얼마나 더 자주 함께 등장하는지를 측정합니다.
    PMI가 높을수록 소비자가 두 브랜드를 비교·연관해서 인식하는 경향이 강합니다.
    </div>""", unsafe_allow_html=True)

    pmi_top = pmi_df[pmi_df["PMI"] > -5].head(15)
    fig = px.bar(
        pmi_top, x="PMI", y=pmi_top["브랜드1"] + " ↔ " + pmi_top["브랜드2"],
        orientation="h", color="PMI",
        color_continuous_scale=["#BEE3F8","#2D6BC4"],
        title="브랜드 간 공출현 PMI 상위 15쌍",
        text=pmi_top["PMI"].apply(lambda x: f"{x:.2f}"),
        labels={"y":"브랜드 쌍"},
    )
    fig.update_traces(textposition="outside")
    chart_style(fig, height=480, showlegend=False)
    st.plotly_chart(fig, use_container_width=True)

    st.markdown("<div class='section-title'>PMI 매트릭스 히트맵</div>", unsafe_allow_html=True)
    pmi_matrix = pmi_df.pivot_table(index="브랜드1", columns="브랜드2",
                                    values="PMI", fill_value=0)
    all_b = sorted(set(pmi_df["브랜드1"].tolist() + pmi_df["브랜드2"].tolist()))
    pmi_full = pd.DataFrame(0.0, index=all_b, columns=all_b)
    for _, row in pmi_df.iterrows():
        pmi_full.loc[row["브랜드1"], row["브랜드2"]] = row["PMI"]
        pmi_full.loc[row["브랜드2"], row["브랜드1"]] = row["PMI"]
    # 발산형 컬러스케일: 음(파랑) - 0(흰색) - 양(빨강)으로 명확히 구분
    pmi_vals = pmi_full.values
    pmi_abs_max = max(abs(pmi_vals.min()), abs(pmi_vals.max())) if pmi_vals.size else 1
    fig2 = px.imshow(
        pmi_full,
        color_continuous_scale=["#2166AC","#67A9CF","#F7F7F7","#EF8A62","#B2182B"],
        color_continuous_midpoint=0,
        zmin=-pmi_abs_max, zmax=pmi_abs_max,
        title="브랜드 공출현 PMI 매트릭스 (빨강=강한 연관, 파랑=약한 연관)",
        text_auto=".2f", aspect="auto",
    )
    fig2.update_traces(textfont=dict(color="#1B2A4A", size=11))
    chart_style(fig2, height=440, showlegend=True)
    st.plotly_chart(fig2, use_container_width=True)

    col1, col2 = st.columns(2)
    with col1:
        st.markdown("""
        <div class='insight-box'><b>프리미엄 소규모 클러스터</b><br>
        테라로사·폴바셋·블루보틀·이디야·할리스가 높은 PMI로 묶임.
        소비자가 이 브랜드들을 함께 '대안 프리미엄 선택지'로 인식.
        파스쿠찌는 이 클러스터에 진입 가능한 이미지 자산 보유.
        </div>""", unsafe_allow_html=True)
    with col2:
        st.markdown("""
        <div class='warning-box'><b>스타벅스의 낮은 PMI</b><br>
        스타벅스는 대량 개별 언급이 많아 PMI가 낮음.
        다른 브랜드와 '함께 비교되기보다' 독자적으로 회자되는 구조.
        불매 후 대안으로는 PMI가 높은 클러스터 브랜드들이 부상.
        </div>""", unsafe_allow_html=True)

    st.markdown("<div class='section-title'>PMI 상세 데이터</div>", unsafe_allow_html=True)
    st.dataframe(pmi_df[pmi_df["공출현수"] > 0].sort_values("PMI", ascending=False),
                 use_container_width=True, hide_index=True)

# ════════════════════════════════════════════════════
# TAB 9 — 포지셔닝맵
# ════════════════════════════════════════════════════
def _rt_pos():
    st.markdown("<div class='section-title'>📍 브랜드 포지셔닝맵</div>", unsafe_allow_html=True)

    def norm01(s):
        mn, mx = s.min(), s.max()
        return (s - mn) / (mx - mn + 1e-9)

    pos_df = brand_df.copy()
    pos_df["norm_total"]  = norm01(pos_df["언급량"])
    pos_df["norm_nss"]    = norm01(pos_df["NSS"])
    pos_df["norm_pos"]    = norm01(pos_df["긍정률"])
    pos_df["neg_risk"]    = norm01(100 - pos_df["긍정률"])

    st.caption("두 장의 사분면 지도로 브랜드 위치를 직관적으로 읽습니다. 점선은 브랜드 평균선이며, 모서리의 회색 라벨이 그 칸(사분면)의 의미입니다.")

    def _quad_labels(_fig, tl, tr, bl, br):
        """사분면 네 모서리에 의미 라벨(paper 좌표)."""
        for _x, _y, _t, _xa, _ya in [
            (0.015, 0.985, tl, "left", "top"), (0.985, 0.985, tr, "right", "top"),
            (0.015, 0.015, bl, "left", "bottom"), (0.985, 0.015, br, "right", "bottom")]:
            _fig.add_annotation(x=_x, y=_y, xref="paper", yref="paper", text=_t,
                                showarrow=False, align=_xa, xanchor=_xa, yanchor=_ya,
                                font=dict(size=10.5, color="#8A98AC"),
                                bgcolor="rgba(247,249,252,0.65)", borderpad=2)

    c1, c2 = st.columns(2)
    with c1:
        fig = px.scatter(
            pos_df, x="norm_total", y="NSS",
            size="언급량", color="브랜드",
            color_discrete_map=BRAND_COLORS, text="브랜드",
            size_max=70, title="① 화제성 × 호감도",
        )
        fig.update_traces(textposition="top center",
                          textfont=dict(size=10, color="#1B2A4A", family="Arial Black"))
        fig.add_hline(y=pos_df["NSS"].mean(), line_dash="dot", line_color="#A0AEC0")
        fig.add_vline(x=pos_df["norm_total"].mean(), line_dash="dot", line_color="#A0AEC0")
        fig.update_layout(
            xaxis_title="◀ 적게 언급　　　화제성 (언급량) — 얼마나 많이　　　많이 언급 ▶",
            yaxis_title="◀ 호감 낮음　　호감도 (NSS) — 얼마나 좋게　　호감 높음 ▶")
        _quad_labels(fig,
            "💎 알짜 호감<br>호감 높음·인지도 낮음",
            "🏆 시장 리더<br>많이 + 좋게 언급",
            "🌱 잠재·약세<br>인지도·호감 모두 낮음",
            "⚠️ 논란형<br>많이 언급되나 호감 약함")
        chart_style(fig, height=450, showlegend=False)
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        fig2 = px.scatter(
            pos_df, x="긍정률", y="부정률",
            size="언급량", color="브랜드",
            color_discrete_map=BRAND_COLORS, text="브랜드",
            size_max=70, title="② 긍정 비율 × 부정 비율(리스크)",
        )
        fig2.update_traces(textposition="top center",
                           textfont=dict(size=10, color="#1B2A4A", family="Arial Black"))
        fig2.add_hline(y=pos_df["부정률"].mean(), line_dash="dot", line_color="#A0AEC0")
        fig2.add_vline(x=pos_df["긍정률"].mean(), line_dash="dot", line_color="#A0AEC0")
        fig2.update_layout(
            xaxis_title="◀ 긍정 적음　　　긍정 비율 (%)　　　긍정 많음 ▶",
            yaxis_title="◀ 리스크 낮음　　부정 비율 (%)　　리스크 높음 ▶")
        _quad_labels(fig2,
            "🚨 평판 위험<br>긍정 적고 부정 많음",
            "🔥 양극화<br>호불호 뚜렷",
            "😐 무관심·중립<br>긍·부정 모두 적음",
            "😀 안정 호감<br>긍정 많고 부정 적음")
        chart_style(fig2, height=450, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)

    st.markdown("""
    <div class='insight-box' style='font-size:0.82rem;'>
    <b>사분면 읽는 법</b> &nbsp;① <b>화제성×호감도</b>: 오른쪽 위일수록 '많이 + 좋게' 언급되는 강자,
    오른쪽 아래는 화제는 많지만 평판이 약한 '논란형'입니다. &nbsp;
    ② <b>긍정×부정</b>: 오른쪽 아래(긍정↑·부정↓)가 가장 안전한 '안정 호감', 왼쪽 위(긍정↓·부정↑)는 '평판 위험'입니다.
    점선은 브랜드 평균선입니다.
    </div>""", unsafe_allow_html=True)

    st.markdown("<div class='section-title'>경쟁 브랜드 경쟁력 레이더차트</div>", unsafe_allow_html=True)
    st.caption("경쟁 브랜드 비교 — 파스쿠찌는 굵은 라인+채움으로 강조 표기 (파스쿠찌 값은 추정 기준)")
    radar_axes = ["SOV 규모","NSS 평판","맛·품질","공간 전문성","가격 포지션","디카페인 기회","브랜드 이미지"]
    radar_data = {
        "스타벅스":  [100, 20, 50, 60, 40, 55, 45],
        "이디야":    [15, 80, 70, 55, 60, 75, 65],
        "테라로사":  [8,  90, 88, 70, 40, 70, 85],
        "메가커피":  [12, 65, 55, 30, 90, 40, 40],
        "투썸플레이스": [25, 55, 72, 65, 50, 50, 70],
    }
    pascucci_vals = [18, 60, 78, 72, 45, 55, 75]  # [v2.3.4] 파스쿠찌 추정치
    fig3 = go.Figure()
    colors_radar = ["#1B6CA8","#E8511D","#2CA02C","#9467BD","#D62728"]
    for (name, vals), color in zip(radar_data.items(), colors_radar):
        v_plot = vals + [vals[0]]
        fig3.add_trace(go.Scatterpolar(
            r=v_plot, theta=radar_axes + [radar_axes[0]],
            fill=None, name=name,
            line=dict(color=color, width=1.8),
            marker=dict(size=4, color=color), opacity=0.7,
        ))
    # [v2.3.4] 파스쿠찌 — 굵은 라인 + 채움으로 최상단 강조
    pv = pascucci_vals + [pascucci_vals[0]]
    fig3.add_trace(go.Scatterpolar(
        r=pv, theta=radar_axes + [radar_axes[0]],
        fill="toself", fillcolor="rgba(192,57,43,0.13)", name="★ 파스쿠찌",
        line=dict(color="#C0392B", width=5),
        marker=dict(size=9, color="#C0392B"),
    ))
    fig3.update_layout(
        polar=dict(
            bgcolor="#FFFFFF",
            radialaxis=dict(visible=True, range=[0, 100],
                            gridcolor="#E2E8F0", tickfont=dict(color="#718096")),
            angularaxis=dict(gridcolor="#E2E8F0",
                             tickfont=dict(color="#2D3748", size=12)),
        ),
        legend=dict(bgcolor="#FFFFFF", bordercolor="#E2E8F0", borderwidth=1,
                    font=dict(color="#2D3748", size=12)),
        title="경쟁 브랜드 경쟁력 레이더 (선형 비교)",
        paper_bgcolor="#F7F9FC", height=500,
    )
    st.plotly_chart(fig3, use_container_width=True)

# ════════════════════════════════════════════════════
# TAB 10 — 키워드 버블
# ════════════════════════════════════════════════════
def _rt_kw():
    st.markdown("<div class='section-title'>☁️ 브랜드별 키워드 버블차트</div>", unsafe_allow_html=True)
    sel_b10 = st.selectbox("브랜드", TARGET_BRANDS, key="kw_brand")
    n_kw10  = st.slider("표시 키워드 수", 15, 40, 25, key="kw_n")

    b10_docs = v[v["brands"].apply(lambda x: sel_b10 in x)]
    tokens10 = [t for ts in b10_docs["tokens"] for t in ts]
    topN     = Counter(tokens10).most_common(n_kw10)

    if topN:
        words, counts = zip(*topN)
        # [v2.5.1] 트리맵 — 면적·색 농도=빈도, 인접 배치로 키워드 군집을 직관 파악
        base_color = BRAND_COLORS.get(sel_b10, "#2D6BC4")
        kw_df = pd.DataFrame({"키워드": list(words), "빈도": [int(c) for c in counts]})
        fig = px.treemap(
            kw_df, path=[px.Constant(f"{sel_b10} 키워드"), "키워드"], values="빈도",
            color="빈도", color_continuous_scale=["#EAF1FB", "#7FA8DC", base_color],
            title=f"{sel_b10} — 키워드 트리맵 (면적·색 = 언급 빈도)",
        )
        fig.update_traces(
            textinfo="label+value",
            textfont=dict(size=15),
            marker=dict(line=dict(width=1.5, color="white")),
            tiling=dict(pad=2),
        )
        fig.update_layout(height=520, margin=dict(t=50, b=10, l=10, r=10),
                          paper_bgcolor="#F7F9FC")
        st.plotly_chart(fig, use_container_width=True)
        st.caption("사각형 크기·색 농도 = 언급 빈도. 인접한 키워드일수록 함께 묶여 보이므로, "
                   "양뿐 아니라 키워드 간 군집·연관 맥락에서 기획 영감을 얻을 수 있습니다.")

    st.markdown("<div class='section-title'>감성별 키워드 비교</div>", unsafe_allow_html=True)
    st.caption(f"선택 브랜드({sel_b10})의 긍정/부정 리뷰 TOP 키워드")
    c1, c2 = st.columns(2)
    for col, sent in zip([c1, c2], ["긍정","부정"]):
        with col:
            s_toks = [t for ts in b10_docs[b10_docs["sentiment"]==sent]["tokens"] for t in ts]
            s_top  = Counter(s_toks).most_common(12)
            if s_top:
                sw, sc = zip(*s_top)
                fig_s = px.bar(
                    x=list(sc), y=list(sw), orientation="h",
                    color=list(sc),
                    color_continuous_scale=["#EBF4FF", SENT_COLORS[sent]],
                    title=f"{sent} 리뷰 TOP 키워드",
                    text=[str(c) for c in sc],
                    labels={"x":"빈도","y":"키워드"},
                )
                fig_s.update_traces(textposition="outside")
                fig_s.update_layout(yaxis=dict(autorange="reversed"))
                chart_style(fig_s, height=380, showlegend=False)
                st.plotly_chart(fig_s, use_container_width=True)

    # [항목 3] 전체 데이터 기준 + 키워드 유형별(속성 카테고리) 감성 키워드 비교
    st.markdown("<div class='section-title'>전체 데이터 — 키워드 유형별 감성 비교</div>", unsafe_allow_html=True)
    st.caption("브랜드 무관 전체 VoC에서 속성 유형(맛·가격·공간 등)별로 긍정/부정 키워드를 구분해 비교")
    kw_type = st.selectbox("키워드 유형 선택", ["전체"] + list(ATTR_DICT.keys()), key="kwtype_sel")
    if kw_type == "전체":
        type_docs = v
    else:
        pat_t = "|".join(ATTR_DICT[kw_type])
        type_docs = v[v["full_text"].str.contains(pat_t, na=False)]
    cc1, cc2 = st.columns(2)
    for col, sent in zip([cc1, cc2], ["긍정","부정"]):
        with col:
            st_toks = [t for ts in type_docs[type_docs["sentiment"]==sent]["tokens"] for t in ts]
            # 유형 선택 시 해당 유형 키워드 위주로 필터
            if kw_type != "전체":
                attr_kws = ATTR_DICT[kw_type]
                st_toks = [t for t in st_toks if any(a in t for a in attr_kws)] or st_toks
            st_top = Counter(st_toks).most_common(12)
            if st_top:
                tw, tc = zip(*st_top)
                fig_t = px.bar(
                    x=list(tc), y=list(tw), orientation="h",
                    color=list(tc),
                    color_continuous_scale=["#EBF4FF", SENT_COLORS[sent]],
                    title=f"[{kw_type}] {sent} 키워드 (전체 {len(type_docs)}건)",
                    text=[str(c) for c in tc],
                    labels={"x":"빈도","y":"키워드"},
                )
                fig_t.update_traces(textposition="outside")
                fig_t.update_layout(yaxis=dict(autorange="reversed"))
                chart_style(fig_t, height=380, showlegend=False)
                st.plotly_chart(fig_t, use_container_width=True)

# ════════════════════════════════════════════════════
# TAB 11 — 소비 맥락
# ════════════════════════════════════════════════════
def _rt_occ():
    st.markdown("<div class='section-title'>🎯 소비 맥락 (Occasion) 분석</div>", unsafe_allow_html=True)

    occ_cnt = v["occasion"].value_counts().reset_index()
    occ_cnt.columns = ["상황","문서수"]
    fig = px.bar(
        occ_cnt, x="문서수", y="상황", orientation="h",
        color="문서수", color_continuous_scale=["#BEE3F8","#2D6BC4"],
        title="소비 상황별 언급량",
        text="문서수",
    )
    fig.update_traces(textposition="outside")
    fig.update_layout(yaxis=dict(autorange="reversed"))
    chart_style(fig, height=320, showlegend=False)
    st.plotly_chart(fig, use_container_width=True)

    st.markdown("<div class='section-title'>소비 상황 × 브랜드 히트맵</div>", unsafe_allow_html=True)
    occs = [o for o in OCCASION_DICT.keys()] + ["기타"]
    occ_brand_mat = pd.DataFrame(0, index=occs, columns=TARGET_BRANDS)
    for occ in occs:
        sub = v[v["occasion"] == occ]
        for b in TARGET_BRANDS:
            cnt = sub["brands"].apply(lambda x: b in x).sum()
            occ_brand_mat.loc[occ, b] = cnt
    fig2 = px.imshow(
        occ_brand_mat,
        color_continuous_scale="Blues",
        title="소비 상황 × 브랜드 언급 히트맵 (건수)",
        text_auto=True, aspect="auto",
    )
    fig2.update_traces(textfont=dict(color="#1B2A4A", size=11))
    chart_style(fig2, height=380, showlegend=False)
    st.plotly_chart(fig2, use_container_width=True)

    st.markdown("<div class='section-title'>소비 상황별 감성 분포</div>", unsafe_allow_html=True)
    occ_sent = v.groupby(["occasion","sentiment"]).size().reset_index(name="수")
    fig3 = px.bar(
        occ_sent, x="occasion", y="수", color="sentiment",
        color_discrete_map=SENT_COLORS, barmode="stack",
        title="소비 상황별 감성 누적",
        labels={"occasion":"소비 상황","수":"문서 수"},
    )
    chart_style(fig3, height=340)
    st.plotly_chart(fig3, use_container_width=True)

    st.markdown("""
    <div class='insight-box'><b>파스쿠찌 소비 상황 전략</b><br>
    공부·업무(카공)와 만남·미팅 상황에서 스타벅스 집중도가 가장 높음.
    이탈 후 대안으로 파스쿠찌가 포지셔닝되려면 '조용하고 넓은 공간' 강조와 함께
    카공족을 위한 콘센트·WiFi 특화 마케팅이 효과적.
    </div>""", unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("<div class='section-title'>🗂️ 소비맥락 세분화 (15유형) — 같은 화면 통합 [v2.3.4]</div>", unsafe_allow_html=True)
    st.markdown("<div class='section-title'>🗂️ 소비맥락 세분화 — 15유형 × 브랜드 결합</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    기존 5~6개 일반 맥락을 <b>15개 세부 상황</b>으로 분화하고, <b>브랜드·감성과 교차</b>합니다.<br>
    "누가, 언제, 왜, 어떤 브랜드를" 선택하는지 입체적으로 파악해 타깃 마케팅 인사이트를 도출합니다.
    </div>""", unsafe_allow_html=True)

    if len(occ_detail_df) == 0:
        st.info("소비맥락 분석을 위한 데이터가 부족합니다.")
    else:
        c1, c2 = st.columns([1.3, 1])
        with c1:
            fig = px.bar(
                occ_detail_df.sort_values("언급량"),
                x="언급량", y="소비맥락", orientation="h",
                color="NSS", color_continuous_scale=["#E74C3C","#FFFFFF","#27AE60"],
                color_continuous_midpoint=0,
                title="15개 세부 소비맥락 언급량 (색상=NSS)",
                text="언급량",
            )
            fig.update_traces(textposition="outside")
            chart_style(fig, height=480, showlegend=False)
            st.plotly_chart(fig, use_container_width=True)
        with c2:
            st.markdown("<div class='section-title'>맥락별 긍정률 TOP</div>", unsafe_allow_html=True)
            for _, row in occ_detail_df.nlargest(8, "긍정률").iterrows():
                st.markdown(f"""
                <div style='background:#F0FFF4;border-left:3px solid #27AE60;
                            border-radius:0 6px 6px 0;padding:7px 11px;margin-bottom:5px;'>
                    <b style='color:#276749;font-size:0.85rem;'>{row['소비맥락']}</b>
                    <span style='float:right;color:#27AE60;font-weight:700;'>긍정 {row['긍정률']}%</span>
                    <div style='color:#718096;font-size:0.72rem;'>{row['언급량']}건 · NSS {row['NSS']:+.1f}</div>
                </div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>소비맥락 × 브랜드 교차 히트맵</div>", unsafe_allow_html=True)
        st.caption("각 소비 상황에서 어떤 브랜드가 주로 선택되는지 — 빈 셀이 파스쿠찌 진입 기회")
        fig2 = px.imshow(
            occ_brand_mat,
            color_continuous_scale="Blues",
            title="소비맥락 × 브랜드 언급 매트릭스",
            text_auto=True, aspect="auto",
        )
        fig2.update_traces(textfont=dict(color="#1B2A4A", size=10))
        chart_style(fig2, height=520, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)

        st.markdown("""
        <div class='success-box'>
        <b>전략 활용</b>: "집중 카공" "디저트 타임" 등 파스쿠찌가 강점을 가질 수 있는 상황을
        선별해 상황 맞춤형 메뉴·공간·프로모션을 기획합니다. 경쟁사가 독점한 상황은 차별화 포인트를 찾습니다.
        </div>""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════
# TAB 12 — 리스크 탐지
# ════════════════════════════════════════════════════
def _rt_risk():
    st.markdown("<div class='section-title'>🚨 리스크 탐지 — 부정 드라이버 맵</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='danger-box'>
    <b>리스크 탐지</b>: 부정 리뷰에서 반복 등장하는 불만 패턴을 자동 감지합니다.
    빈도가 높고 여러 브랜드에 걸쳐 나타나는 키워드일수록 구조적 리스크 신호입니다.
    </div>""", unsafe_allow_html=True)

    neg_docs = v[v["sentiment"] == "부정"].copy()
    neg_tokens_all = [t for ts in neg_docs["tokens"] for t in ts]
    neg_top40 = Counter(neg_tokens_all).most_common(40)

    c1, c2 = st.columns([1.2, 1])
    with c1:
        if neg_top40:
            nwords, ncounts = zip(*neg_top40)
            nc_arr = np.array(ncounts, dtype=float)
            np.random.seed(7)
            angles = np.linspace(0, 4*np.pi, len(nwords))
            radii  = np.linspace(0.1, 1.0, len(nwords))
            nx = radii*np.cos(angles) + np.random.uniform(-0.1, 0.1, len(nwords))
            ny = radii*np.sin(angles) + np.random.uniform(-0.1, 0.1, len(nwords))
            sizes = (nc_arr / nc_arr.max()) * 100 + 15
            fig = go.Figure()
            for word, cnt, xi, yi, sz in zip(nwords, ncounts, nx, ny, sizes):
                op = 0.4 + (cnt / nc_arr.max()) * 0.6
                fig.add_trace(go.Scatter(
                    x=[xi], y=[yi], mode="markers+text",
                    marker=dict(size=sz, color="#E74C3C", opacity=op,
                                line=dict(color="white", width=1.5)),
                    text=[word], textposition="middle center",
                    textfont=dict(size=max(8, min(14, int(sz/5))),
                                  color="#FFFFFF" if sz > 40 else "#742A2A"),
                    hovertemplate=f"<b>{word}</b><br>빈도: {cnt}회<extra></extra>",
                    showlegend=False,
                ))
            fig.update_layout(
                title="부정 리뷰 키워드 버블맵 (크기=빈도)",
                xaxis=dict(visible=False, range=[-1.5, 1.5]),
                yaxis=dict(visible=False, range=[-1.5, 1.5]),
                plot_bgcolor="#FFF5F5", paper_bgcolor="#F7F9FC", height=480,
            )
            st.plotly_chart(fig, use_container_width=True)
    with c2:
        st.markdown("**브랜드별 주요 불만 키워드**")
        for b in TARGET_BRANDS:
            neg_b   = v[(v["brands"].apply(lambda x: b in x)) & (v["sentiment"] == "부정")]
            neg_tok = [t for ts in neg_b["tokens"] for t in ts]
            top5    = Counter(neg_tok).most_common(5)
            if not top5:
                continue
            color   = BRAND_COLORS.get(b, "#2D6BC4")
            kw_str  = "  ·  ".join([f"<b>{w}</b>({c})" for w, c in top5])
            st.markdown(f"""
            <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-left:4px solid {color};
                        border-radius:0 8px 8px 0;padding:8px 12px;margin-bottom:6px'>
                <span style='font-weight:700;color:{color};font-size:0.88rem'>{b}</span>
                <div style='color:#C53030;font-size:0.82rem;margin-top:4px'>{kw_str}</div>
            </div>""", unsafe_allow_html=True)

    st.markdown("<div class='section-title'>리스크 유형별 브랜드 히트맵</div>", unsafe_allow_html=True)
    risk_rows = []
    for risk_type, kws in RISK_PATTERNS.items():
        pat = "|".join(kws)
        for b in TARGET_BRANDS:
            nb  = v[(v["brands"].apply(lambda x: b in x)) & (v["sentiment"] == "부정")]
            cnt = nb["full_text"].str.contains(pat, na=False).sum()
            risk_rows.append({"리스크 유형": risk_type, "브랜드": b, "빈도": cnt})
    risk_df = pd.DataFrame(risk_rows)
    risk_pivot = risk_df.pivot_table(index="브랜드", columns="리스크 유형",
                                     values="빈도", fill_value=0)
    if not risk_pivot.empty:
        fig2 = px.imshow(
            risk_pivot,
            color_continuous_scale=["#FFFFFF","#FED7D7","#E53E3E"],
            title="브랜드 × 리스크 유형 히트맵 (부정 리뷰 빈도)",
            text_auto=True, aspect="auto",
        )
        fig2.update_traces(textfont=dict(color="#1B2A4A", size=12))
        chart_style(fig2, height=320, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)

    st.markdown("<div class='section-title'>부정 리뷰 원문 탐색</div>", unsafe_allow_html=True)
    st.caption("브랜드 필터는 '주요 브랜드(본문 내 최다 언급)' 기준 — 곁가지로 스친 브랜드는 제외되어 정합성을 높였습니다.")
    rcol1, rcol2 = st.columns([1, 1])
    with rcol1:
        rb  = st.selectbox("브랜드 필터", ["전체"] + sorted(TARGET_BRANDS), key="risk_b")
    with rcol2:
        rmode = st.radio("필터 기준", ["주요 브랜드만", "언급된 전체"], horizontal=True, key="risk_mode")
    rkw = st.text_input("키워드 검색", placeholder="예: 환불, 불매, 비싸", key="risk_kw")
    if rb == "전체":
        neg_show = neg_docs
    elif rmode == "주요 브랜드만":
        # [수정] primary_brand가 선택 브랜드인 문서만 — 무관 원문 혼입 방지
        neg_show = neg_docs[neg_docs["primary_brand"] == rb]
    else:
        neg_show = neg_docs[neg_docs["brands"].apply(lambda x: rb in x)]
    if rkw:
        neg_show = neg_show[neg_show["full_text"].str.contains(rkw, na=False)]
    show_cols = ["title","body","primary_brand","year_month","sentiment"]
    show_df = neg_show[show_cols].copy()
    show_df.columns = ["제목","본문","주요브랜드","작성월","감성"]
    show_df["본문"] = show_df["본문"].str[:200]
    st.dataframe(show_df.head(30), use_container_width=True, height=300,
                 column_config={"본문": st.column_config.TextColumn(width="large")})
    st.caption(f"{min(30, len(show_df))}건 표시 / 전체 {len(show_df)}건")

# ════════════════════════════════════════════════════
# TAB — 브랜드 로열티 (찐팬 지수) [고도화 A]
# ════════════════════════════════════════════════════
def _rt_loyalty():
    st.markdown("<div class='section-title'>💎 브랜드 로열티 (찐팬) 지수</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    단순 호불호(NSS)를 넘어, <b>"이 브랜드의 찐팬 비중"</b>을 측정합니다.<br>
    <b>독점충성</b>("여기 아니면 안 가", 가중치 3.0) · <b>강한애착</b>("최애·단골", 2.0) · <b>반복소비</b>("자주·매일", 1.0)
    신호어를 문서에서 탐지해 <b>로열티 지수(0~100)</b>와 <b>찐팬 비율(%)</b>로 산출합니다.
    </div>""", unsafe_allow_html=True)

    if len(loyalty_df) == 0:
        st.info("로열티 분석을 위한 데이터가 부족합니다.")
    else:
        c1, c2 = st.columns([1.3, 1])
        with c1:
            sorted_l = loyalty_df.sort_values("로열티지수")
            fig = px.bar(
                sorted_l, x="로열티지수", y="브랜드", orientation="h",
                color="로열티지수", color_continuous_scale=["#F5E6E0", "#E67E22", "#C0392B"],
                title="브랜드별 로열티 지수 (0~100)",
                text=sorted_l["로열티지수"].apply(lambda x: f"{x:.1f}"),
            )
            fig.update_traces(textposition="outside")
            fig.add_vline(x=loyalty_df["로열티지수"].mean(), line_dash="dot",
                          line_color="#718096",
                          annotation_text=f"평균 {loyalty_df['로열티지수'].mean():.1f}")
            chart_style(fig, height=420, showlegend=False)
            st.plotly_chart(fig, use_container_width=True)
        with c2:
            for _, row in loyalty_df.head(6).iterrows():
                color = BRAND_COLORS.get(row["브랜드"], "#2D6BC4")
                st.markdown(f"""
                <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-left:4px solid {color};
                            border-radius:0 8px 8px 0;padding:9px 12px;margin-bottom:6px;'>
                    <span style='font-weight:700;color:{color};font-size:0.9rem;'>{row['브랜드']}</span>
                    <span style='float:right;font-weight:800;color:{color};'>찐팬 {row['찐팬비율']}%</span>
                    <div style='color:#718096;font-size:0.78rem;margin-top:3px;'>
                        로열티 {row['로열티지수']} · 독점 {row['독점충성']} / 애착 {row['강한애착']} / 반복 {row['반복소비']}
                    </div>
                </div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>로열티 지수 × NSS 매트릭스</div>", unsafe_allow_html=True)
        st.caption("우상단 = 호감도와 충성도 모두 높은 '찐팬 보유 브랜드' / 우하단 = 호감은 있으나 충성 약함(재방문 유도 필요)")
        fig2 = go.Figure()
        for _, row in loyalty_df.iterrows():
            color = BRAND_COLORS.get(row["브랜드"], "#2D6BC4")
            fig2.add_trace(go.Scatter(
                x=[row["NSS"]], y=[row["로열티지수"]],
                mode="markers+text",
                marker=dict(size=max(row["언급량"]/8, 14), color=color, opacity=0.8,
                            line=dict(color="white", width=2)),
                text=[row["브랜드"]], textposition="top center",
                textfont=dict(size=11, color=color),
                name=row["브랜드"],
                hovertemplate=f"<b>{row['브랜드']}</b><br>NSS: {row['NSS']:+.1f}<br>로열티: {row['로열티지수']}<br>찐팬비율: {row['찐팬비율']}%<extra></extra>",
            ))
        fig2.add_hline(y=loyalty_df["로열티지수"].mean(), line_dash="dot", line_color="#A0AEC0")
        fig2.add_vline(x=loyalty_df["NSS"].mean(), line_dash="dot", line_color="#A0AEC0")
        fig2.update_layout(title="로열티 × NSS (버블=언급량)",
                           xaxis_title="NSS (호감도)", yaxis_title="로열티 지수 (충성도)",
                           showlegend=False)
        chart_style(fig2, height=460)
        st.plotly_chart(fig2, use_container_width=True)

        st.markdown("<div class='section-title'>충성 신호 구성 (Stacked)</div>", unsafe_allow_html=True)
        fig3 = px.bar(
            loyalty_df.sort_values("찐팬비율"),
            x="브랜드", y=["독점충성", "강한애착", "반복소비"],
            title="브랜드별 충성 신호 문서 수 구성 (절대량)",
            color_discrete_map={"독점충성": "#C0392B", "강한애착": "#E67E22", "반복소비": "#F5C77E"},
            labels={"value": "문서 수", "variable": "충성 등급"},
        )
        chart_style(fig3, height=350)
        st.plotly_chart(fig3, use_container_width=True)

        # [기능 2] 절대량 ↔ 지수 차이 해설
        st.markdown("""
        <div class='warning-box'>
        <b>왜 스타벅스는 충성 신호 '수'가 가장 많은데 로열티 지수는 낮을까?</b><br>
        충성 신호 구성(Stacked) 차트는 <b>절대 문서 수</b>를 보여줍니다. 스타벅스는
        전체 언급량 자체가 압도적으로 많아(전체의 절반 이상) 충성 신호 문서의 '수'도 자연히 많습니다.<br><br>
        반면 <b>로열티 지수와 찐팬 비율은 '비중'(충성 문서 수 ÷ 전체 언급 수)</b>으로 계산합니다.
        스타벅스는 분모(전체 언급)가 워낙 커서, 충성 신호가 많아도 <b>비율로는 희석</b>됩니다.
        또한 스타벅스는 불매·이슈 등 부정·중립 담론도 많이 섞여 있어 충성 비중이 더 낮아집니다.<br><br>
        즉, <b>절대량(영향력·화제성)은 스타벅스가 1위지만, 충성도 밀도(찐팬 비중)는
        메가·이디야 등 가성비·니치 브랜드가 더 높은</b> 구조입니다. 두 지표를 함께 봐야
        "많이 회자되는 브랜드"와 "찐팬이 두터운 브랜드"를 구분할 수 있습니다.
        </div>""", unsafe_allow_html=True)

        st.dataframe(
            loyalty_df[["브랜드","언급량","로열티지수","찐팬비율","독점충성","강한애착","반복소비","NSS"]],
            use_container_width=True, hide_index=True,
        )


# ════════════════════════════════════════════════════
# TAB — 토픽 클러스터링 (토픽×브랜드×고객 중첩) [고도화 C]
# ════════════════════════════════════════════════════
# ════════════════════════════════════════════════════
# TAB — 트렌드 조기경보 (상대평가) [고도화 D]
# ════════════════════════════════════════════════════
def _render_fnb_trend():
    """[v2.6.5] 식음료 인접 트렌드 조기경보 — 커피 외 디저트·베이커리·식재료 VoC. 브랜드 지표 미반영."""
    st.markdown("---")
    st.markdown("<div class='section-title'>🍰 식음료 인접 트렌드 조기경보 "
                "<span style='font-size:0.62em;color:#94A3B8;font-weight:600'>· 커피 외 디저트·베이커리·식재료 VoC · 브랜드 지표 미반영</span></div>",
                unsafe_allow_html=True)
    if v_adj is None or len(v_adj) == 0:
        st.info("인접 식음료 VoC가 충분하지 않습니다. (이 데이터셋에 커피 외 식음료 글이 적습니다.)")
        return
    kw_df, cat_df = fnb_trends
    c1, c2, c3 = st.columns(3)
    c1.metric("인접 식음료 VoC", f"{len(v_adj):,}건")
    c2.metric("긍정 비중", f"{(v_adj['sentiment'] == '긍정').mean() * 100:.0f}%")
    c3.metric("센싱 카테고리", f"{len(cat_df)}개")
    st.caption("커피 브랜드 언급이 없어 NSS·SOV 등 브랜드 지표에는 반영되지 않습니다. 인접 시장의 흐름을 읽는 상품기획 영감용 레이어입니다.")

    if not kw_df.empty:
        st.markdown("<div class='section-title'>📈 급상승 식음료 키워드 (최근월 vs 기준 평균)</div>", unsafe_allow_html=True)
        rad = kw_df[kw_df["총언급"] >= 8].head(12)
        if rad.empty:
            rad = kw_df.head(10)
        rad = rad.sort_values("모멘텀")
        fig = px.bar(rad, x="모멘텀", y="키워드", color="카테고리", orientation="h",
                     text=rad["총언급"].apply(lambda x: f"{x}건"),
                     title="급상승 트렌드 레이더 (모멘텀 = 최근월 ÷ 기준월 평균)")
        fig.update_traces(textposition="outside")
        chart_style(fig, height=420)
        st.plotly_chart(fig, use_container_width=True)

    if not cat_df.empty:
        cc1, cc2 = st.columns([1.2, 1])
        with cc1:
            st.markdown("<div class='section-title'>🍽️ 카테고리별 볼륨·감성</div>", unsafe_allow_html=True)
            figc = px.bar(cat_df, x="카테고리", y="건수", color="NSS",
                          color_continuous_scale=["#E74C3C", "#F4D03F", "#27AE60"],
                          text="건수", title="인접 식음료 카테고리별 VoC")
            chart_style(figc, height=340, showlegend=False)
            st.plotly_chart(figc, use_container_width=True)
        with cc2:
            st.markdown("<div class='section-title'>대표 소비자 원문</div>", unsafe_allow_html=True)
            shown = 0
            for cat in cat_df["카테고리"].tolist():
                sub = v_adj[v_adj["food_cats"].apply(lambda cs: cat in cs) & v_adj["sentiment"].isin(["긍정", "부정"])]
                sub = sub[sub["body"].fillna("").astype(str).str.len().between(25, 260)]
                if len(sub) == 0:
                    continue
                r = sub.iloc[0]
                sc = "#27AE60" if r["sentiment"] == "긍정" else "#E74C3C"
                body_txt = str(r["body"])[:110].replace("<", "&lt;").replace(">", "&gt;")
                st.markdown(f"<div style='background:#FFF;border:1px solid #E2E8F0;border-left:4px solid {sc};"
                            f"border-radius:0 8px 8px 0;padding:7px 11px;margin-bottom:6px;'>"
                            f"<b style='color:{sc};font-size:0.78rem'>[{cat}·{r['sentiment']}]</b>"
                            f"<div style='color:#475569;font-size:0.76rem;margin-top:3px'>{body_txt}</div></div>",
                            unsafe_allow_html=True)
                shown += 1
                if shown >= 4:
                    break
    if not kw_df.empty:
        st.markdown("<div class='section-title'>식음료 키워드 모멘텀 상세</div>", unsafe_allow_html=True)
        st.dataframe(kw_df.head(20), use_container_width=True, hide_index=True)


def _rt_trend():
    st.markdown("<div class='section-title'>📡 상품 트렌드 조기경보 — 상대평가 기반</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    단순 증감이 아니라 <b>늘 소비되는 기준 제품(아메리카노·라떼) 대비 상대 배수</b>로 유행을 예측합니다.<br>
    <b>상대배수</b>(기준선 대비) + <b>성장 기울기</b> + <b>브랜드 확산도</b>를 결합한 <b>ES점수</b>로
    "뜰 것 같은 제품·식재료"를 사전 포착합니다. (예: 두쫀쿠 열풍 전 카다이프 검색량 급증 패턴)
    </div>""", unsafe_allow_html=True)

    if trend_ew is None:
        st.info("트렌드 분석을 위한 데이터가 부족합니다.")
    else:
        ew_df_full = trend_ew["df"]
        st.markdown(f"""
        <div class='warning-box'>
        <b>기준선(아메리카노·라떼 등) 월평균 언급량 = {trend_ew['baseline_avg']}건</b><br>
        각 신흥 제품의 최근 언급량을 이 기준선과 비교해 상대 배수를 산출합니다.
        </div>""", unsafe_allow_html=True)

        # [항목 10] 카테고리 구분 필터
        cat_options = ["전체"] + list(TREND_CATEGORY.keys())
        sel_cat = st.radio("카테고리 구분", cat_options, horizontal=True, key="trend_cat")
        if sel_cat == "전체":
            ew_df = ew_df_full
        else:
            ew_df = ew_df_full[ew_df_full["카테고리"] == sel_cat]

        # 카테고리별 요약 카드
        ccols = st.columns(len(TREND_CATEGORY))
        cat_emoji = {"커피·음료": "☕", "푸드류": "🥐", "식재료": "🧂"}
        for col, (cat, items) in zip(ccols, TREND_CATEGORY.items()):
            cat_df = ew_df_full[ew_df_full["카테고리"] == cat]
            top_item = cat_df.iloc[0]["제품·식재료"] if len(cat_df) > 0 else "—"
            top_es = cat_df.iloc[0]["ES점수"] if len(cat_df) > 0 else 0
            with col:
                st.markdown(f"""
                <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-top:3px solid #E67E22;
                            border-radius:8px;padding:10px;text-align:center;'>
                    <div style='font-size:1.3rem;'>{cat_emoji.get(cat,"📦")}</div>
                    <div style='font-weight:700;color:#1B2A4A;font-size:0.85rem;'>{cat}</div>
                    <div style='color:#718096;font-size:0.72rem;'>{len(cat_df)}개 항목</div>
                    <div style='color:#E67E22;font-weight:700;font-size:0.78rem;margin-top:3px;'>
                        TOP: {top_item} (ES {top_es:.0f})
                    </div>
                </div>""", unsafe_allow_html=True)

        if len(ew_df) == 0:
            st.info(f"{sel_cat} 카테고리에 해당하는 제품이 없습니다.")
            st.stop()

        c1, c2 = st.columns([1.3, 1])
        with c1:
            sorted_ew = ew_df.sort_values("ES점수")
            alert_color_map = {
                "🚨 강한 신호": "#E74C3C", "📈 성장 조짐": "#E67E22",
                "➡️ 안정": "#2980B9", "📉 둔화": "#7F8C8D",
            }
            fig = px.bar(
                sorted_ew, x="ES점수", y="제품·식재료", orientation="h",
                color="경보", color_discrete_map=alert_color_map,
                title="제품·식재료 Emerging Signal 점수",
                text=sorted_ew["ES점수"].apply(lambda x: f"{x:.0f}"),
            )
            fig.update_traces(textposition="outside")
            chart_style(fig, height=480)
            st.plotly_chart(fig, use_container_width=True)
        with c2:
            st.markdown("<div class='section-title'>조기경보 신호</div>", unsafe_allow_html=True)
            for _, row in ew_df.head(8).iterrows():
                ac = {"🚨 강한 신호":"#E74C3C","📈 성장 조짐":"#E67E22",
                      "➡️ 안정":"#2980B9","📉 둔화":"#7F8C8D"}.get(row["경보"],"#2980B9")
                st.markdown(f"""
                <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-left:4px solid {ac};
                            border-radius:0 8px 8px 0;padding:8px 12px;margin-bottom:5px;'>
                    <b style='color:{ac};font-size:0.85rem;'>{row['경보']} {row['제품·식재료']}</b>
                    <div style='color:#718096;font-size:0.74rem;margin-top:3px;'>
                        상대배수 {row['상대배수']}배 · 기울기 {row['성장기울기']:+.1f} · {row['브랜드확산']}개 브랜드 · NSS {row['NSS']:+.1f}
                    </div>
                </div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>트렌드 스코어 포지셔닝 (상대배수 × 성장기울기)</div>", unsafe_allow_html=True)
        st.caption("우상단 = 이미 많이 언급되고 + 계속 성장 중 / 좌상단 = 아직 작지만 급성장 (조기 진입 기회)")
        fig2 = go.Figure()
        for _, row in ew_df.iterrows():
            ac = {"🚨 강한 신호":"#E74C3C","📈 성장 조짐":"#E67E22",
                  "➡️ 안정":"#2980B9","📉 둔화":"#7F8C8D"}.get(row["경보"],"#2980B9")
            fig2.add_trace(go.Scatter(
                x=[row["상대배수"]], y=[row["성장기울기"]],
                mode="markers+text",
                marker=dict(size=max(row["총언급"]/2, 12), color=ac, opacity=0.75,
                            line=dict(color="white", width=1.5)),
                text=[row["제품·식재료"]], textposition="top center",
                textfont=dict(size=10, color=ac),
                name=row["제품·식재료"],
                hovertemplate=f"<b>{row['제품·식재료']}</b><br>상대배수: {row['상대배수']}<br>기울기: {row['성장기울기']}<br>ES: {row['ES점수']}<extra></extra>",
            ))
        fig2.add_hline(y=0, line_dash="dash", line_color="#A0AEC0")
        fig2.add_vline(x=0.5, line_dash="dot", line_color="#A0AEC0")
        fig2.update_layout(title="상대배수 × 성장기울기 (버블=총언급량)",
                           xaxis_title="기준선 대비 상대배수", yaxis_title="성장 기울기",
                           showlegend=False)
        chart_style(fig2, height=460)
        st.plotly_chart(fig2, use_container_width=True)

        st.markdown("<div class='section-title'>상위 신호 제품 월별 추이</div>", unsafe_allow_html=True)
        top_products = ew_df.head(5)["제품·식재료"].tolist()
        ts_rows = []
        for _, row in ew_df.head(5).iterrows():
            for mi, m in enumerate(trend_ew["months"]):
                ts_rows.append({"제품": row["제품·식재료"], "월": m,
                                "언급량": row["_monthly"][mi] if mi < len(row["_monthly"]) else 0})
        ts_df = pd.DataFrame(ts_rows)
        fig3 = px.line(ts_df, x="월", y="언급량", color="제품", markers=True,
                       title="상위 5개 신호 제품 월별 언급 추이")
        chart_style(fig3, height=360)
        st.plotly_chart(fig3, use_container_width=True)

        st.dataframe(
            ew_df[["제품·식재료","총언급","상대배수","성장기울기","브랜드확산","NSS","ES점수","경보"]],
            use_container_width=True, hide_index=True,
        )

    _render_fnb_trend()


# ════════════════════════════════════════════════════
# TAB — 브랜드별 제품 [항목 4]
# ════════════════════════════════════════════════════
def _rt_product():
    st.markdown("<div class='section-title'>🛍️ 브랜드별 제품 언급·감성 분석</div>", unsafe_allow_html=True)
    st.markdown("""
    <div class='insight-box'>
    브랜드별로 가장 많이 언급되는 <b>제품(커피·음료·디저트·푸드 망라)</b>과 그 제품의
    <b>언급량·감성(NSS)·월별 추이</b>를 확인합니다. 시그니처 제품 발굴과 상품 기획에 활용합니다.
    </div>""", unsafe_allow_html=True)

    sel_bp = st.selectbox("브랜드 선택", TARGET_BRANDS, key="product_brand")
    prod_df = compute_brand_products(v, sel_bp)

    if len(prod_df) == 0:
        st.info(f"{sel_bp}의 제품 언급 데이터가 부족합니다.")
    else:
        c1, c2 = st.columns([1.4, 1])
        with c1:
            top_prod = prod_df.head(15).sort_values("언급량")
            fig = px.bar(
                top_prod, x="언급량", y="제품", orientation="h",
                color="NSS", color_continuous_scale=["#E74C3C","#FFFFFF","#27AE60"],
                color_continuous_midpoint=0,
                title=f"{sel_bp} — 제품별 언급량 (색상=NSS)",
                text="언급량",
                hover_data=["카테고리","긍정","부정"],
            )
            fig.update_traces(textposition="outside")
            chart_style(fig, height=480, showlegend=True)
            st.plotly_chart(fig, use_container_width=True)
        with c2:
            st.markdown("<div class='section-title'>카테고리별 비중</div>", unsafe_allow_html=True)
            cat_sum = prod_df.groupby("카테고리")["언급량"].sum().reset_index()
            fig2 = px.pie(cat_sum, names="카테고리", values="언급량", hole=0.5,
                          color_discrete_sequence=["#2D6BC4","#E67E22","#27AE60"],
                          title="제품 카테고리 비중")
            chart_style(fig2, height=300)
            st.plotly_chart(fig2, use_container_width=True)
            st.markdown("<div class='section-title'>최다 언급 제품 TOP 5</div>", unsafe_allow_html=True)
            for _, row in prod_df.head(5).iterrows():
                nss_color = "#27AE60" if row["NSS"] > 0 else "#E74C3C"
                st.markdown(f"""
                <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-left:3px solid {nss_color};
                            border-radius:0 6px 6px 0;padding:6px 10px;margin-bottom:4px;'>
                    <b style='font-size:0.85rem;'>{row['제품']}</b>
                    <span style='float:right;color:{nss_color};font-weight:700;'>NSS {row['NSS']:+.0f}</span>
                    <div style='color:#718096;font-size:0.72rem;'>{row['카테고리']} · {row['언급량']}건</div>
                </div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>제품별 월별 언급량 · 감성 추이</div>", unsafe_allow_html=True)
        sel_prod = st.selectbox("제품 선택", prod_df["제품"].tolist(), key="product_item")
        ptrend = compute_product_trend(v, sel_bp, sel_prod)
        if len(ptrend) > 0 and ptrend["언급량"].sum() > 0:
            fig3 = make_subplots(specs=[[{"secondary_y": True}]])
            fig3.add_trace(go.Bar(
                x=ptrend["월"], y=ptrend["언급량"], name="언급량",
                marker_color="#2D6BC4", opacity=0.75,
            ), secondary_y=False)
            fig3.add_trace(go.Scatter(
                x=ptrend["월"], y=ptrend["NSS"], name="NSS", mode="lines+markers",
                line=dict(color="#E67E22", width=3), marker=dict(size=7),
            ), secondary_y=True)
            fig3.add_hline(y=0, line_dash="dot", line_color="#A0AEC0", secondary_y=True)
            fig3.update_layout(
                title=f"{sel_bp} — {sel_prod} 월별 언급량 & 감성 추이",
                paper_bgcolor="#F7F9FC", plot_bgcolor="#FFFFFF",
                font=dict(color="#2D3748"), height=400,
                legend=dict(orientation="h", y=1.1),
            )
            fig3.update_yaxes(title_text="언급량", secondary_y=False)
            fig3.update_yaxes(title_text="NSS", secondary_y=True, range=[-100, 100])
            st.plotly_chart(fig3, use_container_width=True)
        else:
            st.info(f"{sel_prod}의 월별 추이 데이터가 부족합니다.")

        st.dataframe(prod_df, use_container_width=True, hide_index=True)


# ════════════════════════════════════════════════════
# TAB — 프로모션 모니터링 [항목 5]
# ════════════════════════════════════════════════════
def _rt_promo():
    st.markdown("<div class='section-title'>🎁 프로모션 모니터링</div>", unsafe_allow_html=True)
    st.markdown(f"""
    <div class='insight-box'>
    통신·카드할인, 이벤트, 1+1, 굿즈 증정 등 <b>촉진(프로모션) 관련 담론</b>을 브랜드별로 모니터링하고,
    "할인 개꿀", "이벤트 필참", "잇템" 등 <b>고객 반응 강도</b>를 측정합니다.
    (전체 {promo_doc_count}건의 프로모션 언급 분석)
    </div>""", unsafe_allow_html=True)

    # [v2.4.1] 브랜드별 프로모션 언급량 시계열
    st.markdown("<div class='section-title'>📈 브랜드별 프로모션 언급량 추이 (시계열)</div>", unsafe_allow_html=True)
    _promo_pat = "|".join(sum(PROMO_DICT.values(), []))
    _pv = v[v["full_text"].str.contains(_promo_pat, na=False)].copy()
    _months_p = sorted([m for m in _pv["year_month"].dropna().unique() if "NaT" not in str(m)])
    if _months_p:
        _ts_rows = []
        for b in TARGET_BRANDS:
            bd = _pv[_pv["brands"].apply(lambda x: b in x)]
            if len(bd) < 3:
                continue
            cnt = bd.groupby("year_month").size()
            for m in _months_p:
                _ts_rows.append({"월": str(m), "브랜드": b, "언급량": int(cnt.get(m, 0))})
        if _ts_rows:
            _ts_df = pd.DataFrame(_ts_rows)
            fig_pts = px.line(_ts_df, x="월", y="언급량", color="브랜드",
                              color_discrete_map=BRAND_COLORS, markers=True,
                              title="브랜드별 월간 프로모션 언급량 추이")
            chart_style(fig_pts, height=380)
            st.plotly_chart(fig_pts, use_container_width=True)
            st.caption("프로모션 키워드 포함 문서를 브랜드·월별 집계 (언급 3건 미만 브랜드 제외)")
        else:
            st.caption("브랜드별 시계열을 표시할 데이터가 부족합니다.")
    else:
        st.caption("월별 데이터가 부족합니다.")

    c1, c2 = st.columns([1, 1])
    with c1:
        st.markdown("<div class='section-title'>프로모션 유형별 언급량</div>", unsafe_allow_html=True)
        fig = px.bar(
            promo_total_df.sort_values("언급량"), x="언급량", y="프로모션",
            orientation="h", color="언급량",
            color_continuous_scale=["#FDE9D9","#E67E22"],
            title="촉진 유형별 언급량", text="언급량",
        )
        fig.update_traces(textposition="outside")
        chart_style(fig, height=340, showlegend=False)
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        st.markdown("<div class='section-title'>촉진 반응 강도 분포</div>", unsafe_allow_html=True)
        fig2 = px.bar(
            promo_react_df, x="반응", y="언급량",
            color="반응",
            color_discrete_map={"강한 긍정":"#C0392B","긍정":"#E67E22","부정":"#7F8C8D"},
            title="프로모션 반응 강도 분포", text="언급량",
        )
        fig2.update_traces(textposition="outside")
        chart_style(fig2, height=340, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)

    # [v2.3.1] 최근 급증 프로모션 + 브랜드 cross 체크
    st.markdown("<div class='section-title'>🔥 최근 급증한 프로모션 활동</div>", unsafe_allow_html=True)
    st.caption("최근 1개월 언급량을 직전 3개월 월평균과 비교 — 🔥는 언급 2건 이상 & 50% 이상 급증")
    if len(promo_surge_df) > 0:
        cc1, cc2 = st.columns(2)
        with cc1:
            st.dataframe(promo_surge_df, use_container_width=True, hide_index=True, height=330)
        with cc2:
            fig_s = px.bar(promo_surge_df.sort_values("증감률"), x="증감률", y="프로모션",
                           orientation="h", color="증감률",
                           color_continuous_scale=["#2980B9", "#EDF1F5", "#E74C3C"],
                           color_continuous_midpoint=0, title="프로모션별 최근 증감률 (%)")
            chart_style(fig_s, height=330)
            st.plotly_chart(fig_s, use_container_width=True)
    st.markdown("**프로모션 × 브랜드 Cross 체크 (최근 2개월) — 어떤 브랜드가 어떤 촉진을 주도했는가**")
    if hasattr(promo_cross_df, "values") and promo_cross_df.values.sum() > 0:
        fig_c = px.imshow(promo_cross_df, text_auto=True, aspect="auto",
                          color_continuous_scale="Oranges",
                          labels=dict(x="브랜드", y="프로모션", color="언급"))
        chart_style(fig_c, height=380)
        st.plotly_chart(fig_c, use_container_width=True)
    else:
        st.caption("최근 2개월 프로모션 교차 데이터가 충분하지 않습니다.")

    # [기능 4] 최근 많이 언급된 프로모션 + 주력 브랜드
    st.markdown("<div class='section-title'>📅 최근 가장 많이 언급된 프로모션 & 주력 브랜드</div>", unsafe_allow_html=True)
    st.caption("최근 2개월 기준 — 어떤 촉진 활동이 화제이고, 어떤 브랜드가 그 활동을 주도하는지")
    if len(promo_recent_df) > 0:
        rc1, rc2 = st.columns([1.3, 1])
        with rc1:
            fig_r = px.bar(
                promo_recent_df.sort_values("최근언급"),
                x="최근언급", y="프로모션", orientation="h",
                color="최근언급", color_continuous_scale=["#FDE9D9","#E67E22"],
                title="최근 2개월 프로모션 언급량", text="최근언급",
            )
            fig_r.update_traces(textposition="outside")
            chart_style(fig_r, height=320, showlegend=False)
            st.plotly_chart(fig_r, use_container_width=True)
        with rc2:
            st.markdown("**프로모션별 주력 브랜드**")
            for _, row in promo_recent_df.head(7).iterrows():
                bcolor = BRAND_COLORS.get(row["주력브랜드"], "#64748B")
                st.markdown(f"""
                <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-left:3px solid {bcolor};
                            border-radius:0 6px 6px 0;padding:6px 10px;margin-bottom:4px;'>
                    <b style='font-size:0.82rem;'>{row['프로모션']}</b>
                    <span style='float:right;color:{bcolor};font-weight:700;font-size:0.8rem;'>{row['주력브랜드']}</span>
                    <div style='color:#718096;font-size:0.7rem;'>최근 {row['최근언급']}건 · 주력 {row['주력브랜드_언급']}건</div>
                </div>""", unsafe_allow_html=True)

    # [기능 4] 브랜드별 프로모션 적극도
    st.markdown("<div class='section-title'>브랜드별 프로모션 적극도</div>", unsafe_allow_html=True)
    st.caption("프로모션 언급 비중 = 해당 브랜드 프로모션 언급 ÷ 전체 언급. 프로모션을 가장 활발히 하는 브랜드 식별")
    if len(brand_promo_df) > 0:
        fig_bp = px.bar(
            brand_promo_df.sort_values("프로모션언급"),
            x="프로모션언급", y="브랜드", orientation="h",
            color="프로모션비중", color_continuous_scale=["#FDE9D9","#C0392B"],
            title="브랜드별 프로모션 언급량 (색상=프로모션 비중%)",
            text=brand_promo_df.sort_values("프로모션언급").apply(
                lambda r: f"{r['프로모션언급']}건 ({r['프로모션비중']}%)", axis=1),
        )
        fig_bp.update_traces(textposition="outside")
        chart_style(fig_bp, height=340, showlegend=True)
        st.plotly_chart(fig_bp, use_container_width=True)

    st.markdown("<div class='section-title'>프로모션 유형 × 브랜드 히트맵</div>", unsafe_allow_html=True)
    st.caption("어떤 브랜드가 어떤 촉진 활동으로 가장 많이 회자되는지 — 빈 셀은 파스쿠찌 차별화 기회")
    fig3 = px.imshow(
        promo_brand_mat, color_continuous_scale="Oranges",
        title="프로모션 유형 × 브랜드 언급 매트릭스",
        text_auto=True, aspect="auto",
    )
    fig3.update_traces(textfont=dict(color="#1B2A4A", size=11))
    chart_style(fig3, height=360, showlegend=False)
    st.plotly_chart(fig3, use_container_width=True)

    # [기능 4] 촉진 활동별 반응 강도 — 0~100 지표화
    st.markdown("<div class='section-title'>촉진 활동별 반응 강도 지표 (0~100)</div>", unsafe_allow_html=True)
    st.caption("반응강도지표 = (강한긍정×2 + 긍정×1 − 부정×1)/언급량 을 0~100으로 정규화. 지속 관리용 지표")
    if len(promo_type_react_df) > 0:
        fig4 = px.bar(
            promo_type_react_df.sort_values("반응강도지표"),
            x="반응강도지표", y="프로모션", orientation="h",
            color="반응강도지표", color_continuous_scale=["#E74C3C","#F5C77E","#27AE60"],
            range_color=[0, 100],
            title="촉진 유형별 반응 강도 지표 (0~100)",
            text=promo_type_react_df.sort_values("반응강도지표")["반응강도지표"].apply(lambda x: f"{x:.0f}"),
            hover_data=["언급량","강한긍정","긍정","부정"],
        )
        fig4.update_traces(textposition="outside")
        chart_style(fig4, height=340, showlegend=True)
        st.plotly_chart(fig4, use_container_width=True)

    # [기능 5] 촉진 활동 세부 + 원문 출처 확인
    st.markdown("<div class='section-title'>🔎 프로모션 활동 세부 & 원문 확인</div>", unsafe_allow_html=True)
    st.caption("프로모션 유형을 선택하면 해당 키워드가 실제 포함된 문장을 발췌해 보여줍니다(무관 글 제외, 관련도순). [v2.5.2 정확도 개선]")
    sel_promo = st.selectbox("프로모션 유형 선택", list(PROMO_DICT.keys()), key="promo_detail_sel")
    detail_docs = get_promo_detail_docs(v, sel_promo)
    if len(detail_docs) == 0:
        st.info(f"{sel_promo} 관련 언급이 없습니다.")
    else:
        dc1, dc2 = st.columns([1, 2])
        with dc1:
            st.markdown(f"**{sel_promo} — 브랜드별 언급**")
            bc = Counter()
            for _, r in detail_docs.iterrows():
                for b in r["brands"]:
                    bc[b] += 1
            if bc:
                bc_df = pd.DataFrame([{"브랜드": b, "언급": c} for b, c in bc.most_common()])
                fig_d = px.bar(
                    bc_df.sort_values("언급"), x="언급", y="브랜드", orientation="h",
                    color="언급", color_continuous_scale=["#FDE9D9","#E67E22"],
                    text="언급",
                )
                fig_d.update_traces(textposition="outside")
                chart_style(fig_d, height=300, showlegend=False)
                st.plotly_chart(fig_d, use_container_width=True)
        with dc2:
            st.markdown(f"**{sel_promo} — 원문 출처 ({len(detail_docs)}건)**")
            _cols = ["title","발췌","primary_brand","year_month"] + (["post_url"] if "post_url" in detail_docs.columns else [])
            detail_show = detail_docs[_cols].copy()
            detail_show.columns = ["제목","원문 발췌(키워드 문장)","주요브랜드","작성월"] + (["원문링크"] if "post_url" in detail_docs.columns else [])
            st.dataframe(detail_show.head(20), use_container_width=True, height=320,
                         column_config={
                             "원문 발췌(키워드 문장)": st.column_config.TextColumn(width="large"),
                             "원문링크": st.column_config.LinkColumn(width="small") if "post_url" in detail_docs.columns else None,
                         })
            st.caption(f"{min(20, len(detail_show))}건 표시 / 전체 {len(detail_show)}건")

    st.markdown("""
    <div class='success-box'>
    <b>전략 활용</b>: 반응 강도 지표가 높은 촉진 유형(굿즈·1+1 등)은 파스쿠찌도 적극 도입,
    경쟁사가 독점한 촉진 영역은 차별화된 방식으로 접근. 최근 화제 프로모션은 즉시 벤치마킹.
    </div>""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════
# TAB — 고객 클러스터링 & 페르소나 [항목 7]
# ════════════════════════════════════════════════════
def _rt_cluster():
    st.markdown("<div class='section-title'>👥 고객 클러스터링 & 페르소나</div>", unsafe_allow_html=True)
    st.markdown("""
    <div style='background:#F4F6FA;border:1px dashed #B6C2D6;border-radius:8px;
                padding:10px 14px;margin:2px 0 10px;color:#5A6678;font-size:0.86rem;'>
    <b style='color:#7A8699;'>ℹ️ 분석 설명</b> &nbsp; 고객을 <b>성향·특성·관심사</b>로 군집화해 <b>페르소나</b>를 도출하고,
    각 페르소나의 카페·커피 U&A, 관심사, 주요 브랜드를 정리합니다.
    <span style='color:#9AA6B8;'>(TF-IDF 텍스트 70% + 감성·충성도·브랜드수·디카페인 관심 등 행동 피처 30%)</span>
    </div>""", unsafe_allow_html=True)

    if customer_cluster is None:
        st.info("고객 클러스터링을 위한 데이터가 부족합니다. (30건 이상 필요)")
    else:
        sub_c = customer_cluster["sub"]
        profiles = customer_cluster["profiles"]
        pca_var = customer_cluster["pca_var"]

        st.markdown("<div class='section-title'>고객 세그먼트 분포 (PCA 2D)</div>", unsafe_allow_html=True)
        cluster_palette = ["#2D6BC4","#E67E22","#27AE60","#8E44AD","#E74C3C","#16A085",
                           "#D4A017","#C0392B","#5D6D7E","#117A65"]
        fig = go.Figure()
        for prof in profiles:
            c = prof["cluster"]
            cdocs = sub_c[sub_c["cluster"] == c]
            persona = make_persona(prof)
            fig.add_trace(go.Scatter(
                x=cdocs["pca_x"], y=cdocs["pca_y"], mode="markers",
                name=f"C{c}: {persona['name']}",
                marker=dict(size=8, color=cluster_palette[c % len(cluster_palette)],
                            opacity=0.6, line=dict(color="white", width=0.5)),
                hovertemplate=f"<b>{persona['name']}</b><br>%{{text}}<extra></extra>",
                text=cdocs["title"].str[:30],
            ))
        bc_cent = customer_cluster.get("brand_centroids", {})
        if bc_cent:
            fig.add_trace(go.Scatter(
                x=[xy[0] for xy in bc_cent.values()], y=[xy[1] for xy in bc_cent.values()],
                mode="markers+text", text=list(bc_cent.keys()), textposition="top center",
                textfont=dict(size=11, color="#1B2A4A"),
                marker=dict(symbol="diamond", size=14, color="#1B2A4A",
                            line=dict(color="white", width=1.2)),
                name="브랜드 위치", hoverinfo="text"))
        fig.update_layout(
            title=f"고객 세그먼트 분포 + 인접 브랜드 (PC1 {pca_var[0]}% · PC2 {pca_var[1]}%)",
            xaxis_title="주성분 1", yaxis_title="주성분 2",
            paper_bgcolor="#F7F9FC", plot_bgcolor="#FFFFFF",
            font=dict(color="#2D3748"), height=480,
            legend=dict(bgcolor="#FFFFFF", bordercolor="#E2E8F0", borderwidth=1),
        )
        st.plotly_chart(fig, use_container_width=True)

        st.markdown("<div class='section-title'>도출된 고객 페르소나</div>", unsafe_allow_html=True)
        st.caption("각 페르소나 카드 아래 '📂 상세 분석 보기'를 펼치면 마케팅 활용 정보와 감성·키워드 차트를 볼 수 있습니다.")
        for prof in profiles:
            persona = make_persona(prof)
            c = prof["cluster"]
            color = cluster_palette[c % len(cluster_palette)]
            brands_str = ", ".join(persona["brands"][:3]) if persona["brands"] else "특정 브랜드 없음"
            kw_str = ", ".join(persona["kw"][:6])
            life_str = " · ".join(persona["lifestyle"]) if persona["lifestyle"] else "—"
            co_str = ", ".join(persona["co_brands"]) if persona["co_brands"] else "—"
            st.markdown(f"""
            <div style='background:#FFFFFF;border:1px solid #E2E8F0;border-left:5px solid {color};
                        border-radius:10px;padding:14px 16px;margin:10px 0 0;
                        box-shadow:0 2px 6px rgba(0,0,0,0.05);'>
                <div style='color:{color};font-weight:800;font-size:1.05rem;'>C{c}. {persona['name']}</div>
                <div style='color:#718096;font-size:0.76rem;margin:2px 0 8px;'>👥 n={persona['n']}명 · 😊 NSS {persona['nss']:+.0f} · 💎 충성 {persona['loyalty']}% · 🌙 디카페인 {persona['decaf']}%</div>
                <div style='color:#2D3748;font-size:0.86rem;margin-bottom:8px;'>{persona['desc']}</div>
                <div style='font-size:0.78rem;color:#4A5568;'>
                    🏷️ <b>주요 브랜드</b>: {brands_str}　🌿 <b>라이프스타일</b>: {life_str}<br>
                    🔑 <b>관심 키워드</b>: {kw_str}
                </div>
            </div>""", unsafe_allow_html=True)
            with st.expander(f"📂 C{c}. {persona['name']} — 상세 분석 보기"):
                dc1, dc2 = st.columns([1, 1])
                with dc1:
                    st.markdown(f"📝 **설명** — {persona['desc']}")
                    st.markdown(f"☕ **카페·커피 U&A** — {persona['ua']}")
                    st.markdown(f"🎯 **추천 마케팅 액티비티** — {persona.get('marketing', '—')}")
                    st.markdown(f"😣 **Pain-point** — {persona.get('pain', '—')}")
                    st.markdown(f"🛒 **Key Buying Factor** — {persona.get('kbf', '—')}")
                    st.markdown(f"🌿 **라이프스타일·관심도** — {' · '.join(persona['lifestyle']) or '—'}")
                    st.markdown(f"📍 **주요 소비상황** — {', '.join(persona['occ']) or '—'}")
                    st.markdown(f"🏷️ **커피 브랜드** — {', '.join(persona['brands'][:3]) or '—'}")
                    st.markdown(f"🤝 **커피 외 동시언급 브랜드** — {', '.join(persona['co_brands']) or '—'}")
                    st.markdown(f"🔑 **관심 키워드** — {', '.join(persona['kw'][:8])}")
                with dc2:
                    cdocs = sub_c[sub_c["cluster"] == c]
                    sent_cnt = cdocs["sentiment"].value_counts()
                    if len(sent_cnt) > 0:
                        figp = px.pie(values=sent_cnt.values, names=sent_cnt.index, hole=0.5,
                                      color=sent_cnt.index, color_discrete_map=SENT_COLORS,
                                      title="세그먼트 감성 분포")
                        chart_style(figp, height=240)
                        st.plotly_chart(figp, use_container_width=True, key=f"persona_sent_{c}")
                    _alltok = [t for ts in cdocs["tokens"] for t in ts]
                    _allow = set(_kw_clean(_alltok, 10000))
                    kwc = Counter([t for t in _alltok if t in _allow]).most_common(8)
                    if kwc:
                        figk = px.bar(x=[val for _, val in kwc][::-1], y=[k for k, _ in kwc][::-1],
                                      orientation="h", title="대표 키워드 Top 8 (정제)",
                                      color_discrete_sequence=["#8E44AD"])
                        chart_style(figk, height=240)
                        st.plotly_chart(figk, use_container_width=True, key=f"persona_kw_{c}")

        st.markdown("<div class='section-title'>세그먼트 × 브랜드 언급 히트맵</div>", unsafe_allow_html=True)
        seg_brand = pd.DataFrame(0, index=[f"C{p['cluster']}" for p in profiles], columns=TARGET_BRANDS)
        for prof in profiles:
            c = prof["cluster"]
            cdocs = sub_c[sub_c["cluster"] == c]
            for b in TARGET_BRANDS:
                seg_brand.loc[f"C{c}", b] = int(cdocs["brands"].apply(lambda x: b in x).sum())
        fig2 = px.imshow(
            seg_brand, color_continuous_scale="Purples",
            title="고객 세그먼트 × 브랜드 언급 매트릭스",
            text_auto=True, aspect="auto",
        )
        fig2.update_traces(textfont=dict(color="#1B2A4A", size=11))
        chart_style(fig2, height=320, showlegend=False)
        st.plotly_chart(fig2, use_container_width=True)


# ════════════════════════════════════════════════════
# [v2.5.1] 카테고리(segmented control) → 세부 탭 디스패치
# ════════════════════════════════════════════════════
_CATS = {
    "📊 시장·브랜드 진단": [
        ("📊 Overview", _rt_ov), ("📈 NSS 평판", _rt_nss), ("📣 SOV 분석", _rt_sov),
        ("💎 브랜드 로열티", _rt_loyalty), ("📍 포지셔닝맵", _rt_pos),
    ],
    "🧠 고객 이해": [
        ("👥 고객 클러스터링", _rt_cluster), ("🎯 소비 맥락", _rt_occ),
        ("🔍 드라이버 분석", _rt_absa), ("🗺️ 브랜드 이미지맵", _rt_ca),
    ],
    "🛍️ 상품기획 지원": [
        ("🛍️ 브랜드별 제품", _rt_product), ("📡 트렌드 조기경보", _rt_trend),
        ("☁️ 키워드 버블", _rt_kw),
    ],
    "📣 커뮤니케이션 인텔리전스": [
        ("🎯 LDA·토픽 클러스터", _rt_lda), ("🎁 프로모션 모니터링", _rt_promo),
        ("🔗 PMI 분석", _rt_pmi),
    ],
    "🚨 이슈 캐처": [
        ("⚡ Burst 탐지", _rt_burst), ("🚨 리스크 탐지", _rt_risk),
    ],
}
_CAT_DESC = {
    "📊 시장·브랜드 진단": "브랜드별 Status는? — 언급량·평판·SOV·로열티·포지셔닝 등",
    "🧠 고객 이해": "우리의 고객은 어떤 사람들인가? — 고객 프로파일링, 성향·특성·관심사별 군집분석, 만족/불만족 동인, 브랜드 이미지 매핑 등",
    "🛍️ 상품기획 지원": "지금 사람들의 관심 제품은 무엇인가? — 제품 담론 분석, 제품 유형별 트렌드 캐칭, 브랜드별 주요 키워드 등",
    "📣 커뮤니케이션 인텔리전스": "고객을 움직이는 것은 무엇인가? — 주제별 담론 분석, 프로모션 모니터, 경쟁 인식 분석 등",
    "🚨 이슈 캐처": "무엇을 조심할까? — 이상 급증·리스크 신호",
}
_cat_names = list(_CATS.keys())
# [v2.5.3] 메인 최상단 타이틀 + 대형 카테고리 버튼 + 세부탭 스타일
st.markdown("""
<style>
/* [v2.5.4] 카테고리 대형 버튼 — 글씨 키우고 볼드(라벨 내부 p/div까지 강제) */
div[data-testid="stHorizontalBlock"] .stButton button {
    min-height: 72px; font-size: 1.12rem; font-weight: 800;
    white-space: normal; line-height: 1.25; border-radius: 12px; padding: 8px 8px;
}
div[data-testid="stHorizontalBlock"] .stButton button p,
div[data-testid="stHorizontalBlock"] .stButton button div {
    font-weight: 800 !important;
}
/* [D2] 활성(선택) 버튼 — 진한 남색 대신 옅은 블루 배경 + 진한 글씨로 가독성 확보 */
div[data-testid="stHorizontalBlock"] .stButton button[kind="primary"],
div[data-testid="stHorizontalBlock"] .stButton button[data-testid="baseButton-primary"] {
    background-color: #DCE9F7 !important; border: 2px solid #4A89DC !important; color: #1B2A4A !important;
}
div[data-testid="stHorizontalBlock"] .stButton button[kind="primary"]:hover,
div[data-testid="stHorizontalBlock"] .stButton button[data-testid="baseButton-primary"]:hover {
    background-color: #CFE0F3 !important; border-color: #4A89DC !important; color: #1B2A4A !important;
}
/* [D2] 세부 분석 탭 — 영역/폰트 키우고 볼드 */
div[data-baseweb="tab-list"] { gap: 6px; }
button[data-baseweb="tab"] { padding-top: 12px !important; padding-bottom: 12px !important; }
button[data-baseweb="tab"] div[data-testid="stMarkdownContainer"] p,
button[data-baseweb="tab"] p {
    font-size: 1.06rem !important; font-weight: 700 !important;
}
</style>
""", unsafe_allow_html=True)
# [D4] 메인 화면 최상단 타이틀
st.markdown("<div style='font-weight:800;font-size:1.6rem;color:#1B2A4A;"
            "margin:0 0 12px;letter-spacing:0.3px;'>Consumer Intelligence Dashboard</div>",
            unsafe_allow_html=True)
if st.session_state.get("cat_sel") not in _cat_names:
    st.session_state["cat_sel"] = _cat_names[0]
st.markdown("<div style='font-weight:800;font-size:1.08rem;color:#1B2A4A;margin:2px 0 8px;'>"
            "🧭 분석 영역</div>", unsafe_allow_html=True)
_bcols = st.columns(len(_cat_names))
for _i, _ck in enumerate(_cat_names):
    with _bcols[_i]:
        _active = (st.session_state["cat_sel"] == _ck)
        if st.button(_ck, key=f"catbtn_{_i}", use_container_width=True,
                     type=("primary" if _active else "secondary")):
            st.session_state["cat_sel"] = _ck
            (st.rerun if hasattr(st, "rerun") else st.experimental_rerun)()
_cat = st.session_state["cat_sel"]
# [D1] 카테고리 설명 — 풋노트(inform) 형태로 축소
st.markdown(
    f"<div style='background:#F4F7FB;border-left:3px solid #9DB8DC;border-radius:6px;"
    f"padding:5px 11px;margin:6px 0 12px;font-size:0.78rem;color:#6B7888;'>"
    f"ⓘ <b style='color:#5A6678;'>{_cat}</b> · {_CAT_DESC.get(_cat, '')}</div>",
    unsafe_allow_html=True)
_items = _CATS[_cat]
_subtabs = st.tabs([_lbl for _lbl, _ in _items])
for (_lbl, _fn), _tb in zip(_items, _subtabs):
    with _tb:
        _fn()


# ── 푸터 ──────────────────────────────────────────────
st.markdown("---")
st.markdown(
    f"<p style='text-align:center;color:#A0AEC0;font-size:11px;'>"
    f"{COPYRIGHT}"
    f"</p>",
    unsafe_allow_html=True,
)